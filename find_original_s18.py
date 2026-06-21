import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

# v4에서 슬라이드 18 (0-based 17) 확인
for path, label in [
    (r'C:\Users\jay\Downloads\AI교육_완성판_v4.pptx', 'v4'),
    (r'C:\Users\jay\Downloads\AI교육_완성판_v3.pptx', 'v3'),
]:
    try:
        prs = Presentation(path)
        print(f'\n=== {label} — 슬라이드 18 (idx=17) ===')
        slide = prs.slides[17]
        for i, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = para.text.strip()
                    if txt:
                        print(f'  [{i}] "{txt}"')
    except Exception as e:
        print(f'  {label} 파일 없음/오류: {e}')

# 원본 파일
try:
    prs0 = Presentation(r'C:\Users\jay\Downloads\AI교육 최종.pptx')
    print(f'\n=== 원본 — 슬라이드 17 (idx=16) ===')
    slide = prs0.slides[16]
    for i, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                txt = para.text.strip()
                if txt:
                    print(f'  [{i}] "{txt}"')
except Exception as e:
    print(f'  원본 파일 오류: {e}')
