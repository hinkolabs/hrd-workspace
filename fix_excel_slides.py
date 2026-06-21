"""
Slide 19: 페이지 번호 | 12 → | 19
Slide 20: AI 프롬프트 → 엑셀 강사파일 v3 기준 상세 버전, 이모티콘 제거
Slide 21: 오류 조건 6개 → 5개 (필수값 누락 제거), 조건 텍스트 정밀 수정
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SRC = r'C:\Users\jay\Downloads\AI교육_완성판_v11.pptx'
DST = r'C:\Users\jay\Downloads\AI교육_완성판_v12.pptx'
prs = Presentation(SRC)

FONT = '맑은 고딕'


def set_run(shape, text, size=None, bold=None, color=None, italic=None):
    tf = shape.text_frame
    p  = tf.paragraphs[0]
    from pptx.oxml.ns import qn
    # 기존 runs 내용 교체
    if p.runs:
        run = p.runs[0]
        # 나머지 runs 제거
        for r in p.runs[1:]:
            r.text = ''
    else:
        r_el = p._p.makeelement(qn('a:r'))
        p._p.append(r_el)
        run = p.runs[0]
    run.text = text
    run.font.name = FONT
    if size   is not None: run.font.size      = Pt(size)
    if bold   is not None: run.font.bold      = bold
    if color  is not None: run.font.color.rgb = color
    if italic is not None: run.font.italic    = italic


def set_multiline(shape, text, size=None, bold=None, color=None):
    """텍스트박스 전체를 새 텍스트로 교체 (줄바꿈 포함)"""
    tf = shape.text_frame
    tf.word_wrap = True
    from pptx.oxml.ns import qn
    # 기존 paragraphs 모두 제거
    txBody = tf._txBody
    for p_el in txBody.findall(qn('a:p')):
        txBody.remove(p_el)
    # 새 paragraphs 추가
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p_el = txBody.makeelement(qn('a:p'))
        txBody.append(p_el)
        if line:
            r_el = p_el.makeelement(qn('a:r'))
            rPr  = r_el.makeelement(qn('a:rPr'), {'lang': 'ko-KR', 'dirty': '0'})
            if bold:
                rPr.set('b', '1')
            r_el.insert(0, rPr)
            t_el = r_el.makeelement(qn('a:t'))
            t_el.text = line
            r_el.append(t_el)
            # 폰트 설정
            latin = rPr.makeelement(qn('a:latin'), {'typeface': FONT})
            rPr.append(latin)
            if size:
                rPr.set('sz', str(int(size * 100)))
            if color:
                solidFill = rPr.makeelement(qn('a:solidFill'))
                srgbClr   = solidFill.makeelement(qn('a:srgbClr'),
                                                  {'val': f'{color.rgb:06X}'})
                solidFill.append(srgbClr)
                rPr.append(solidFill)
            p_el.append(r_el)
        else:
            # 빈 줄
            r_el = p_el.makeelement(qn('a:r'))
            t_el = r_el.makeelement(qn('a:t'))
            t_el.text = ''
            r_el.append(t_el)
            p_el.append(r_el)


# ══════════════════════════════════════════════════════════
#  Slide 19 — 페이지 번호 수정
# ══════════════════════════════════════════════════════════
slide19 = prs.slides[18]
shapes19 = list(slide19.shapes)
for shape in shapes19:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if para.text.strip() == '| 12':
                for run in para.runs:
                    if '12' in run.text:
                        run.text = run.text.replace('12', '19')
                        print('  [Slide 19] 페이지 번호 12 → 19')
                        break

# ══════════════════════════════════════════════════════════
#  Slide 20 — AI 프롬프트 교체 + 이모티콘 제거
# ══════════════════════════════════════════════════════════
slide20  = prs.slides[19]
shapes20 = list(slide20.shapes)

# 이모티콘 제거: 섹션 헤더 레이블
emoji_fixes = {
    8:  '거래원장_RAW 시트 구조 (A~J열)',
    50: '결과작성_STEP1 시트 — 집계표 3종',
    61: 'AI 프롬프트 (STEP1)',
    65: '체크포인트',
}
for idx, txt in emoji_fixes.items():
    set_run(shapes20[idx], txt)
    print(f'  [Slide 20] shape[{idx}] 이모티콘 제거: "{txt}"')

# AI 프롬프트 전체 교체 (shape[63])
new_prompt = (
    '아래 거래원장 시트 구조를 참고해서 보고용 집계표를 만들려고 합니다.\n'
    '\n'
    '컬럼:\n'
    'A 거래일자, B 지점코드, C 직원ID, D 고객ID,\n'
    'E 상품유형, F 거래유형, G 거래금액,\n'
    'H 수수료율, I 수수료금액, J 처리상태\n'
    '\n'
    '만들어야 할 표:\n'
    '① 지점별(B001~B005)\n'
    '   처리상태=정상인 거래건수·거래금액합계·수수료합계,\n'
    '   취소+오류 건수\n'
    '② 상품유형별(주식/펀드/채권/ELS/RP)\n'
    '   처리상태=정상인 거래건수·거래금액합계·수수료합계,\n'
    '   전체 정상 거래금액 대비 비중\n'
    '③ 직원별(E001~E020)\n'
    '   직원마스터 시트에서 직원명·지점코드 VLOOKUP,\n'
    '   정상 거래건수·정상 수수료합계·취소+오류 건수\n'
    '\n'
    '결과작성_STEP1_보고용집계 시트에 들어갈\n'
    'Excel 수식으로 만들어줘.\n'
    'RAW 데이터는 거래원장_RAW 시트\n'
    '5행부터 시작, 204행까지 데이터 있습니다.'
)
set_multiline(shapes20[63], new_prompt, size=9.5)
print('  [Slide 20] AI 프롬프트 엑셀 v3 기준 교체 완료')

# ══════════════════════════════════════════════════════════
#  Slide 21 — 오류 조건 수정 + 필수값 누락(06) 제거
# ══════════════════════════════════════════════════════════
slide21  = prs.slides[20]
shapes21 = list(slide21.shapes)

# 조건 텍스트 업데이트 (엑셀 강사파일 기준)
cond_updates = {
    16: '처리상태=정상\nAND 거래금액(G열) ≤ 0',
    23: '처리상태=정상\nAND 수수료율(H열) ≤ 0\n또는 0.5% 초과',
    30: '처리상태=정상\nAND 수수료금액과\n거래금액×수수료율\n차이 1원 초과',
    37: '처리상태=정상\nAND 주식·펀드·채권·\nELS·RP 외 값',
    44: '처리상태=정상\nAND 매수·매도·가입·\n해지 외 값',
}
for idx, txt in cond_updates.items():
    set_multiline(shapes21[idx], txt, size=10)
    print(f'  [Slide 21] shape[{idx}] 조건 텍스트 수정')

# 필수값 누락 카드 (shapes 45~51) 제거
sp_tree = slide21.shapes._spTree
for i in range(45, 52):
    el = shapes21[i]._element
    el.getparent().remove(el)
print('  [Slide 21] 필수값 누락(06) 카드 제거 완료')

# 이모티콘 제거: 프롬프트 헤더 (shape[52] 이후 인덱스 재계산)
# 52번 이후 shapes 재로드
shapes21r = list(slide21.shapes)
for shape in shapes21r:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            txt = para.text.strip()
            if '💬' in txt:
                for run in para.runs:
                    run.text = run.text.replace('💬 ', '').replace('💬', '')
                print(f'  [Slide 21] 이모티콘 제거: {txt[:40]}')

# ══════════════════════════════════════════════════════════
prs.save(DST)
print(f'\n✓ 저장: {DST}')
