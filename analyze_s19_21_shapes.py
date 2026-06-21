import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

SRC = r'C:\Users\jay\Downloads\AI교육_완성판_v11.pptx'
prs = Presentation(SRC)

for si in [19, 20, 21]:
    slide = prs.slides[si - 1]
    print(f'\n=== Slide {si} ({len(slide.shapes)} shapes) ===')
    for i, shape in enumerate(slide.shapes):
        L = shape.left/914400
        T = shape.top/914400
        W = shape.width/914400
        H = shape.height/914400
        txt = ''
        if shape.has_text_frame:
            lines = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
            if lines:
                txt = ' | ' + ' /// '.join(lines)[:120]
        print(f'  [{i:2d}] T={T:.2f}" L={L:.2f}" W={W:.2f}" H={H:.2f}"{txt}')
