"""
student_v4 기준 엑셀 실습 슬라이드 19~21 전면 재구성
- Slide 19: 섹션 디바이더 (내용 업데이트)
- Slide 20: STEP1 — 직원별 집계표 (VLOOKUP + COUNTIFS)
- Slide 21: STEP2 — 2개 오류 조건 (거래금액 / 상품유형)
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

SRC = r'C:\Users\jay\Downloads\AI교육_완성판_v12.pptx'
DST = r'C:\Users\jay\Downloads\AI교육_완성판_v13.pptx'
prs = Presentation(SRC)

GREEN       = RGBColor(0x00, 0x91, 0x78)
DARK_GREEN  = RGBColor(0x00, 0x4E, 0x42)
LIGHT_GREEN = RGBColor(0xD7, 0xED, 0xE6)
LIGHT_GRAY  = RGBColor(0xF7, 0xF9, 0xF8)
DARK_TEXT   = RGBColor(0x22, 0x22, 0x22)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY        = RGBColor(0xCC, 0xCC, 0xCC)
ORANGE      = RGBColor(0xE8, 0x6B, 0x00)
FONT        = '맑은 고딕'
SLIDE_W     = prs.slide_width
SLIDE_H     = prs.slide_height


def rgb_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def add_rect(slide, l, t, w, h, rgb):
    s = slide.shapes.add_shape(1, l, t, w, h)
    rgb_fill(s, rgb)
    return s


def add_tb(slide, l, t, w, h, text, size=12, bold=False,
           color=DARK_TEXT, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = FONT
    return tb


def add_tb_multiline(slide, l, t, w, h, lines, size=10,
                     bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT):
    """여러 줄 텍스트박스"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line:
            run = p.add_run()
            run.text = line
            run.font.size  = Pt(size)
            run.font.bold  = bold
            run.font.color.rgb = color
            run.font.name  = FONT
    return tb


def remove_content(slide):
    to_remove = []
    for shape in slide.shapes:
        l = shape.left / 914400
        t = shape.top  / 914400
        w = shape.width / 914400
        if l < 0.02 and w < 0.20 and t < 0.5:
            continue
        if t > 6.90:
            continue
        to_remove.append(shape._element)
    for el in to_remove:
        el.getparent().remove(el)


def common_header(slide, section_label, title, page_num):
    """공통 헤더: 녹색 세로바 + 섹션라벨 + 타이틀 + 구분선"""
    add_rect(slide, 0, 0, Inches(0.14), SLIDE_H, GREEN)
    add_tb(slide, Inches(0.47), Inches(0.24), Inches(7.87), Inches(0.28),
           section_label, size=10, color=GREEN)
    add_tb(slide, Inches(0.47), Inches(0.48), Inches(9.06), Inches(0.54),
           title, size=22, bold=True, color=DARK_TEXT)
    add_rect(slide, Inches(0.47), Inches(1.07), Inches(9.06), Inches(0.02), GREEN)
    add_rect(slide, 0, Inches(7.05), SLIDE_W, Inches(0.45), GREEN)
    add_tb(slide, Inches(1.57), Inches(7.11), Inches(6.30), Inches(0.33),
           'Copyright ⓒ 2025 Hana Financial Group. All rights reserved.',
           size=6, color=WHITE)
    add_tb(slide, Inches(8.98), Inches(7.09), Inches(0.91), Inches(0.37),
           f'| {page_num}', size=8, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════
#  Slide 19 — 섹션 디바이더 업데이트
# ══════════════════════════════════════════════════════════
def update_slide19():
    slide = prs.slides[18]
    shapes = list(slide.shapes)
    # subtitle 텍스트 교체 (shape[6])
    for shape in shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if '보고용 집계표' in t or '숨겨진 오류' in t:
                    for run in para.runs:
                        run.text = ''
                    if para.runs:
                        para.runs[0].text = 'AI와 함께 직원별 거래 현황을 집계하고, 정상으로 표시된 오류 데이터를 찾아낸다'
                    print('  [Slide 19] 섹션 디바이더 서브타이틀 수정')
                if '| 12' in t or '| 19' in t:
                    for run in para.runs:
                        if '12' in run.text:
                            run.text = '| 19'
                        elif '19' in run.text:
                            pass  # already correct


# ══════════════════════════════════════════════════════════
#  Slide 20 — STEP1 직원별 집계표 (VLOOKUP + COUNTIFS)
# ══════════════════════════════════════════════════════════
def rebuild_slide20():
    slide = prs.slides[19]
    remove_content(slide)
    common_header(slide,
                  '데이터 분석 활용(EXCEL)',
                  'STEP 1 · 직원별 거래 현황 집계  (10~25분)',
                  20)

    # ── 도입 배너
    add_rect(slide, Inches(0.47), Inches(1.20), Inches(9.06), Inches(0.38), DARK_GREEN)
    add_rect(slide, Inches(0.47), Inches(1.20), Inches(0.06), Inches(0.38), GREEN)
    add_tb(slide, Inches(0.62), Inches(1.25), Inches(8.80), Inches(0.28),
           '직원마스터에서 직원명·지점코드를 조회하고, 거래원장_RAW에서 처리상태=정상 기준으로 집계합니다.',
           size=11, color=WHITE)

    # ── 좌측: RAW 시트 컬럼 구조
    lx = Inches(0.47)
    col_t = Inches(1.72)
    col_w = Inches(4.35)

    # 컬럼 헤더 배너
    add_rect(slide, lx, col_t, col_w, Inches(0.40), GREEN)
    add_tb(slide, lx + Inches(0.10), col_t + Inches(0.05), col_w - Inches(0.18), Inches(0.30),
           '거래원장_RAW 시트 구조 (A~J열)', size=12, bold=True, color=WHITE)

    cols = [
        ('A', '거래일자'),
        ('B', '지점코드'),
        ('C', '직원ID'),
        ('D', '고객ID'),
        ('E', '상품유형'),
        ('F', '거래유형'),
        ('G', '거래금액'),
        ('H', '수수료율'),
        ('I', '수수료금액'),
        ('J', '처리상태  ← 핵심'),
    ]
    row_h = Inches(0.27)
    for i, (col, name) in enumerate(cols):
        ty = col_t + Inches(0.40) + i * row_h
        bg_color = LIGHT_GREEN if i % 2 == 0 else LIGHT_GRAY
        add_rect(slide, lx, ty, col_w, row_h, bg_color)
        add_rect(slide, lx, ty, Inches(0.36), row_h, DARK_GREEN if col == 'J' else GREEN)
        add_tb(slide, lx, ty, Inches(0.36), row_h,
               col, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        fc = DARK_GREEN if col == 'J' else DARK_TEXT
        add_tb(slide, lx + Inches(0.42), ty + Inches(0.03),
               col_w - Inches(0.50), row_h - Inches(0.04),
               name, size=11, color=fc, bold=(col == 'J'))

    # ── 우측: 결과 시트 구조 + AI 프롬프트
    rx = Inches(4.98)
    rw = Inches(4.58)

    # 결과 시트 헤더
    add_rect(slide, rx, col_t, rw, Inches(0.40), GREEN)
    add_tb(slide, rx + Inches(0.10), col_t + Inches(0.05), rw - Inches(0.18), Inches(0.30),
           '결과작성_STEP1 시트 — 직원별 집계표', size=12, bold=True, color=WHITE)

    # 결과 테이블 컬럼 헤더
    result_cols = [
        ('직원ID', '예: E001'),
        ('직원명', 'VLOOKUP ← 직원마스터'),
        ('지점코드', 'VLOOKUP ← 직원마스터'),
        ('정상 거래건수', 'COUNTIFS 조건 2개'),
        ('정상 거래금액', 'SUMIFS 조건 2개'),
        ('정상 수수료합계', 'SUMIFS 조건 2개'),
        ('취소·오류 건수', 'COUNTIFS 조건 2개'),
    ]
    row_h2 = Inches(0.34)
    for i, (col_nm, hint) in enumerate(result_cols):
        ty = col_t + Inches(0.40) + i * row_h2
        bg_color = LIGHT_GREEN if i % 2 == 0 else LIGHT_GRAY
        add_rect(slide, rx, ty, rw, row_h2, bg_color)
        add_rect(slide, rx, ty, Inches(0.06), row_h2, GREEN)
        add_tb(slide, rx + Inches(0.12), ty + Inches(0.03),
               Inches(1.60), row_h2 - Inches(0.04),
               col_nm, size=10, bold=True, color=DARK_GREEN)
        add_tb(slide, rx + Inches(1.78), ty + Inches(0.03),
               rw - Inches(1.86), row_h2 - Inches(0.04),
               hint, size=9.5, color=DARK_TEXT, italic=True)

    # AI 프롬프트 예시 박스
    pr_t = col_t + Inches(0.40) + len(result_cols) * row_h2 + Inches(0.10)
    add_rect(slide, rx, pr_t, rw, Inches(0.34), DARK_GREEN)
    add_tb(slide, rx + Inches(0.10), pr_t + Inches(0.06), rw - Inches(0.18), Inches(0.24),
           'AI 프롬프트 힌트', size=11, bold=True, color=WHITE)

    pr_body_t = pr_t + Inches(0.34)
    pr_body_h = Inches(6.95) - pr_body_t/914400 * 914400   # remaining space
    add_rect(slide, rx, pr_body_t, rw, Inches(1.00), LIGHT_GRAY)
    add_rect(slide, rx, pr_body_t, Inches(0.06), Inches(1.00), GREEN)
    prompt_lines = [
        '"A5셀(E001)의 직원명을 직원마스터 시트에서',
        ' 찾아서 B5에 넣는 VLOOKUP 수식 알려줘"',
        '"거래원장_RAW에서 C열=A5, J열=정상인',
        ' 거래건수를 세는 COUNTIFS 수식 알려줘"',
        '"위 방식으로 D5:G5 수식을 한꺼번에 만들어줘"',
    ]
    add_tb_multiline(slide, rx + Inches(0.14), pr_body_t + Inches(0.08),
                     rw - Inches(0.22), Inches(0.88),
                     prompt_lines, size=9.5, color=DARK_GREEN)

    print('  [Slide 20] STEP1 직원별 집계표 재구성 완료')


# ══════════════════════════════════════════════════════════
#  Slide 21 — STEP2 2개 오류 조건
# ══════════════════════════════════════════════════════════
def rebuild_slide21():
    slide = prs.slides[20]
    remove_content(slide)
    common_header(slide,
                  '데이터 분석 활용(EXCEL)',
                  'STEP 2 · 정상으로 표시된 오류 데이터 찾기  (25~42분)',
                  21)

    # ── 도입 설명 배너
    add_rect(slide, Inches(0.47), Inches(1.20), Inches(9.06), Inches(0.42), DARK_GREEN)
    add_rect(slide, Inches(0.47), Inches(1.20), Inches(0.06), Inches(0.42), GREEN)
    add_tb(slide, Inches(0.62), Inches(1.26), Inches(8.80), Inches(0.32),
           '처리상태 = 정상 인데도 실제 데이터값에 오류가 있는 행을 찾아냅니다. '
           'AI로 K열(오류유형_작성) 수식을 작성합니다.',
           size=11, color=WHITE)

    # ── 오류 조건 2개 (큰 카드)
    card_w = Inches(4.40)
    card_h = Inches(2.20)
    card_t = Inches(1.75)
    gap    = Inches(0.22)

    cards = [
        ('01', '거래금액 오류',
         '처리상태 = 정상\nAND 거래금액(G열) ≤ 0',
         '거래가 정상으로 등록됐는데\n거래금액이 0 또는 음수인 경우.\n\n'
         '허용값: 0 초과\n실습 데이터에 여러 건 숨겨져 있습니다.'),
        ('02', '상품유형 오류',
         '처리상태 = 정상\nAND 상품유형이 허용값 외',
         '허용 상품유형: 주식 / 펀드 / 채권 / ELS / RP\n\n'
         '실습 데이터에는 FX, FWD, ETF 등\n허용되지 않은 상품유형이 포함됩니다.'),
    ]

    for i, (num, title, cond, desc) in enumerate(cards):
        lx = Inches(0.47) + i * (card_w + gap)
        # 카드 배경
        add_rect(slide, lx, card_t, card_w, card_h, LIGHT_GRAY)
        add_rect(slide, lx, card_t, card_w, Inches(0.06), GREEN)
        # 번호 배지
        add_rect(slide, lx + Inches(0.12), card_t + Inches(0.14),
                 Inches(0.52), Inches(0.52), GREEN)
        add_tb(slide, lx + Inches(0.12), card_t + Inches(0.14),
               Inches(0.52), Inches(0.52),
               num, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # 제목
        add_tb(slide, lx + Inches(0.74), card_t + Inches(0.16),
               card_w - Inches(0.88), Inches(0.42),
               title, size=16, bold=True, color=DARK_GREEN)
        # 조건 박스
        add_rect(slide, lx + Inches(0.12), card_t + Inches(0.72),
                 card_w - Inches(0.22), Inches(0.62), DARK_GREEN)
        add_tb_multiline(slide, lx + Inches(0.22), card_t + Inches(0.78),
                         card_w - Inches(0.40), Inches(0.52),
                         cond.split('\n'), size=11, bold=True, color=WHITE)
        # 설명
        add_tb_multiline(slide, lx + Inches(0.22), card_t + Inches(1.42),
                         card_w - Inches(0.34), Inches(0.70),
                         desc.split('\n'), size=10, color=DARK_TEXT)

    # ── K열 작성 안내
    k_t = card_t + card_h + Inches(0.18)
    add_rect(slide, Inches(0.47), k_t, Inches(9.06), Inches(0.36), GREEN)
    add_tb(slide, Inches(0.60), k_t + Inches(0.06), Inches(8.80), Inches(0.26),
           'K열(오류유형_작성): 오류가 없으면 빈칸, 여러 조건이 겹치면 대표 오류 하나만 표시',
           size=11, bold=True, color=WHITE)

    # ── AI 프롬프트 3개 (가로 배치)
    pr_t  = k_t + Inches(0.46)
    pr_w  = Inches(2.90)
    pr_h  = Inches(1.62)
    pr_gap = Inches(0.18)

    prompts = [
        ('K열 — 거래금액 오류',
         ['"거래원장_RAW에서 J열(처리상태)이',
          ' 정상인데 G열(거래금액)이 0 이하인',
          ' 행을 찾아 K열에',
          ' \'거래금액 오류\'를 표시하는',
          ' 수식 알려줘"']),
        ('K열 — 상품유형 오류',
         ['"처리상태=정상인데 E열(상품유형)이',
          ' 주식/펀드/채권/ELS/RP가 아닌 행을',
          ' 찾아 K열에 \'상품유형 오류\'를',
          ' 표시하는',
          ' IF+ISNUMBER+MATCH 수식 알려줘"']),
        ('K열 — 통합 수식',
         ['"위 두 조건을 합쳐서',
          ' K열에 한꺼번에 표시하는',
          ' 수식 만들어줘.',
          ' (거래금액 오류가 먼저,',
          '  상품유형 오류가 나중)"']),
    ]

    for i, (ph_title, ph_lines) in enumerate(prompts):
        px = Inches(0.47) + i * (pr_w + pr_gap)
        add_rect(slide, px, pr_t, pr_w, pr_h, LIGHT_GRAY)
        add_rect(slide, px, pr_t, pr_w, Inches(0.34), DARK_GREEN)
        add_tb(slide, px + Inches(0.10), pr_t + Inches(0.06),
               pr_w - Inches(0.18), Inches(0.24),
               ph_title, size=10.5, bold=True, color=WHITE)
        add_tb_multiline(slide, px + Inches(0.10), pr_t + Inches(0.40),
                         pr_w - Inches(0.18), pr_h - Inches(0.46),
                         ph_lines, size=9.5, color=DARK_GREEN)

    # ── 검증 안내 배너
    vf_t = pr_t + pr_h + Inches(0.10)
    add_rect(slide, Inches(0.47), vf_t, Inches(9.06), Inches(0.34), LIGHT_GREEN)
    add_rect(slide, Inches(0.47), vf_t, Inches(0.06), Inches(0.34), GREEN)
    add_tb(slide, Inches(0.62), vf_t + Inches(0.06), Inches(8.80), Inches(0.26),
           '검증: 오류로 잡힌 행을 직접 눈으로 확인하세요. AI 수식이 정상 처리건을 K열에 표시하지 않는지 반드시 체크.',
           size=10, color=DARK_GREEN, italic=True)

    print('  [Slide 21] STEP2 2개 오류 조건으로 재구성 완료')


# ══════════════════════════════════════════════════════════
if __name__ == '__main__':
    update_slide19()
    rebuild_slide20()
    rebuild_slide21()
    prs.save(DST)
    print(f'\n✓ 저장: {DST}')
