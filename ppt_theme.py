"""
ppt_theme.py — 하나증권 공통 PPT 디자인 토큰 & 헬퍼

색상 기준 (웹앱 hana 테마와 동일):
  PRIMARY  = #009591  하나그린  — 헤더바·섹션배경·표헤더·불릿
  ACCENT   = #ED1651  하나레드  — 강조선·배지·포인트
  NAVY     = #1E4D9B  네이비    — 보조 강조 (분석·전략 맥락)
  DARK     = #231F20  다크       — 기본 본문 텍스트
  GRAY     = #666666  회색       — 보조 텍스트·각주
  LGRAY    = #F0FAFA  연한 민트  — 카드 배경
  WHITE    = #FFFFFF

폰트 기준:
  제목: 하나2.0 B  (없으면 맑은 고딕)
  본문: 하나2.0 M  (없으면 맑은 고딕)
"""

import sys, io, os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

# ── 색상 팔레트 (하나증권 CI) ───────────────────────────────────────────────────
PRIMARY  = RGBColor(0x00, 0x95, 0x91)   # 하나그린
ACCENT   = RGBColor(0xED, 0x16, 0x51)   # 하나레드
NAVY     = RGBColor(0x1E, 0x4D, 0x9B)   # 네이비
DARK_BG  = RGBColor(0x00, 0x7A, 0x77)   # 다크 그린 (섹션 배경)
DARK     = RGBColor(0x23, 0x1F, 0x20)   # 기본 텍스트
GRAY     = RGBColor(0x66, 0x66, 0x66)   # 보조 텍스트
MGRAY    = RGBColor(0x9C, 0xA3, 0xAF)   # 중간 회색
LGRAY    = RGBColor(0xF0, 0xFA, 0xFA)   # 연한 민트 (카드 배경)
LGRAY2   = RGBColor(0xE6, 0xF7, 0xF7)   # 표 짝수 행
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

# ── 폰트 ───────────────────────────────────────────────────────────────────────
FONT_TITLE = "하나2.0 B"
FONT_BODY  = "하나2.0 M"

# ── 슬라이드 규격 (와이드 16:9) ────────────────────────────────────────────────
SLIDE_W = Inches(13.33)   # 10" pptxgenjs → 13.33" python-pptx (widescreen)
SLIDE_H = Inches(7.50)

# ── 로고 경로 ──────────────────────────────────────────────────────────────────
_DIR = os.path.dirname(os.path.abspath(__file__))
LOGO_PATH = os.path.join(_DIR, "public", "hana-mark.png")


# ══════════════════════════════════════════════════════════════════════════════════
# 새 Presentation 객체 생성 헬퍼
# ══════════════════════════════════════════════════════════════════════════════════

def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    """빈 레이아웃(6번) 슬라이드 추가"""
    return prs.slides.add_slide(prs.slide_layouts[6])


# ══════════════════════════════════════════════════════════════════════════════════
# 도형 헬퍼
# ══════════════════════════════════════════════════════════════════════════════════

def add_rect(slide, left, top, width, height,
             fill: RGBColor | None = None,
             line: RGBColor | None = None,
             line_width: Pt = Pt(0)) -> object:
    """직사각형 추가. fill=None → 투명, line=None → 선 없음"""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height,
                text: str,
                font_size: Pt = Pt(11),
                bold: bool = False,
                italic: bool = False,
                color: RGBColor = DARK,
                align: PP_ALIGN = PP_ALIGN.LEFT,
                font_face: str | None = None,
                wrap: bool = True,
                line_spacing: float | None = None) -> object:
    """텍스트박스 추가"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        from pptx.util import Pt as _Pt
        p.line_spacing = _Pt(line_spacing)
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if font_face:
        run.font.name = font_face
    return txBox


def add_multi_para_textbox(slide, left, top, width, height, paras: list[dict],
                            wrap: bool = True) -> object:
    """여러 단락을 가진 텍스트박스. paras = [{text, font_size, bold, color, align, italic, space_before, space_after, font_face}]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    for i, p_info in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p_info.get("align", PP_ALIGN.LEFT)
        if p_info.get("space_before"):
            para.space_before = p_info["space_before"]
        if p_info.get("space_after"):
            para.space_after = p_info["space_after"]
        run = para.add_run()
        run.text = p_info.get("text", "")
        run.font.size = p_info.get("font_size", Pt(11))
        run.font.bold = p_info.get("bold", False)
        run.font.italic = p_info.get("italic", False)
        run.font.color.rgb = p_info.get("color", DARK)
        if p_info.get("font_face"):
            run.font.name = p_info["font_face"]
    return txBox


# ══════════════════════════════════════════════════════════════════════════════════
# 공통 크롬 (헤더 바 / 푸터)
# ══════════════════════════════════════════════════════════════════════════════════

def add_header(slide, title: str, subtitle: str = "", date: str = "",
               compact: bool = False) -> float:
    """
    하나증권 공통 헤더 바 추가.
    compact=True → 임원보고 원페이지용 좁은 헤더 (h=0.68")
    compact=False → 일반 슬라이드 헤더 (h=0.85")
    반환값: 헤더 높이(inches, float)
    """
    h = Inches(0.68) if compact else Inches(0.85)
    H = 0.68 if compact else 0.85

    # 헤더 배경
    add_rect(slide, 0, 0, SLIDE_W, h, fill=PRIMARY)
    # 좌측 레드 액센트 스트라이프
    add_rect(slide, 0, 0, Inches(0.14), h, fill=ACCENT)
    # 하단 레드 액센트 선
    add_rect(slide, 0, Inches(H - 0.05), SLIDE_W, Inches(0.05), fill=ACCENT)

    # 타이틀
    title_fs = Pt(16) if compact else Pt(18)
    add_textbox(slide, Inches(0.22), 0, Inches(10), h,
                title, font_size=title_fs, bold=True, color=WHITE,
                font_face=FONT_TITLE, align=PP_ALIGN.LEFT)

    # 부제목 / 날짜 (우측)
    if subtitle:
        add_textbox(slide, Inches(10.5), Inches(0.05), Inches(2.6), Inches(H * 0.45),
                    subtitle, font_size=Pt(9.5), color=WHITE, font_face=FONT_BODY,
                    align=PP_ALIGN.RIGHT)
    if date:
        add_textbox(slide, Inches(10.5), Inches(H * 0.5), Inches(2.6), Inches(H * 0.42),
                    date, font_size=Pt(9.5), color=ACCENT, font_face=FONT_BODY,
                    align=PP_ALIGN.RIGHT)

    # 로고 (우측 상단, 헤더 안)
    if os.path.exists(LOGO_PATH):
        try:
            slide.shapes.add_picture(LOGO_PATH, Inches(12.9), Inches(0.08),
                                     height=Inches(H - 0.16))
        except Exception:
            pass  # 로고 파일 오류 시 조용히 넘김

    return H


def add_footer(slide, left_text: str = "", right_text: str = "") -> None:
    """슬라이드 하단 풋터 바"""
    footer_y = Inches(7.1)
    footer_h = Inches(0.4)
    add_rect(slide, 0, footer_y, SLIDE_W, footer_h, fill=DARK)
    if left_text:
        add_textbox(slide, Inches(0.3), footer_y, Inches(8), footer_h,
                    left_text, font_size=Pt(9), color=MGRAY, font_face=FONT_BODY)
    if right_text:
        add_textbox(slide, Inches(9), footer_y, Inches(4.1), footer_h,
                    right_text, font_size=Pt(9), color=MGRAY, font_face=FONT_BODY,
                    align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════════
# 표 헬퍼
# ══════════════════════════════════════════════════════════════════════════════════

def add_table(slide, left, top, width, height,
              headers: list[str], rows: list[list[str]],
              header_fs: Pt = Pt(10), cell_fs: Pt = Pt(9.5)) -> object:
    """
    하나증권 CI 스타일 표 추가.
    헤더: 하나그린 배경 + 흰 텍스트
    짝수 행: 연한 민트(LGRAY2), 홀수 행: 흰색
    """
    cols = len(headers)
    data_rows = len(rows)
    tbl = slide.shapes.add_table(data_rows + 1, cols, left, top, width, height).table

    # 열 너비 균등 분배
    col_w = int(width / cols)
    for c in range(cols):
        tbl.columns[c].width = col_w

    def _cell(r, c, text, bg: RGBColor, fg: RGBColor, bold: bool, fs: Pt, align=PP_ALIGN.CENTER):
        cell = tbl.cell(r, c)
        cell.text = text
        para = cell.text_frame.paragraphs[0]
        para.alignment = align
        run = para.runs[0] if para.runs else para.add_run()
        run.font.size = fs
        run.font.bold = bold
        run.font.color.rgb = fg
        if para.runs:
            para.runs[0].font.name = FONT_TITLE if bold else FONT_BODY
        from pptx.oxml.ns import qn
        from lxml import etree
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
        srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", str(bg).upper())

    # 헤더 행
    for c, h in enumerate(headers):
        _cell(0, c, h, PRIMARY, WHITE, True, header_fs)

    # 데이터 행
    for r, row in enumerate(rows):
        bg = LGRAY2 if r % 2 == 1 else WHITE
        for c, val in enumerate(row):
            _cell(r + 1, c, val, bg, DARK, False, cell_fs,
                  align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)

    return tbl
