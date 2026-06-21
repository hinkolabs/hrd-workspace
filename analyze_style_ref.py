import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

SRC = r'C:\Users\jay\Downloads\AI교육_완성판_v8.pptx'
prs = Presentation(SRC)

def show(idx):
    slide = prs.slides[idx]
    print(f'\n=== Slide {idx+1} ({len(slide.shapes)} shapes) ===')
    for i, shape in enumerate(slide.shapes):
        L = shape.left/914400
        T = shape.top/914400
        W = shape.width/914400
        H = shape.height/914400
        fill_rgb = None
        try:
            fill_rgb = shape.fill.fore_color.rgb
        except: pass
        txt = ''
        if shape.has_text_frame:
            lines = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
            if lines:
                sz = None
                try: sz = shape.text_frame.paragraphs[0].runs[0].font.size
                except: pass
                txt = f' | sz={sz and sz//12700} | ' + ' / '.join(lines)[:80]
        rgb_str = f' fill={fill_rgb}' if fill_rgb else ''
        print(f'  [{i:2d}] T={T:.2f}" L={L:.2f}" W={W:.2f}" H={H:.2f}"{rgb_str}{txt}')

# 스타일 참고: 실제 잘 구성된 콘텐츠 슬라이드들
for idx in [13, 14, 15]:
    show(idx)
