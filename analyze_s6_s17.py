import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt

SRC = r'C:\Users\jay\Downloads\AI교육_완성판_v5.pptx'
prs = Presentation(SRC)

def analyze_slide(idx):
    slide = prs.slides[idx]
    print(f'\n=== Slide {idx+1} ({len(slide.shapes)} shapes) ===')
    for i, shape in enumerate(slide.shapes):
        print(f'  [{i}] type={shape.shape_type} name="{shape.name}"')
        print(f'       pos=L{shape.left/914400:.2f}" T{shape.top/914400:.2f}" W{shape.width/914400:.2f}" H{shape.height/914400:.2f}"')
        if shape.has_text_frame:
            for j, para in enumerate(shape.text_frame.paragraphs):
                txt = para.text.strip()
                if txt:
                    print(f'       para[{j}]: "{txt}"')

analyze_slide(5)   # slide 6
analyze_slide(16)  # slide 17
