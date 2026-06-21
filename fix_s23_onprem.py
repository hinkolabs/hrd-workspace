"""
Slide 23 — 온프레미스(On-Premise) 환경에 맞게 내용 전면 수정
하나증권 도입 환경: Alli Works 온프렘, 외부 클라우드 미사용
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SRC = r'C:\Users\jay\Downloads\AI교육_완성판_v9.pptx'
DST = r'C:\Users\jay\Downloads\AI교육_완성판_v10.pptx'
prs = Presentation(SRC)

GREEN       = RGBColor(0x00, 0x91, 0x78)
DARK_GREEN  = RGBColor(0x00, 0x4E, 0x42)
LIGHT_GREEN = RGBColor(0xD7, 0xED, 0xE6)
LIGHT_GRAY  = RGBColor(0xF7, 0xF9, 0xF8)
DARK_TEXT   = RGBColor(0x22, 0x22, 0x22)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY        = RGBColor(0xCC, 0xCC, 0xCC)
FONT        = '맑은 고딕'


def rgb_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def add_rect(slide, l, t, w, h, rgb):
    s = slide.shapes.add_shape(1, l, t, w, h)
    rgb_fill(s, rgb)
    return s


def add_tb(slide, l, t, w, h, text, size=12, bold=False,
           color=DARK_TEXT, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = FONT
    return tb


def add_card(slide, l, t, w, h, num_str, title, body, num_size=18):
    add_rect(slide, l, t, w, h, LIGHT_GRAY)
    add_rect(slide, l, t, Inches(0.65), h, GREEN)
    add_tb(slide, l, t, Inches(0.65), h,
           num_str, size=num_size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    inner_l = l + Inches(0.72)
    inner_w = w - Inches(0.80)
    add_tb(slide, inner_l, t + Inches(0.10), inner_w, Inches(0.38),
           title, size=13, bold=True, color=DARK_GREEN)
    add_rect(slide, inner_l, t + Inches(0.52), inner_w, Inches(0.01), GRAY)
    add_tb(slide, inner_l, t + Inches(0.58), inner_w, h - Inches(0.66),
           body, size=10, color=DARK_TEXT, wrap=True)


def remove_all_content(slide):
    to_remove = []
    for shape in slide.shapes:
        l = shape.left / 914400
        t = shape.top / 914400
        w = shape.width / 914400
        if l < 0.02 and w < 0.20 and t < 0.5:
            continue
        if t > 6.90:
            continue
        to_remove.append(shape._element)
    for el in to_remove:
        el.getparent().remove(el)


def rebuild_slide23_onprem():
    slide = prs.slides[22]
    remove_all_content(slide)

    # ── 헤더
    add_tb(slide, Inches(0.47), Inches(0.24), Inches(7.87), Inches(0.28),
           '2. 실습 · 나만의 챗봇 업무 비서 만들기', size=10, color=GREEN)
    add_tb(slide, Inches(0.47), Inches(0.48), Inches(9.06), Inches(0.54),
           'Alli Works & RAG — 기본 개념 이해',
           size=24, bold=True, color=DARK_TEXT)
    add_rect(slide, Inches(0.47), Inches(1.07), Inches(9.06), Inches(0.02), GREEN)

    # ── 온프렘 환경 안내 배너
    add_rect(slide, Inches(0.47), Inches(1.20), Inches(9.06), Inches(0.40), DARK_GREEN)
    add_rect(slide, Inches(0.47), Inches(1.20), Inches(0.06), Inches(0.40), GREEN)
    add_tb(slide, Inches(0.62), Inches(1.25), Inches(8.80), Inches(0.30),
           '하나증권 Alli Works는 사내 서버에서 직접 운영되는 온프레미스(On-Premise) 환경입니다. '
           '모든 데이터는 사내 망 안에서만 처리됩니다.',
           size=11, color=WHITE)

    # ── 2×2 카드 (온프렘 맞춤)
    card_w = Inches(4.40)
    card_h = Inches(1.55)
    gap    = Inches(0.18)
    row1_t = Inches(1.73)
    row2_t = row1_t + card_h + gap
    col1_l = Inches(0.47)
    col2_l = col1_l + card_w + gap

    cards = [
        (col1_l, row1_t, '01', 'Alli Works — 온프레미스 AI 플랫폼',
         '사내 서버에 직접 설치·운영되는 AI 챗봇 제작 플랫폼.\n'
         '외부 클라우드(ChatGPT·Gemini 등)와 달리 데이터가 회사 네트워크 밖으로 나가지 않습니다.'),
        (col2_l, row1_t, '02', '온프렘의 핵심 장점',
         '• 데이터 보안: MNPI·고객 정보가 외부 서버에 전송되지 않음\n'
         '• 내부 통제: IT·컴플라이언스 기준 충족\n'
         '• 인터넷 없이도 사내망에서 안정적 운영'),
        (col1_l, row2_t, '03', 'RAG(검색 증강 생성)란?',
         'AI가 질문을 받을 때마다 사내 서버에 저장된 문서를 실시간으로 검색해\n'
         '관련 단락을 찾아 답변에 활용하는 기술.\n'
         '"미리 외운 답" 대신 "그때그때 찾은 답"을 제공합니다.'),
        (col2_l, row2_t, '04', 'RAG 동작 흐름',
         '① 질문 입력  →  ② 사내 서버 문서 검색 (관련 단락 추출)\n'
         '→  ③ 질문 + 단락을 온프렘 AI 모델에 전달\n'
         '→  ④ 문서 기반 답변 생성 + 출처 제시'),
    ]
    for lx, ty, num, title, body in cards:
        add_card(slide, lx, ty, card_w, card_h, num, title, body)

    # ── 핵심 메모 배너
    note_t = row2_t + card_h + Inches(0.15)
    add_rect(slide, Inches(0.47), note_t, Inches(9.06), Inches(0.36), LIGHT_GREEN)
    add_rect(slide, Inches(0.47), note_t, Inches(0.06), Inches(0.36), GREEN)
    add_tb(slide, Inches(0.62), note_t + Inches(0.06), Inches(8.80), Inches(0.28),
           '핵심: Alli Works는 AI가 사내 문서를 "암기"하는 것이 아닙니다. '
           '질문마다 사내 서버에서 검색해 그때그때 참고하는 구조입니다.',
           size=10, color=DARK_GREEN, italic=True)

    # ── 하단 3개 핵심 기능 카드
    feat_t   = note_t + Inches(0.46)
    feat_h   = Inches(1.55)
    feat_w   = Inches(2.90)
    feat_gap = Inches(0.18)

    feats = [
        ('데이터 사내 완결 처리',
         '업로드된 문서와 AI 답변 모두 사내 서버 내에서만 처리·저장.\n외부 API 호출 없이 완전 폐쇄망 운영.'),
        ('업무 맞춤 챗봇 제작',
         '부서 FAQ·업무 프로세스·사규 문서를 업로드하고\n시스템 프롬프트로 어조·답변 범위 설정.'),
        ('보안 및 접근 권한 관리',
         '부서별·직급별 접근 권한 설정 가능.\nMNPI·개인정보가 포함된 자료는 업로드 금지.'),
    ]
    for i, (ftitle, fbody) in enumerate(feats):
        fx = Inches(0.47) + i * (feat_w + feat_gap)
        add_rect(slide, fx, feat_t, feat_w, feat_h, LIGHT_GRAY)
        add_rect(slide, fx, feat_t, feat_w, Inches(0.04), GREEN)
        add_tb(slide, fx + Inches(0.10), feat_t + Inches(0.08),
               feat_w - Inches(0.18), Inches(0.34),
               ftitle, size=11, bold=True, color=DARK_GREEN)
        add_rect(slide, fx + Inches(0.10), feat_t + Inches(0.46),
                 feat_w - Inches(0.18), Inches(0.01), GRAY)
        add_tb(slide, fx + Inches(0.10), feat_t + Inches(0.52),
               feat_w - Inches(0.18), feat_h - Inches(0.58),
               fbody, size=9.5, color=DARK_TEXT, wrap=True)

    print('  [Slide 23] 온프렘 환경 기준으로 내용 재구성 완료')


rebuild_slide23_onprem()
prs.save(DST)
print(f'\n✓ 저장: {DST}')
