# -*- coding: utf-8 -*-
"""
build_1pptx.py  —  하나증권 CI 단일 디자인 시스템 v3
=======================================================
모든 슬라이드가 동일한 58pt 레드 헤더 + 흰 본문 구조를 공유합니다.
표지·간지의 대제목도 흰 본문 영역에 배치하여 시각 통일.
Output: C:/Users/jay/Downloads/1_rebuilt.pptx
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── 슬라이드 규격 ────────────────────────────────────────
W = Pt(960)
H = Pt(540)

# ─── 하나증권 CI 색상 ──────────────────────────────────────
RED   = RGBColor(0xED, 0x16, 0x51)
DARK  = RGBColor(0x23, 0x1F, 0x20)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xDD, 0xDD, 0xDD)
MGRAY = RGBColor(0x88, 0x88, 0x88)
BGRAY = RGBColor(0xF5, 0xF5, 0xF5)

# ─── 폰트 ─────────────────────────────────────────────────
FONT_B = "하나2.0 B"
FONT_M = "하나2.0 M"

# ─── 고정 레이아웃 상수 ────────────────────────────────────
HDR_H   = 58    # 전 슬라이드 공통 헤더 높이 (pt)
BODY_Y  = HDR_H # 본문 시작 Y
BODY_H  = 540 - HDR_H
NAV_X   = 906   # 우측 섹션 내비 X
FTR_Y   = 505   # 하단 구분선 Y


# ═══════════════════════════════════════════════════════════
# 헬퍼 유틸리티
# ═══════════════════════════════════════════════════════════

def px(pt: float) -> int:
    return int(pt * 12700)


def rect(slide, x, y, w, h, fill=None):
    s = slide.shapes.add_shape(1, px(x), px(y), px(w), px(h))
    s.fill.solid() if fill else s.fill.background()
    if fill:
        s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def txt(slide, x, y, w, h, text,
        font=FONT_M, size=12, bold=False,
        color=DARK, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    _font(r, font, size, bold, color)
    tb.line.fill.background()
    return tb


def _font(run, name, size, bold, color):
    f = run.font
    f.name = name
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color


def _add_para(tf, text, font=FONT_M, size=12, bold=False,
              color=DARK, align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    r = p.add_run()
    r.text = text
    _font(r, font, size, bold, color)
    return p


def hline(slide, x, y, w, color=LGRAY):
    return rect(slide, x, y, w, 0.6, fill=color)


def badge(slide, x, y, text, size=22, bg=RED, fg=WHITE):
    rect(slide, x, y, size, size, fill=bg)
    txt(slide, x, y + 2, size, size - 4,
        text, FONT_B, size * 0.58, True, fg, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# 공통 크롬 (전 슬라이드 동일)
# ═══════════════════════════════════════════════════════════

def draw_chrome(slide, header_label, page_num=None,
                sections=None, footnote="(2026. 04 기준)"):
    """
    58pt 레드 헤더 + 좌측 4pt 다크 강조선 + 흰 본문 배경
    + 우측 섹션 내비 (sections 있을 때) + 하단 푸터
    모든 슬라이드가 이 함수 하나만 사용합니다.
    """
    # ── 헤더 (58pt, 전 슬라이드 동일) ──────────────────────
    rect(slide, 0, 0, 960, HDR_H, fill=RED)
    rect(slide, 0, 0, 6, HDR_H, fill=DARK)           # 좌측 강조선
    txt(slide, 18, 0, 840, HDR_H,
        header_label, FONT_B, 20, True, WHITE, PP_ALIGN.LEFT)

    # ── 흰 본문 배경 ─────────────────────────────────────────
    rect(slide, 0, BODY_Y, 960, BODY_H, fill=WHITE)

    # ── 우측 섹션 내비 ───────────────────────────────────────
    if sections:
        each_h = 420 / len(sections)
        y0 = BODY_Y + 10
        for i, (label, active) in enumerate(sections):
            col = RED if active else LGRAY
            rect(slide, NAV_X, y0 + i * each_h, 48, each_h - 3, fill=col)
            txt(slide, NAV_X - 1, y0 + i * each_h + 3, 50, each_h - 6,
                label, FONT_M, 6.5, active,
                WHITE if active else DARK, PP_ALIGN.CENTER)

    # ── 하단 푸터 ────────────────────────────────────────────
    hline(slide, 18, FTR_Y, 870)
    if footnote:
        txt(slide, 18, FTR_Y + 5, 500, 20,
            footnote, FONT_M, 8, False, MGRAY)
    if page_num:
        txt(slide, 860, FTR_Y + 5, 80, 20,
            f"{page_num} / 11", FONT_M, 8, False, MGRAY, PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════
# 카드 컴포넌트
# ═══════════════════════════════════════════════════════════

def draw_card(slide, x, y, w, h, num, title, bullets):
    """3단 카드"""
    rect(slide, x, y, w, h, fill=BGRAY)
    rect(slide, x, y, w, 3, fill=RED)          # 카드 상단 강조선
    badge(slide, x + 12, y + 12, num, size=22)
    txt(slide, x + 40, y + 13, w - 52, 26,
        title, FONT_B, 13, True, DARK)
    hline(slide, x + 12, y + 42, w - 24)
    tb = slide.shapes.add_textbox(px(x + 12), px(y + 48),
                                   px(w - 24), px(h - 60))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for bul in bullets:
        is_sub = bul.startswith("→")
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(2)
        r = p.add_run()
        r.text = ("  " if is_sub else "") + bul
        _font(r, FONT_M, 9.5 if is_sub else 11,
              False, MGRAY if is_sub else DARK)


def draw_row_card(slide, x, y, w, h, num, group, title, detail):
    """가로 행 카드 (우수사례용)"""
    rect(slide, x, y, w, h, fill=BGRAY)
    rect(slide, x, y, 4, h, fill=RED)
    badge(slide, x + 12, y + (h - 22) // 2, num, size=22)
    txt(slide, x + 44, y + 10, 110, 22,
        group, FONT_B, 12, True, RED)
    txt(slide, x + 162, y + 10, w - 174, 22,
        title, FONT_B, 12, True, DARK)
    hline(slide, x + 44, y + 36, w - 56)
    tb = slide.shapes.add_textbox(px(x + 44), px(y + 42),
                                   px(w - 56), px(h - 50))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in detail.split("\n"):
        is_sub = line.startswith("→")
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(2)
        r = p.add_run()
        r.text = ("  " if is_sub else "") + line
        _font(r, FONT_M, 9.5 if is_sub else 11,
              False, MGRAY if is_sub else DARK)


# ═══════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════

def build_cover(prs):
    """1p  표지"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    draw_chrome(slide, "하나증권  경영지원그룹", page_num=None, footnote="")

    # ── 흰 본문 영역: 왼쪽 대제목 블록 ──────────────────────
    # 회사명 + 레드 강조선
    hline(slide, 40, BODY_Y + 48, 320, RED)
    txt(slide, 40, BODY_Y + 58, 700, 100,
        "관계사 협업 운영안",
        FONT_B, 38, True, DARK)
    txt(slide, 40, BODY_Y + 164, 600, 36,
        "정기환 전무  |  경영지원그룹",
        FONT_M, 15, False, MGRAY)
    txt(slide, 40, BODY_Y + 202, 300, 26,
        "2026. 04. 19",
        FONT_M, 12, False, MGRAY)

    # ── 우측 장식 블록 ───────────────────────────────────────
    # 레드 구분 바
    rect(slide, 700, BODY_Y + 30, 4, 360, fill=RED)
    # 회색 정보 박스
    rect(slide, 716, BODY_Y + 30, 220, 360, fill=BGRAY)
    txt(slide, 726, BODY_Y + 44, 200, 300,
        "하나증권\n경영지원그룹\n\n관계사 협업\n운영 방향\n\n2026. 04",
        FONT_M, 11, False, MGRAY)


def build_toc(prs):
    """2p  목차"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    draw_chrome(slide, "CONTENTS", page_num=2, footnote="")

    sections = [
        ("Ⅰ", "추진 배경"),
        ("Ⅱ", "세부 실행 방안"),
        ("Ⅲ", "협업 우수사례"),
        ("Ⅳ", "향후 일정 (참고)"),
    ]
    item_h = 96
    y0 = BODY_Y + 18

    for i, (roman, label) in enumerate(sections):
        y = y0 + i * item_h
        bg = WHITE if i % 2 == 0 else BGRAY
        rect(slide, 40, y, 880, item_h - 4, fill=bg)
        rect(slide, 40, y, 62, item_h - 4, fill=RED)
        txt(slide, 40, y + 18, 62, item_h - 36,
            roman, FONT_B, 24, True, WHITE, PP_ALIGN.CENTER)
        txt(slide, 118, y + 26, 720, 44,
            label, FONT_M, 20, False, DARK)
        if i < len(sections) - 1:
            hline(slide, 40, y + item_h - 5, 880)


def build_divider(prs, roman, section_title, sub_items, page_num):
    """3p · 5p · 9p  간지"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 헤더 라벨: 섹션 번호 + 섹션명 (얇은 58pt 바 안에 표기)
    draw_chrome(slide,
                f"{roman}  {section_title}",
                page_num=page_num,
                footnote="")

    # ── 흰 본문 영역: 대형 로마숫자 + 섹션 정보 ─────────────
    # 대형 로마숫자 (좌측)
    txt(slide, 30, BODY_Y + 20, 200, 300,
        roman, FONT_B, 120, True, LGRAY, PP_ALIGN.LEFT)

    # 섹션 제목 + 하위 목차 (우측)
    hline(slide, 260, BODY_Y + 32, 4, RED)   # 수직선 대신 레드 두꺼운 구분선
    rect(slide, 258, BODY_Y + 30, 4, 340, fill=RED)
    txt(slide, 276, BODY_Y + 44, 640, 56,
        section_title, FONT_B, 28, True, DARK)
    hline(slide, 276, BODY_Y + 106, 400, RED)
    y_sub = BODY_Y + 120
    for item in sub_items:
        txt(slide, 292, y_sub, 580, 30,
            item, FONT_M, 14, False, MGRAY)
        y_sub += 36


def build_timeline(prs):
    """4p  1.1 캠페인 타임라인 (3열)"""
    SECS = [("Ⅰ 추진배경", True), ("Ⅱ 실행방안", False),
            ("Ⅲ 우수사례", False), ("Ⅳ 향후일정", False)]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    draw_chrome(slide, "1.1  캠페인 추진 타임라인", 4, SECS)

    COL_X = [30, 326, 622]
    COL_W = 276
    HDR_Y = BODY_Y + 14
    BOD_Y = HDR_Y + 32

    columns = [
        ("2023 ~ 2024", "소통과 참여, 변화와 혁신의 시작",
         "One Team, One Spirit",
         ["→ 관계사 협업 인식 제고 캠페인 착수",
          "→ 칭찬 웨이브·격려 이벤트 시범 운영"]),
        ("2025", "하나금융 DNA로 '다시, 하나답게'",
         "캠페인 본격 확산",
         ["→ 그룹 전 계열사 확대 적용",
          "→ 분기 표창 체계 구축"]),
        ("2026", "역동적인 하나증권",
         "Cheer Up! Level Up! 협업문화 장착",
         ["→ 콜라보북 배포 / 콜라보 DAY 정례화",
          "→ 주간 CEO 실적 점검 체계 운영"]),
    ]

    for i, (year, theme, sub, bullets) in enumerate(columns):
        x = COL_X[i]
        rect(slide, x, HDR_Y, COL_W, 28, fill=RED)
        txt(slide, x, HDR_Y, COL_W, 28,
            year, FONT_B, 13, True, WHITE, PP_ALIGN.CENTER)
        rect(slide, x, BOD_Y, COL_W, FTR_Y - BOD_Y - 4, fill=BGRAY)
        rect(slide, x, BOD_Y, COL_W, 3, fill=RED)
        tb = slide.shapes.add_textbox(
            px(x + 10), px(BOD_Y + 10),
            px(COL_W - 20), px(FTR_Y - BOD_Y - 20))
        tf = tb.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        r = p0.add_run()
        r.text = theme
        _font(r, FONT_B, 13, True, DARK)
        _add_para(tf, sub, FONT_M, 11, False, MGRAY, space_before=2)
        for bul in bullets:
            _add_para(tf, bul, FONT_M, 10, False, MGRAY, space_before=6)
        if i < 2:
            rect(slide, COL_X[i + 1] - 6, HDR_Y, 2, FTR_Y - HDR_Y - 4,
                 fill=LGRAY)


def build_culture(prs):
    """6p  2.1 기업문화"""
    SECS = [("Ⅰ 추진배경", False), ("Ⅱ 실행방안", True),
            ("Ⅲ 우수사례", False), ("Ⅳ 향후일정", False)]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    draw_chrome(slide, "2.1  기업문화 — 협업 분위기 조성", 6, SECS)

    cards = [
        ("①", "칭찬 웨이브", [
            "매주 협업 우수사례 공유",
            "→ 상위 우수직원 칭찬 포인트 지급 (월·분기·기 단위 포상)",
            "→ 사내복지포인트 전환 / 휴양소 지원 / CES 교육 기회 부여",
        ]),
        ("②", "현장 격려 이벤트", [
            "증권-은행 복합점포 등 거점 점포 중심 격려 이벤트 추진",
            "→ 커피·간식 차량이 현장에 직접 방문하여 응원",
        ]),
        ("③", "관계사 협업 = 본업 캠페인", [
            "사내 POP, 플랙카드, 별프로 활용 캠페인",
            "→ 매월 첫 번째 수요일 '콜라보 DAY' 운영",
            "→ 아이디어 공모를 통한 자발적 참여 유도",
        ]),
    ]
    CW, CH = 272, FTR_Y - BODY_Y - 22
    CY = BODY_Y + 10
    for i, (num, title, bullets) in enumerate(cards):
        draw_card(slide, 28 + i * (CW + 14), CY, CW, CH, num, title, bullets)


def build_hr(prs):
    """7p  2.2 인사 — 협업 信賞"""
    SECS = [("Ⅰ 추진배경", False), ("Ⅱ 실행방안", True),
            ("Ⅲ 우수사례", False), ("Ⅳ 향후일정", False)]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    draw_chrome(slide, "2.2  인사 — 협업 信賞 체계", 7, SECS)

    items = [
        ("①", "인사 가점 부여",
         "분기별 표창을 통해 CEO 표창에 준하는 수준의 인사 가점 부여\n"
         "→ 협업 실적 우수자 전보·승진 시 우선 배려"),
        ("②", "우수 외부연수 기회 우선 배정",
         "배정 인원의 30% 우선 배정\n"
         "→ 국내외 우수 연수 프로그램 선발 시 협업 실적 반영"),
    ]
    CW = 412
    CH = FTR_Y - BODY_Y - 22
    GAP = 26
    Y = BODY_Y + 10

    for i, (num, title, detail) in enumerate(items):
        x = 28 + i * (CW + GAP)
        rect(slide, x, Y, CW, CH, fill=BGRAY)
        rect(slide, x, Y, CW, 3, fill=RED)
        badge(slide, x + 12, Y + 12, num, size=22)
        txt(slide, x + 40, Y + 13, CW - 52, 26,
            title, FONT_B, 14, True, DARK)
        hline(slide, x + 12, Y + 44, CW - 24)
        tb = slide.shapes.add_textbox(
            px(x + 12), px(Y + 52),
            px(CW - 24), px(CH - 64))
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for line in detail.split("\n"):
            is_sub = line.startswith("→")
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
                p.space_before = Pt(5)
            r = p.add_run()
            r.text = ("  " if is_sub else "") + line
            _font(r, FONT_M, 10 if is_sub else 12,
                  False, MGRAY if is_sub else DARK)


def build_synergy(prs):
    """8p  2.3 시너지"""
    SECS = [("Ⅰ 추진배경", False), ("Ⅱ 실행방안", True),
            ("Ⅲ 우수사례", False), ("Ⅳ 향후일정", False)]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    draw_chrome(slide, "2.3  시너지 — 그룹별 협업 마케팅 활성화", 8, SECS)

    cards = [
        ("①", "콜라보북 제작 및 전직원 배포", [
            '"하나로 연결된 모두의 금융" 실천',
            "→ 관계사 상품정보·리플렛 담긴 콜라보 포켓북",
            "→ 손님 상담 시 관계사 상품 자동 제안 문화 만들기",
        ]),
        ("②", "그룹별 협업 마케팅 구조화", [
            "WM그룹: 증권 패밀리오피스 자산가 전용 하나카드 제작",
            "→ 자산가의 증권 로열티 함양",
            "IB그룹: Deal 발주 단계 관계사 상품 패키지 involve",
            "S&T그룹: 법인금상 손님 → 하나저축은행 금리상품 확대",
        ]),
        ("③", "관계사 협업 실적 주간 점검", [
            "CEO께서 직접 관리 (주간 단위)",
            "→ 그룹·팀별 실적 취합 후 대시보드 공유",
            "→ 우수 그룹 인센티브, 미달 그룹 개선 피드백",
        ]),
    ]
    CW, CH = 272, FTR_Y - BODY_Y - 22
    CY = BODY_Y + 10
    for i, (num, title, bullets) in enumerate(cards):
        draw_card(slide, 28 + i * (CW + 14), CY, CW, CH, num, title, bullets)


def build_cases(prs):
    """10p  3.1 협업 우수사례"""
    SECS = [("Ⅰ 추진배경", False), ("Ⅱ 실행방안", False),
            ("Ⅲ 우수사례", True), ("Ⅳ 향후일정", False)]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    draw_chrome(slide, "3.1  협업 우수사례", 10, SECS)

    cases = [
        ("①", "WM그룹", "하나카드 제안으로 150억 예금 유치",
         "The 센터필드W 등 증권 패밀리오피스 자산가 전용 하나카드 제작\n"
         "→ 카드 발급 연계 예금 상품 제안으로 150억 신규 예금 유치\n"
         "→ 자산가의 하나증권 로열티 제고 효과 확인"),
        ("②", "IB그룹", "디엠티 브릿지론 인수금융 딜에 관계사 상품 involve",
         "인프라·부동산·물류창고 공사 Deal 발주 단계 관계사 상품 패키지 구성\n"
         "→ 하나은행·하나카드 공동 패키지로 Deal 경쟁력 제고\n"
         "→ 관계사 수수료 수익 공유 모델 확립"),
        ("③", "S&T그룹", "하나저축은행 금리상품 판매 확대",
         "당사 법인금상 손님 대상 하나저축은행 금리형 상품 판매\n"
         "→ 금리형 상품 600억 매각 달성\n"
         "→ 시너지 데이터 수집 중 — 추가 실적 업데이트 예정"),
    ]
    RH = (FTR_Y - BODY_Y - 30) // 3
    RW = 846
    Y0 = BODY_Y + 8
    for i, (num, group, title, detail) in enumerate(cases):
        draw_row_card(slide, 30, Y0 + i * (RH + 8), RW, RH,
                      num, group, title, detail)


def build_closing(prs):
    """11p  마무리"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    draw_chrome(slide, "하나증권  경영지원그룹", page_num=11, footnote="")

    # 흰 본문 — 표지와 동일한 구조
    hline(slide, 40, BODY_Y + 48, 320, RED)
    txt(slide, 40, BODY_Y + 62, 700, 100,
        "감사합니다",
        FONT_B, 38, True, DARK)
    txt(slide, 40, BODY_Y + 168, 600, 36,
        "정기환 전무  |  경영지원그룹",
        FONT_M, 15, False, MGRAY)
    txt(slide, 40, BODY_Y + 206, 300, 26,
        "2026. 04. 19",
        FONT_M, 12, False, MGRAY)

    rect(slide, 700, BODY_Y + 30, 4, 360, fill=RED)
    rect(slide, 716, BODY_Y + 30, 220, 360, fill=BGRAY)
    txt(slide, 726, BODY_Y + 44, 200, 300,
        "협업이 본업입니다\n\n하나증권\n경영지원그룹\n\n2026. 04",
        FONT_M, 11, False, MGRAY)


# ═══════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    build_cover(prs)                                                   # 1
    build_toc(prs)                                                     # 2
    build_divider(prs, "Ⅰ", "추진 배경",
                  ["1.1  캠페인 추진 타임라인"], page_num=3)           # 3
    build_timeline(prs)                                                # 4
    build_divider(prs, "Ⅱ", "세부 실행 방안",
                  ["2.1  기업문화 — 협업 분위기 조성",
                   "2.2  인사 — 협업 信賞 체계",
                   "2.3  시너지 — 그룹별 협업 마케팅"],
                  page_num=5)                                           # 5
    build_culture(prs)                                                 # 6
    build_hr(prs)                                                      # 7
    build_synergy(prs)                                                 # 8
    build_divider(prs, "Ⅲ", "협업 우수사례",
                  ["3.1  WM·IB·S&T 그룹 우수사례"], page_num=9)        # 9
    build_cases(prs)                                                   # 10
    build_closing(prs)                                                 # 11

    out = r"C:\Users\jay\Downloads\1_rebuilt.pptx"
    prs.save(out)
    print(f"Saved → {out}")

    print("\n[검증]")
    prs2 = Presentation(out)
    print(f"슬라이드 수: {len(prs2.slides)},  크기: {prs2.slide_width.pt} x {prs2.slide_height.pt} pt")
    for idx, sl in enumerate(prs2.slides, 1):
        texts = []
        for sh in sl.shapes:
            if sh.has_text_frame:
                for par in sh.text_frame.paragraphs:
                    t = par.text.strip()
                    if t:
                        texts.append(t[:36])
            if len(texts) >= 3:
                break
        print(f"  Slide {idx:2d}: {' | '.join(texts[:3])}")


if __name__ == "__main__":
    main()
