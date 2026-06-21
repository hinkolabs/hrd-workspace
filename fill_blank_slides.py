"""
AI교육 최종.pptx 빈칸 슬라이드 채우기
- Slide 6  (대본용 → AI 규정 안내 슬라이드)
- Slide 16 (AI 질문 샘플 채우기)
- Slide 19 (엑셀 데이터 채우기 실습)
- Slide 20 (엑셀 오류 검증 실습)
- Slide 22 (AI 챗봇 플랫폼 소개)
- Slide 23 (챗봇 만들기 실습)
- Slide 24 (마무리)
"""

import sys, os, shutil
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pptx.oxml.ns import qn

SRC  = r'C:\Users\jay\Downloads\AI교육 최종.pptx'
DST  = r'C:\Users\jay\Downloads\AI교육_완성판.pptx'
shutil.copy2(SRC, DST)
prs  = Presentation(DST)

# ── 공통 색상 ──────────────────────────────────────────────────────
GREEN      = RGBColor(0x00, 0x91, 0x78)
DARK_GREEN = RGBColor(0x00, 0x4E, 0x42)
LIGHT_GREEN= RGBColor(0xD7, 0xED, 0xE6)
LIGHT_GRAY = RGBColor(0xF7, 0xF9, 0xF8)
DARK_TEXT  = RGBColor(0x22, 0x22, 0x22)
RED        = RGBColor(0xCC, 0x00, 0x00)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY   = RGBColor(0x66, 0x66, 0x66)

FONT = '맑은 고딕'
SLIDE_W = prs.slide_width   # 10.0 인치
SLIDE_H = prs.slide_height  # 7.5  인치


# ══════════════════════════════════════════════════════════
#  헬퍼 함수
# ══════════════════════════════════════════════════════════

def rgb_fill(shape, rgb: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def add_rect(slide, left, top, width, height, rgb: RGBColor):
    box = slide.shapes.add_shape(1, left, top, width, height)
    rgb_fill(box, rgb)
    return box


def add_textbox(slide, left, top, width, height,
                text, size=12, bold=False, color=DARK_TEXT,
                align=PP_ALIGN.LEFT, wrap=True, line_spacing=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        from pptx.util import Pt as _Pt
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    return tb


def common_header(slide, section_label, title, slide_num):
    """왼쪽 녹색 바 + 섹션 라벨 + 제목 + 구분선"""
    add_rect(slide, 0, 0, Inches(0.14), SLIDE_H, GREEN)
    add_textbox(slide, Inches(0.47), Inches(0.24), Inches(7.87), Inches(0.28),
                section_label, size=10, color=GREEN)
    add_textbox(slide, Inches(0.47), Inches(0.48), Inches(9.06), Inches(0.54),
                title, size=26, bold=True, color=DARK_TEXT)
    add_rect(slide, Inches(0.47), Inches(1.07), Inches(9.06), Inches(0.02), GREEN)


def common_footer(slide, page_num):
    """하단 녹색 바 + copyright + 페이지 번호"""
    add_rect(slide, 0, Inches(7.05), SLIDE_W, Inches(0.45), GREEN)
    add_textbox(slide, Inches(1.57), Inches(7.11), Inches(6.30), Inches(0.33),
                'Copyright ⓒ 2025 Hana Financial Group. All rights reserved.',
                size=6, color=WHITE)
    add_textbox(slide, Inches(8.98), Inches(7.09), Inches(0.91), Inches(0.37),
                f'| {page_num}', size=8, bold=True, color=WHITE,
                align=PP_ALIGN.RIGHT)


def add_numbered_row(slide, num_str, label, top_y, row_h=Inches(1.26)):
    """슬라이드 10 스타일의 번호+내용 row"""
    add_rect(slide, Inches(0.47), top_y, Inches(0.98), row_h, GREEN)
    add_textbox(slide, Inches(0.47), top_y, Inches(0.98), row_h,
                num_str, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(1.50), top_y, Inches(8.03), row_h, LIGHT_GRAY)
    tb = add_textbox(slide, Inches(1.65), top_y + Inches(0.40), Inches(7.80), Inches(0.46),
                     label, size=14, bold=True, color=DARK_TEXT)
    return tb


def add_numbered_row_with_desc(slide, num_str, label, desc, top_y, row_h=Inches(1.0)):
    """번호 + 제목 + 설명이 있는 row"""
    add_rect(slide, Inches(0.47), top_y, Inches(0.98), row_h, GREEN)
    add_textbox(slide, Inches(0.47), top_y, Inches(0.98), row_h,
                num_str, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(1.50), top_y, Inches(8.03), row_h, LIGHT_GRAY)
    add_textbox(slide, Inches(1.65), top_y + Inches(0.10), Inches(7.80), Inches(0.32),
                label, size=13, bold=True, color=DARK_GREEN)
    add_textbox(slide, Inches(1.65), top_y + Inches(0.42), Inches(7.80), row_h - Inches(0.52),
                desc, size=10, color=DARK_TEXT, wrap=True)


def add_principle_card(slide, num_str, heading, bad, good, tip, left_x, top_y,
                       card_w=Inches(4.41), card_h=Inches(1.65)):
    """슬라이드 8 스타일의 원칙 카드 (좌/우 배치용)"""
    num_w = Inches(0.63)
    add_rect(slide, left_x, top_y, card_w, card_h, LIGHT_GRAY)
    add_rect(slide, left_x, top_y, num_w, card_h, GREEN)
    add_textbox(slide, left_x, top_y, num_w, card_h,
                num_str, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, left_x + num_w, top_y + Inches(0.06), card_w - num_w, Inches(0.35),
                heading, size=13, bold=True, color=DARK_GREEN)
    add_rect(slide, left_x + num_w, top_y + Inches(0.41), card_w - num_w, Inches(0.02), GRAY)
    if bad:
        add_textbox(slide, left_x + num_w, top_y + Inches(0.45), card_w - num_w, Inches(0.30),
                    f'❌  {bad}', size=10, color=RED)
    if good:
        add_textbox(slide, left_x + num_w, top_y + Inches(0.75), card_w - num_w, Inches(0.50),
                    f'✅  {good}', size=10, color=DARK_GREEN, wrap=True)
    if tip:
        add_rect(slide, left_x + num_w, top_y + card_h - Inches(0.35),
                 card_w - num_w, Inches(0.33), LIGHT_GREEN)
        add_textbox(slide, left_x + num_w + Inches(0.08), top_y + card_h - Inches(0.35),
                    card_w - num_w - Inches(0.08), Inches(0.33),
                    f'→  {tip}', size=9, color=DARK_GREEN, wrap=True)


# ══════════════════════════════════════════════════════════
#  Slide 6 — AI 업무 규정 안내 (기존 "대본용" 슬라이드 교체)
# ══════════════════════════════════════════════════════════
def fill_slide6():
    slide = prs.slides[5]
    # 기존 텍스트박스 제거
    sp = slide.shapes[0]._element
    sp.getparent().remove(sp)

    # 배경 흰색
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    common_header(slide, '1교시 · AI를 대하는 태도', 'AI 업무 활용 규정 안내', 6)
    common_footer(slide, 6)

    rows = [
        ('01', '하나금융그룹 AI 업무 가이드라인 (제정 예정)',
         '임직원 소셜미디어 가이드라인처럼, AI 활용에 관한 내부 규정을 연내 제정·배포 예정입니다. '
         '발표 즉시 내용을 숙지하고 실무에 반영해야 합니다.'),
        ('02', '지금 당장 지켜야 할 기본 원칙',
         '① 고객 실명·계좌번호 등 개인정보 입력 금지  '
         '② MNPI(미공개 중요 정보) 절대 입력 금지  '
         '③ 발간 전 리서치 초안 등 대외비 자료 입력 금지'),
        ('03', '모든 AI 결과물은 최종 제출 전 반드시 검증',
         'AI 출력물에는 오류·환각이 포함될 수 있습니다. '
         '수치·법규 등 중요 내용은 원문 자료와 반드시 대조 확인하세요.'),
    ]

    row_h = Inches(1.10)
    gap   = Inches(0.15)
    for i, (num, label, desc) in enumerate(rows):
        top_y = Inches(1.30) + i * (row_h + gap)
        add_numbered_row_with_desc(slide, num, label, desc, top_y, row_h=row_h)

    # 하단 강조 박스
    add_rect(slide, Inches(0.47), Inches(5.90), Inches(9.06), Inches(0.75), LIGHT_GREEN)
    add_rect(slide, Inches(0.47), Inches(5.90), Inches(0.06), Inches(0.75), GREEN)
    add_textbox(slide, Inches(0.63), Inches(6.05), Inches(8.80), Inches(0.45),
                '⚠️  AI는 강력한 도구이지만, 최종 판단과 책임은 항상 여러분에게 있습니다.',
                size=11, bold=True, color=DARK_GREEN)

    print('  [Slide 6] AI 규정 안내 슬라이드 완료')


# ══════════════════════════════════════════════════════════
#  Slide 16 — AI 질문 샘플 채우기
# ══════════════════════════════════════════════════════════
def fill_slide16():
    slide = prs.slides[15]

    # 기존 AUTO_SHAPE(말풍선) 제거 — Shape[4]
    to_remove = []
    for shape in slide.shapes:
        if shape.name == '모서리가 둥근 사각형 설명선 2':
            to_remove.append(shape._element)
    for el in to_remove:
        el.getparent().remove(el)

    # "AI 질문 샘플" 라벨 위치 확인용으로 기존 TextBox 11 아래에 내용 추가
    # 현재 슬라이드 레이아웃: L:0.47 T:0.24 (width), 사진 영역이 오른쪽에 있음

    # AI 질문 샘플 박스 영역 (이미지 오른쪽에 위치하도록)
    add_rect(slide, Inches(4.80), Inches(2.20), Inches(4.85), Inches(4.50), LIGHT_GRAY)
    add_rect(slide, Inches(4.80), Inches(2.20), Inches(4.85), Inches(0.38), GREEN)
    add_textbox(slide, Inches(4.90), Inches(2.23), Inches(4.60), Inches(0.35),
                '💬 AI 질문 샘플', size=13, bold=True, color=WHITE)

    questions = [
        ('Q1', 'SK하이닉스의 2025년 매출액은 얼마이며, 전년 대비 증감률은?'),
        ('Q2', '최근 3개년 영업이익률 추이와 수익성 개선·악화 요인은?'),
        ('Q3', '사업보고서에서 가장 중요한 위험 요인 2가지는 무엇인가?'),
        ('Q4', '2025년 CAPEX 규모와 그것이 향후 실적에 미치는 가능성은?'),
        ('Q5', '부채비율은 몇 %이며, 이는 경쟁사 대비 어떤 수준인가?'),
    ]

    for i, (q_num, q_text) in enumerate(questions):
        y = Inches(2.68) + i * Inches(0.74)
        add_rect(slide, Inches(4.90), y, Inches(0.42), Inches(0.55), DARK_GREEN)
        add_textbox(slide, Inches(4.90), y, Inches(0.42), Inches(0.55),
                    q_num, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(5.40), y + Inches(0.08), Inches(4.15), Inches(0.45),
                    q_text, size=10, color=DARK_TEXT, wrap=True)
        add_rect(slide, Inches(4.90), y + Inches(0.58), Inches(4.55), Inches(0.01), GRAY)

    # 하단 안내
    add_rect(slide, Inches(4.80), Inches(6.43), Inches(4.85), Inches(0.40), LIGHT_GREEN)
    add_textbox(slide, Inches(4.90), Inches(6.48), Inches(4.65), Inches(0.35),
                '→ 답변을 AI에 입력하면 즉각 채점·피드백 제공', size=9, color=DARK_GREEN)

    print('  [Slide 16] AI 질문 샘플 완료')


# ══════════════════════════════════════════════════════════
#  Slide 19 — 엑셀 실습 1: 여러 자료 → 하나의 표
# ══════════════════════════════════════════════════════════
def fill_slide19():
    slide = prs.slides[18]
    # 기존 텍스트박스 제거
    sp = slide.shapes[0]._element
    sp.getparent().remove(sp)

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    common_header(slide, '3교시 · 데이터 분석 활용 (EXCEL)',
                  '여러 개의 자료를 하나의 표로 데이터값 채우기', 19)
    common_footer(slide, 19)

    # 왼쪽 : 실습 목표 + 프롬프트
    add_rect(slide, Inches(0.47), Inches(1.30), Inches(5.25), Inches(5.30), LIGHT_GRAY)
    add_rect(slide, Inches(0.47), Inches(1.30), Inches(5.25), Inches(0.38), GREEN)
    add_textbox(slide, Inches(0.58), Inches(1.33), Inches(5.00), Inches(0.35),
                '실습 목표', size=13, bold=True, color=WHITE)

    steps = [
        ('STEP 1', '엑셀 파일과 AI 창을 함께 열어 놓기'),
        ('STEP 2', '거래원장_RAW 시트 구조를 AI에게 설명'),
        ('STEP 3', '원하는 집계/정리 수식 요청'),
        ('STEP 4', '결과값 직접 확인 및 검증'),
    ]
    for i, (label, desc) in enumerate(steps):
        y = Inches(1.78) + i * Inches(1.14)
        add_rect(slide, Inches(0.55), y, Inches(1.05), Inches(0.95), GREEN)
        add_textbox(slide, Inches(0.55), y, Inches(1.05), Inches(0.95),
                    label, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.70), y + Inches(0.22), Inches(3.90), Inches(0.55),
                    desc, size=11, color=DARK_TEXT, wrap=True)
        if i < len(steps) - 1:
            add_rect(slide, Inches(0.55), y + Inches(0.98), Inches(5.05), Inches(0.01), GRAY)

    # 오른쪽 : 프롬프트 예시
    add_rect(slide, Inches(5.90), Inches(1.30), Inches(3.70), Inches(5.30), LIGHT_GRAY)
    add_rect(slide, Inches(5.90), Inches(1.30), Inches(3.70), Inches(0.38), DARK_GREEN)
    add_textbox(slide, Inches(6.00), Inches(1.33), Inches(3.50), Inches(0.35),
                '프롬프트 예시', size=13, bold=True, color=WHITE)

    prompt_text = (
        '너는 엑셀 전문가야.\n\n'
        '거래원장_RAW 시트:\n'
        '  A열=날짜, B열=지점코드,\n'
        '  E열=상품유형, G열=거래금액,\n'
        '  J열=처리상태\n\n'
        '처리상태가 "정상"인 거래만\n'
        '상품유형별 합계를 구하는\n'
        'SUMIF 수식을\n'
        '결과작성_STEP1 시트 C2셀부터\n'
        '만들어줘.\n\n'
        '수식의 범위는 절대참조로\n'
        '고정해줘.'
    )
    add_textbox(slide, Inches(6.00), Inches(1.78), Inches(3.50), Inches(4.70),
                prompt_text, size=9.5, color=DARK_GREEN, wrap=True)

    print('  [Slide 19] 엑셀 실습 1 완료')


# ══════════════════════════════════════════════════════════
#  Slide 20 — 엑셀 실습 2: 오류 검증
# ══════════════════════════════════════════════════════════
def fill_slide20():
    slide = prs.slides[19]
    sp = slide.shapes[0]._element
    sp.getparent().remove(sp)

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    common_header(slide, '3교시 · 데이터 분석 활용 (EXCEL)',
                  '엑셀 데이터의 오류 검증 및 수정', 20)
    common_footer(slide, 20)

    # 오류 유형 4가지 카드 (2×2 그리드)
    error_types = [
        ('01', '금액 이상값', '거래금액이 0 이하이거나, 최대값과 비교해\n극단적으로 크거나 작은 건 탐지'),
        ('02', '날짜 오류', '미래 날짜 입력, 공백, 비정상 형식\n(예: 숫자가 아닌 텍스트로 저장된 날짜)'),
        ('03', '상품 코드 불일치', '기준 코드 목록에 없는 값 또는\n오타·대소문자 불일치 탐지'),
        ('04', '중복 거래', '동일 날짜+지점+금액 조합이\n2회 이상 반복되는 건 탐지'),
    ]
    positions = [
        (Inches(0.47), Inches(1.30)),
        (Inches(5.22), Inches(1.30)),
        (Inches(0.47), Inches(3.60)),
        (Inches(5.22), Inches(3.60)),
    ]
    card_w = Inches(4.57)
    card_h = Inches(2.10)

    for (num, title, desc), (lx, ty) in zip(error_types, positions):
        add_rect(slide, lx, ty, card_w, card_h, LIGHT_GRAY)
        add_rect(slide, lx, ty, Inches(0.72), card_h, GREEN)
        add_textbox(slide, lx, ty, Inches(0.72), card_h,
                    num, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, lx + Inches(0.80), ty + Inches(0.20), card_w - Inches(0.90), Inches(0.38),
                    title, size=13, bold=True, color=DARK_GREEN)
        add_rect(slide, lx + Inches(0.80), ty + Inches(0.62), card_w - Inches(0.90), Inches(0.01), GRAY)
        add_textbox(slide, lx + Inches(0.80), ty + Inches(0.70), card_w - Inches(0.90), Inches(1.25),
                    desc, size=10, color=DARK_TEXT, wrap=True)

    # 하단 프롬프트 힌트
    add_rect(slide, Inches(0.47), Inches(5.90), Inches(9.06), Inches(0.85), LIGHT_GREEN)
    add_rect(slide, Inches(0.47), Inches(5.90), Inches(0.06), Inches(0.85), GREEN)
    add_textbox(slide, Inches(0.63), Inches(5.95), Inches(8.80), Inches(0.35),
                '💡 AI 활용 프롬프트: "이 데이터에서 금액이 0 이하인 건, 날짜 오류, 중복 거래를 찾아서 몇 번 행인지 알려줘."',
                size=10, color=DARK_GREEN, wrap=True)

    print('  [Slide 20] 엑셀 실습 2 완료')


# ══════════════════════════════════════════════════════════
#  Slide 22 — AI 챗봇 플랫폼 소개
# ══════════════════════════════════════════════════════════
def fill_slide22():
    slide = prs.slides[21]
    sp = slide.shapes[0]._element
    sp.getparent().remove(sp)

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    common_header(slide, '4교시 · 나만의 챗봇 업무 비서 만들기',
                  '당사 도입 예정 AI 챗봇 플랫폼 소개', 22)
    common_footer(slide, 22)

    # Alli 플랫폼 소개 헤더
    add_rect(slide, Inches(0.47), Inches(1.30), Inches(9.06), Inches(0.65), GREEN)
    add_textbox(slide, Inches(0.60), Inches(1.38), Inches(8.80), Inches(0.50),
                'Alli Works — 사내 업무 지식 기반 AI 챗봇 플랫폼',
                size=16, bold=True, color=WHITE)

    # 3단 특징 카드
    features = [
        ('RAG 기반\n지식 검색',
         '사전 학습 모델이 아닌, 사내 문서·매뉴얼을 직접 학습하여 정확한 답변 제공'),
        ('업무 맞춤\n챗봇 제작',
         '부서별 FAQ, 업무 프로세스, 규정집 등을 학습시켜 전용 챗봇을 직접 만들 수 있음'),
        ('보안 및\n접근 통제',
         'MNPI·개인정보가 없는 내부 자료만 학습 허용. 부서 권한별 접근 범위 설정 가능'),
    ]

    feat_w = Inches(2.85)
    gap    = Inches(0.25)
    for i, (title, desc) in enumerate(features):
        lx = Inches(0.47) + i * (feat_w + gap)
        ty = Inches(2.15)
        add_rect(slide, lx, ty, feat_w, Inches(2.90), LIGHT_GRAY)
        add_rect(slide, lx, ty, feat_w, Inches(0.06), GREEN)
        add_textbox(slide, lx + Inches(0.15), ty + Inches(0.20), feat_w - Inches(0.30), Inches(0.60),
                    title, size=13, bold=True, color=DARK_GREEN, align=PP_ALIGN.CENTER)
        add_rect(slide, lx + Inches(0.15), ty + Inches(0.85), feat_w - Inches(0.30), Inches(0.01), GRAY)
        add_textbox(slide, lx + Inches(0.15), ty + Inches(0.95), feat_w - Inches(0.30), Inches(1.85),
                    desc, size=10, color=DARK_TEXT, wrap=True, align=PP_ALIGN.LEFT)

    # 활용 가능 업무 목록
    add_rect(slide, Inches(0.47), Inches(5.30), Inches(9.06), Inches(0.38), DARK_GREEN)
    add_textbox(slide, Inches(0.60), Inches(5.33), Inches(8.80), Inches(0.33),
                '활용 가능 업무 예시', size=12, bold=True, color=WHITE)

    use_cases = '① 반복 업무 FAQ 자동 답변   ② 업무 프로세스 안내   ③ 내부 규정 검색   ④ 회의록·보고서 초안 작성'
    add_textbox(slide, Inches(0.47), Inches(5.78), Inches(9.06), Inches(0.60),
                use_cases, size=11, color=DARK_TEXT, wrap=True)

    print('  [Slide 22] AI 챗봇 플랫폼 소개 완료')


# ══════════════════════════════════════════════════════════
#  Slide 23 — 챗봇 만들기 실습
# ══════════════════════════════════════════════════════════
def fill_slide23():
    slide = prs.slides[22]
    sp = slide.shapes[0]._element
    sp.getparent().remove(sp)

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    common_header(slide, '4교시 · 나만의 챗봇 업무 비서 만들기',
                  '챗봇 만들기 실습', 23)
    common_footer(slide, 23)

    steps_full = [
        ('STEP 1', '챗봇 목적 설정',
         '어떤 질문에 답할 챗봇인지 결정합니다.\n예) "신입사원 업무 FAQ 답변 챗봇"'),
        ('STEP 2', '시스템 프롬프트 작성',
         '역할·어조·답변 범위·금지 사항을 지정합니다.\n예) "너는 하나증권 신입사원 안내 도우미야. 업무 외 질문은 정중히 거절해."'),
        ('STEP 3', '지식 베이스 업로드',
         '사내 FAQ, 업무 매뉴얼 등 PDF·TXT 파일을 업로드합니다.\n(단, MNPI·고객 개인정보 포함 자료 제외)'),
        ('STEP 4', '테스트 및 품질 검증',
         '실제 질문을 10개 이상 입력해 답변 품질을 확인합니다.\n오답·환각 발견 시 시스템 프롬프트를 수정합니다.'),
    ]

    step_h = Inches(1.22)
    gap    = Inches(0.10)
    for i, (snum, stitle, sdesc) in enumerate(steps_full):
        ty = Inches(1.28) + i * (step_h + gap)
        # 번호 원형 (녹색 박스)
        add_rect(slide, Inches(0.47), ty, Inches(0.95), step_h, GREEN)
        add_textbox(slide, Inches(0.47), ty, Inches(0.95), step_h,
                    snum, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, Inches(1.47), ty, Inches(8.06), step_h, LIGHT_GRAY)
        add_textbox(slide, Inches(1.60), ty + Inches(0.12), Inches(7.80), Inches(0.35),
                    stitle, size=13, bold=True, color=DARK_GREEN)
        add_rect(slide, Inches(1.60), ty + Inches(0.50), Inches(7.80), Inches(0.01), GRAY)
        add_textbox(slide, Inches(1.60), ty + Inches(0.55), Inches(7.80), step_h - Inches(0.60),
                    sdesc, size=10, color=DARK_TEXT, wrap=True)

    print('  [Slide 23] 챗봇 만들기 실습 완료')


# ══════════════════════════════════════════════════════════
#  Slide 24 — 마무리
# ══════════════════════════════════════════════════════════
def fill_slide24():
    slide = prs.slides[23]
    sp = slide.shapes[0]._element
    sp.getparent().remove(sp)

    # 배경 짙은 녹색
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x00, 0x4E, 0x42)

    # 장식용 원 (오른쪽 상단, 왼쪽 하단)
    for (cx, cy, r) in [(9.5, -0.5, 3.0), (-1.5, 7.0, 2.5)]:
        c = slide.shapes.add_shape(3,
                                   Inches(cx - r), Inches(cy - r),
                                   Inches(r * 2), Inches(r * 2))
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE
        c.line.fill.background()

    # 타이틀
    add_textbox(slide, Inches(0.80), Inches(1.50), Inches(11.20), Inches(1.20),
                '오늘 배운 것들을 정리해볼까요?', size=34, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)

    # 구분선
    add_rect(slide, Inches(2.50), Inches(2.80), Inches(7.30), Inches(0.05), GREEN)

    # 핵심 메시지 3가지
    messages = [
        ('🎯', 'AI는 도구입니다', '최종 판단과 책임은 언제나 여러분에게 있습니다.'),
        ('🔍', '검증은 필수입니다', 'AI 결과물을 그대로 제출하지 마세요. 반드시 원문과 대조하세요.'),
        ('🛡️', '윤리가 먼저입니다', 'MNPI·개인정보 보호 원칙을 항상 준수하세요.'),
    ]

    card_w = Inches(2.85)
    for i, (icon, title, desc) in enumerate(messages):
        lx = Inches(0.47) + i * (card_w + Inches(0.22))
        ty = Inches(3.05)
        add_rect(slide, lx, ty, card_w, Inches(2.85), RGBColor(0x00, 0x66, 0x5A))
        add_rect(slide, lx, ty, card_w, Inches(0.06), GREEN)
        add_textbox(slide, lx + Inches(0.15), ty + Inches(0.25), card_w - Inches(0.30), Inches(0.55),
                    icon, size=26, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, lx + Inches(0.10), ty + Inches(0.88), card_w - Inches(0.20), Inches(0.45),
                    title, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, lx + Inches(0.15), ty + Inches(1.38), card_w - Inches(0.30), Inches(0.02),
                 RGBColor(0x00, 0x91, 0x78))
        add_textbox(slide, lx + Inches(0.10), ty + Inches(1.50), card_w - Inches(0.20), Inches(1.20),
                    desc, size=10, color=RGBColor(0xCC, 0xEE, 0xE8), align=PP_ALIGN.CENTER, wrap=True)

    # 마무리 문구
    add_textbox(slide, Inches(0.80), Inches(6.15), Inches(11.20), Inches(0.60),
                'AI를 올바르게 활용하는 하나인이 되어주세요.',
                size=15, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    # Copyright 풋터
    add_textbox(slide, Inches(1.57), Inches(7.11), Inches(6.30), Inches(0.33),
                'Copyright ⓒ 2025 Hana Financial Group. All rights reserved.',
                size=6, color=RGBColor(0xAA, 0xAA, 0xAA))
    add_textbox(slide, Inches(8.98), Inches(7.09), Inches(0.91), Inches(0.37),
                '| 24', size=8, bold=True, color=RGBColor(0xAA, 0xAA, 0xAA),
                align=PP_ALIGN.RIGHT)

    print('  [Slide 24] 마무리 슬라이드 완료')


# ══════════════════════════════════════════════════════════
#  실행
# ══════════════════════════════════════════════════════════
if __name__ == '__main__':
    print('=== 빈칸 채우기 시작 ===\n')
    fill_slide6()
    fill_slide16()
    fill_slide19()
    fill_slide20()
    fill_slide22()
    fill_slide23()
    fill_slide24()

    prs.save(DST)
    print(f'\n✓ 저장 완료: {DST}')
