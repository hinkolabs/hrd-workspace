import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches

SRC = r'C:\Users\jay\Downloads\AI교육_완성판_v8.pptx'
prs = Presentation(SRC)

def show(idx):
    slide = prs.slides[idx]
    print(f'\n=== Slide {idx+1} ({len(slide.shapes)} shapes) ===')
    for i, shape in enumerate(slide.shapes):
        L = shape.left/914400
        T = shape.top/914400
        W = shape.width/914400
        H = shape.height/914400
        txt = ''
        if shape.has_text_frame:
            lines = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
            if lines:
                txt = ' | ' + ' / '.join(lines)[:100]
        print(f'  [{i:2d}] type={shape.shape_type} T={T:.2f}" L={L:.2f}" W={W:.2f}" H={H:.2f}"{txt}')

show(22)   # slide 23
show(23)   # slide 24

# 스타일 참고용: 잘 만들어진 콘텐츠 슬라이드 (예: slide 7, 8, 22)
print('\n=== 스타일 참고: Slide 22 shapes ===')
show(21)
