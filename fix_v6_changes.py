"""
1. Slide 6  item 02 → 윤리경영 기준 내용으로 교체
2. Slide 18 → 원본 5개 항목 복원, 이모티콘 제거
3. 전체 슬라이드 교시 → 목차 섹션 번호로 교체
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SRC = r'C:\Users\jay\Downloads\AI교육_완성판_v6.pptx'
DST = r'C:\Users\jay\Downloads\AI교육_완성판_v7.pptx'
prs = Presentation(SRC)

GREEN       = RGBColor(0x00, 0x91, 0x78)
DARK_GREEN  = RGBColor(0x00, 0x4E, 0x42)
LIGHT_GREEN = RGBColor(0xD7, 0xED, 0xE6)
LIGHT_GRAY  = RGBColor(0xF7, 0xF9, 0xF8)
DARK_TEXT   = RGBColor(0x22, 0x22, 0x22)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY        = RGBColor(0xCC, 0xCC, 0xCC)
FONT        = '맑은 고딕'


def rgb_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def add_rect(slide, l, t, w, h, rgb):
    s = slide.shapes.add_shape(1, l, t, w, h)
    rgb_fill(s, rgb)
    return s


def add_tb(slide, l, t, w, h, text, size=12, bold=False,
           color=DARK_TEXT, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = FONT
    return tb


def set_text_run(shape, text, size=None, bold=None, color=None):
    """기존 shape의 첫 번째 paragraph 첫 번째 run 텍스트 교체"""
    tf = shape.text_frame
    p = tf.paragraphs[0]
    if p.runs:
        run = p.runs[0]
        run.text = text
        if size:  run.font.size = Pt(size)
        if bold is not None: run.font.bold = bold
        if color: run.font.color.rgb = color
    else:
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        if size:  run.font.size = Pt(size)
        if bold is not None: run.font.bold = bold
        if color: run.font.color.rgb = color


def remove_shapes_by_idx(slide, indices):
    shapes_list = list(slide.shapes)
    sp_tree = slide.shapes._spTree
    for i in sorted(set(indices), reverse=True):
        if i < len(shapes_list):
            el = shapes_list[i]._element
            el.getparent().remove(el)


# ══════════════════════════════════════════════════════════
#  1. Slide 6 — item 02 윤리경영 기준으로 교체
# ══════════════════════════════════════════════════════════
def fix_slide6():
    slide = prs.slides[5]
    shapes = list(slide.shapes)

    # 현재 item 02 shapes 찾기: "02" 텍스트를 포함하는 그룹
    # shapes[20]=rect badge bg, [21]=tb "02", [22]=rect body bg,
    # [23]=tb title "사내 허용 AI 서비스 안내", [24]=tb content

    # 제목 교체
    set_text_run(shapes[23], '윤리경영 기준 적용', size=13, bold=True, color=DARK_GREEN)

    # 내용 교체
    new_content = (
        '하나금융그룹 윤리경영 원칙에 따라 AI 사용 시에도 반드시 준수해야 합니다. '
        '① 고객 개인정보·계좌번호·MNPI(미공개 중요 정보) AI 입력 절대 금지  '
        '② 발간 전 리서치 초안·대외비 자료 입력 금지  '
        '③ AI를 이용한 부당 이익 취득 금지'
    )
    set_text_run(shapes[24], new_content)

    print('  [Slide 6] item 02 → 윤리경영 기준 적용 완료')


# ══════════════════════════════════════════════════════════
#  2. Slide 18 — 원본 5개 항목 복원, 이모티콘 제거
# ══════════════════════════════════════════════════════════
def rebuild_slide18_items():
    slide = prs.slides[17]

    # 제거 대상: 중복 배경 shapes [4]-[15] + 기존 content [23]-[46]
    remove_shapes_by_idx(slide, list(range(4, 16)) + list(range(23, 47)))

    # 원본 5개 항목 (원문 그대로)
    items = [
        ('AI는 속도, 인사이트는 여러분의 몫',
         'AI는 자료의 정리 속도를 높여주지만, 의미 있는 인사이트를 찾는 것은 여러분의 역할입니다.'),
        ('지식 베이스를 꾸준히 쌓으세요',
         '신뢰성 있는 증권사 리포트, 산업 분석 자료 등을 읽고, AI 엔진에 꾸준히 추가해보세요.'),
        ('경쟁사 자료도 함께 비교하세요',
         '경쟁사의 자료도 입력하여 비교해야 업종 내 상대적인 우위를 확인 가능합니다.'),
        ('기초 질문을 두려워 말 것',
         'AI에게 질문할 때 가장 좋은 점: 기초적이고 쉬운 질문을 부끄러워 마세요. 궁금한 모든 것들을 물어보세요.'),
        ('좋은 질문이 진정한 AI 활용의 가치',
         '좋은 질문을 던지고 의미 있는 인사이트를 찾아가는 것이 AI 활용의 진정한 가치입니다.'),
    ]

    # 5개 항목 배치: T=1.20 ~ T=7.00 (= 5.80"), 항목당 1.04" (H=0.94, gap=0.10)
    item_h = Inches(0.94)
    gap    = Inches(0.10)
    step   = item_h + gap  # 1.04"
    start_t = Inches(1.20)
    lx     = Inches(4.50)
    item_w = Inches(5.08)

    for i, (heading, body) in enumerate(items):
        t = start_t + i * step

        # 배경
        bg = add_rect(slide, lx, t, item_w, item_h, LIGHT_GRAY)
        # 좌측 상단 녹색 바
        add_rect(slide, lx, t, item_w, Inches(0.04), GREEN)

        # 제목 (이모티콘 없이 바로 텍스트)
        add_tb(slide, lx + Inches(0.10), t + Inches(0.06),
               item_w - Inches(0.18), Inches(0.32),
               heading, size=12, bold=True, color=DARK_GREEN)

        # 구분선
        add_rect(slide, lx + Inches(0.10), t + Inches(0.42),
                 item_w - Inches(0.18), Inches(0.01), GRAY)

        # 내용
        add_tb(slide, lx + Inches(0.10), t + Inches(0.48),
               item_w - Inches(0.18), Inches(0.42),
               body, size=10, color=DARK_TEXT, wrap=True)

    print('  [Slide 18] 원본 5개 항목 복원 완료 (이모티콘 제거)')


# ══════════════════════════════════════════════════════════
#  3. 전체 슬라이드 교시 → 목차 섹션 번호로 교체
#  목차: 1=이론, 2=실습, 3=마무리
# ══════════════════════════════════════════════════════════
GYOSI_MAP = {
    # (slide_idx 0-based, shape_idx): (old_text, new_text)
    (9,  2): ('1교시 핵심 정리',                     '1. 이론 · 핵심 정리'),
    (15, 22): ('2교시 · AI 활용 실습',               '2. 실습 · AI 활용 실습'),
    (18, 2): ('3교시',                               '2. 실습 · 엑셀 데이터 분석 실습'),
    (19, 1): ('3교시 · 엑셀 데이터 분석 실습',       '2. 실습 · 엑셀 데이터 분석 실습'),
    (20, 1): ('3교시 · 엑셀 데이터 분석 실습',       '2. 실습 · 엑셀 데이터 분석 실습'),
    (22, 1): ('4교시 · 나만의 챗봇 업무 비서 만들기','2. 실습 · 나만의 챗봇 업무 비서 만들기'),
    (23, 1): ('4교시 · 나만의 챗봇 업무 비서 만들기','2. 실습 · 나만의 챗봇 업무 비서 만들기'),
}


def fix_gyosi():
    for (si, shi), (old, new) in GYOSI_MAP.items():
        slide  = prs.slides[si]
        shapes = list(slide.shapes)
        if shi >= len(shapes):
            print(f'  [Slide {si+1}] shape[{shi}] 없음 — 스킵')
            continue
        shape  = shapes[shi]
        if not shape.has_text_frame:
            print(f'  [Slide {si+1}] shape[{shi}] 텍스트 없음 — 스킵')
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if old in run.text:
                    run.text = run.text.replace(old, new)
                    print(f'  [Slide {si+1}] "{old}" → "{new}"')
                    break
            else:
                # run이 없거나 텍스트가 para에 split된 경우 para 전체 교체
                full = para.text
                if old in full:
                    # 기존 runs 제거 후 재생성
                    for run in para.runs:
                        run.text = ''
                    if para.runs:
                        para.runs[0].text = new
                    else:
                        from pptx.oxml.ns import qn
                        r_el = para._p.makeelement(qn('a:r'))
                        rPr = para._p.makeelement(qn('a:rPr'), {'lang': 'ko-KR'})
                        r_el.insert(0, rPr)
                        t_el = para._p.makeelement(qn('a:t'))
                        t_el.text = new
                        r_el.append(t_el)
                        para._p.append(r_el)
                    print(f'  [Slide {si+1}] (para rebuild) "{old}" → "{new}"')

    print('  [전체] 교시 → 섹션 번호 교체 완료')


# ══════════════════════════════════════════════════════════
if __name__ == '__main__':
    fix_slide6()
    rebuild_slide18_items()
    fix_gyosi()
    prs.save(DST)
    print(f'\n✓ 저장: {DST}')
