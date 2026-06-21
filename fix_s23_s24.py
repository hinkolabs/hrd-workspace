"""
Slide 23 — Alli Works & RAG 설명 → 기존 PPT 2×2 카드 스타일로 재구성
Slide 24 — 챗봇 빌더 스크린샷 삽입 (빌더 플로우 + 대화 결과)
"""
import sys, os
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SRC  = r'C:\Users\jay\Downloads\AI교육_완성판_v8.pptx'
DST  = r'C:\Users\jay\Downloads\AI교육_완성판_v9.pptx'
IMG1 = r'C:\Users\jay\.cursor\projects\c-dev-hrd-workspace\assets\c__Users_jay_AppData_Roaming_Cursor_User_workspaceStorage_14cd220dbe531e9d0f57e5f249b63965_images_image-a4a71571-cdb5-48cc-a63c-08384dd8f246.png'
IMG2 = r'C:\Users\jay\.cursor\projects\c-dev-hrd-workspace\assets\c__Users_jay_AppData_Roaming_Cursor_User_workspaceStorage_14cd220dbe531e9d0f57e5f249b63965_images_image-44dd6f62-b146-4471-a67f-8e5d2c9e4bc4.png'

prs = Presentation(SRC)

GREEN       = RGBColor(0x00, 0x91, 0x78)
DARK_GREEN  = RGBColor(0x00, 0x4E, 0x42)
LIGHT_GREEN = RGBColor(0xD7, 0xED, 0xE6)
LIGHT_GRAY  = RGBColor(0xF7, 0xF9, 0xF8)
DARK_TEXT   = RGBColor(0x22, 0x22, 0x22)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY        = RGBColor(0xCC, 0xCC, 0xCC)
FONT        = '맑은 고딕'


def rgb_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def add_rect(slide, l, t, w, h, rgb):
    s = slide.shapes.add_shape(1, l, t, w, h)
    rgb_fill(s, rgb)
    return s


def add_tb(slide, l, t, w, h, text, size=12, bold=False,
           color=DARK_TEXT, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = FONT
    return tb


def remove_all_content(slide):
    """헤더/푸터 요소(녹색바, 저작권, 페이지번호)만 남기고 모두 제거"""
    keep_texts = {
        'Copyright', '| 23', '| 24',
    }
    sp_tree = slide.shapes._spTree
    to_remove = []
    for shape in slide.shapes:
        # 항상 보존: 첫 번째 녹색 세로바(L=0, W≈0.14), 하단 녹색바(T≈7.05), 저작권, 페이지번호
        l = shape.left / 914400
        t = shape.top / 914400
        w = shape.width / 914400
        # 좌측 녹색 세로바
        if l < 0.02 and w < 0.20 and t < 0.5:
            continue
        # 하단 녹색바 + 저작권 + 페이지
        if t > 6.90:
            continue
        to_remove.append(shape._element)
    for el in to_remove:
        el.getparent().remove(el)


def add_card(slide, l, t, w, h, num_str, title, body, num_size=18):
    """기존 PPT 스타일 카드: 좌측 녹색 번호 스트라이프 + 제목 + 구분선 + 본문"""
    # 배경
    add_rect(slide, l, t, w, h, LIGHT_GRAY)
    # 번호 영역 (녹색)
    add_rect(slide, l, t, Inches(0.65), h, GREEN)
    add_tb(slide, l, t, Inches(0.65), h,
           num_str, size=num_size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    inner_l = l + Inches(0.72)
    inner_w = w - Inches(0.80)
    # 제목
    add_tb(slide, inner_l, t + Inches(0.10), inner_w, Inches(0.38),
           title, size=13, bold=True, color=DARK_GREEN)
    # 구분선
    add_rect(slide, inner_l, t + Inches(0.52), inner_w, Inches(0.01), GRAY)
    # 본문
    add_tb(slide, inner_l, t + Inches(0.58), inner_w, h - Inches(0.66),
           body, size=10, color=DARK_TEXT, wrap=True)


# ══════════════════════════════════════════════════════════
#  Slide 23 — Alli Works & RAG 2×2 카드 레이아웃
# ══════════════════════════════════════════════════════════
def rebuild_slide23():
    slide = prs.slides[22]
    remove_all_content(slide)

    # ── 헤더 재생성
    add_tb(slide, Inches(0.47), Inches(0.24), Inches(7.87), Inches(0.28),
           '2. 실습 · 나만의 챗봇 업무 비서 만들기',
           size=10, color=GREEN)
    add_tb(slide, Inches(0.47), Inches(0.48), Inches(9.06), Inches(0.54),
           'Alli Works & RAG — 기본 개념 이해',
           size=24, bold=True, color=DARK_TEXT)
    add_rect(slide, Inches(0.47), Inches(1.07), Inches(9.06), Inches(0.02), GREEN)

    # ── 도입 배너 (한 줄 설명)
    add_rect(slide, Inches(0.47), Inches(1.20), Inches(9.06), Inches(0.40), DARK_GREEN)
    add_rect(slide, Inches(0.47), Inches(1.20), Inches(0.06), Inches(0.40), GREEN)
    add_tb(slide, Inches(0.62), Inches(1.25), Inches(8.80), Inches(0.32),
           'Alli Works는 사내 문서를 업로드해 AI 챗봇을 만드는 플랫폼, RAG는 그 핵심 기술입니다.',
           size=11, color=WHITE)

    # ── 2×2 카드
    card_w = Inches(4.40)
    card_h = Inches(1.55)
    gap    = Inches(0.18)
    row1_t = Inches(1.73)
    row2_t = row1_t + card_h + gap
    col1_l = Inches(0.47)
    col2_l = col1_l + card_w + gap

    cards = [
        (col1_l, row1_t, '01', 'Alli Works란?',
         '코딩·AI 전문 지식 없이도 누구나 만드는 사내 AI 업무 비서 플랫폼.\n'
         '사내 문서를 업로드하면 그 내용을 바탕으로 질문에 답하는 챗봇이 자동 생성됩니다.'),
        (col2_l, row1_t, '02', '무엇을 만들 수 있나?',
         '• 반복 업무 FAQ 자동 답변\n'
         '• 내부 규정·프로세스 안내\n'
         '• 신입사원 온보딩 안내\n'
         '• 회의록·보고서 초안 작성 보조'),
        (col1_l, row2_t, '03', 'RAG(검색 증강 생성)란?',
         'AI가 질문을 받을 때마다 업로드된 사내 문서를 실시간으로 검색해\n'
         '관련 단락을 찾아 답변에 활용하는 기술.\n'
         '미리 외우지 않고 "그때그때 찾아서" 답합니다.'),
        (col2_l, row2_t, '04', 'RAG 동작 4단계',
         '① 질문 입력  →  ② 문서 검색 (관련 단락 추출)\n'
         '→  ③ 컨텍스트 결합 (질문 + 단락을 AI에 전달)\n'
         '→  ④ 답변 생성 (문서 기반의 정확한 답변 + 출처)'),
    ]
    for lx, ty, num, title, body in cards:
        add_card(slide, lx, ty, card_w, card_h, num, title, body)

    # ── 핵심 메모 배너
    note_t = row2_t + card_h + Inches(0.15)
    add_rect(slide, Inches(0.47), note_t, Inches(9.06), Inches(0.36), LIGHT_GREEN)
    add_rect(slide, Inches(0.47), note_t, Inches(0.06), Inches(0.36), GREEN)
    add_tb(slide, Inches(0.62), note_t + Inches(0.06), Inches(8.80), Inches(0.28),
           '핵심: AI가 문서를 "암기(재학습)"하는 것이 아닙니다. 질문마다 문서를 검색해서 그때그때 참고하는 구조입니다.',
           size=10, color=DARK_GREEN, italic=True)

    # ── 핵심 기능 3개 카드 (하단)
    feat_t = note_t + Inches(0.46)
    feat_h = Inches(1.55)
    feat_w = Inches(2.90)
    feat_gap = Inches(0.18)
    feats = [
        ('RAG 기반 실시간 검색',
         '질문이 오면 업로드된 문서에서 관련 단락을 즉시 검색·추출.\n"외운 것"이 아닌 "찾은 것"으로 답합니다.'),
        ('업무 맞춤 챗봇 제작',
         '부서 FAQ·프로세스·규정집 업로드 후 시스템 프롬프트로\n어조·답변 범위 설정. 코딩 없이 전용 챗봇 완성.'),
        ('보안 및 접근 통제',
         'MNPI·개인정보 없는 사내 자료만 업로드 허용.\n부서 권한별 접근 범위 설정 가능.'),
    ]
    for i, (ftitle, fbody) in enumerate(feats):
        fx = Inches(0.47) + i * (feat_w + feat_gap)
        add_rect(slide, fx, feat_t, feat_w, feat_h, LIGHT_GRAY)
        add_rect(slide, fx, feat_t, feat_w, Inches(0.04), GREEN)
        add_tb(slide, fx + Inches(0.10), feat_t + Inches(0.08),
               feat_w - Inches(0.18), Inches(0.34),
               ftitle, size=11, bold=True, color=DARK_GREEN)
        add_rect(slide, fx + Inches(0.10), feat_t + Inches(0.46),
                 feat_w - Inches(0.18), Inches(0.01), GRAY)
        add_tb(slide, fx + Inches(0.10), feat_t + Inches(0.52),
               feat_w - Inches(0.18), feat_h - Inches(0.58),
               fbody, size=9.5, color=DARK_TEXT, wrap=True)

    print('  [Slide 23] 2×2 카드 + 핵심 기능 3카드 재구성 완료')


# ══════════════════════════════════════════════════════════
#  Slide 24 — 챗봇 완성 예시 스크린샷 삽입
# ══════════════════════════════════════════════════════════
def rebuild_slide24():
    slide = prs.slides[23]
    remove_all_content(slide)

    # ── 헤더
    add_tb(slide, Inches(0.47), Inches(0.24), Inches(7.87), Inches(0.28),
           '2. 실습 · 나만의 챗봇 업무 비서 만들기',
           size=10, color=GREEN)
    add_tb(slide, Inches(0.47), Inches(0.48), Inches(9.06), Inches(0.54),
           '나만의 업무 비서 챗봇 — 완성 예시',
           size=24, bold=True, color=DARK_TEXT)
    add_rect(slide, Inches(0.47), Inches(1.07), Inches(9.06), Inches(0.02), GREEN)

    # ── 소제목 라인
    add_rect(slide, Inches(0.47), Inches(1.18), Inches(9.06), Inches(0.34), DARK_GREEN)
    add_rect(slide, Inches(0.47), Inches(1.18), Inches(0.06), Inches(0.34), GREEN)
    add_tb(slide, Inches(0.62), Inches(1.22), Inches(8.80), Inches(0.26),
           'Alli Works에서 챗봇 플로우를 구성하고, 실제 업무 내용을 질문하면 아래와 같이 답변합니다.',
           size=10.5, color=WHITE)

    # ── 이미지 레이아웃 계산
    # 사용 가능 영역: T=1.62 ~ T=6.85, L=0.47 ~ L=9.53
    # 좌: 플로우 다이어그램 (가로가 더 넓음) - W5.10 H3.70
    # 우: 대화 결과 (세로가 더 김) - W3.70 H4.95
    img_top  = Inches(1.62)
    gap      = Inches(0.18)

    left_w   = Inches(5.10)
    left_h   = Inches(3.70)
    right_w  = Inches(3.70)
    right_h  = Inches(4.95)

    left_l   = Inches(0.47)
    right_l  = left_l + left_w + gap

    img1_exists = os.path.exists(IMG1)
    img2_exists = os.path.exists(IMG2)

    if img1_exists:
        slide.shapes.add_picture(IMG1, left_l, img_top, left_w, left_h)
        print(f'    → 이미지1(플로우) 삽입: {left_w/914400:.1f}" × {left_h/914400:.1f}"')
    else:
        # 이미지가 없으면 placeholder 박스
        ph = add_rect(slide, left_l, img_top, left_w, left_h, LIGHT_GRAY)
        add_rect(slide, left_l, img_top, left_w, Inches(0.04), GREEN)
        add_tb(slide, left_l + Inches(0.18), img_top + Inches(1.50),
               left_w - Inches(0.35), Inches(0.70),
               '[ 챗봇 빌더 플로우 스크린샷 ]',
               size=14, color=GRAY, align=PP_ALIGN.CENTER, bold=True)
        add_tb(slide, left_l + Inches(0.18), img_top + Inches(2.20),
               left_w - Inches(0.35), Inches(0.40),
               '이 공간에 직접 이미지를 삽입해주세요',
               size=10, color=GRAY, align=PP_ALIGN.CENTER, italic=True)
        print(f'    → 이미지1 없음: Placeholder 박스 생성')

    if img2_exists:
        slide.shapes.add_picture(IMG2, right_l, img_top, right_w, right_h)
        print(f'    → 이미지2(대화) 삽입: {right_w/914400:.1f}" × {right_h/914400:.1f}"')
    else:
        ph2 = add_rect(slide, right_l, img_top, right_w, right_h, LIGHT_GRAY)
        add_rect(slide, right_l, img_top, right_w, Inches(0.04), GREEN)
        add_tb(slide, right_l + Inches(0.18), img_top + Inches(2.00),
               right_w - Inches(0.35), Inches(0.70),
               '[ 챗봇 대화 결과 스크린샷 ]',
               size=14, color=GRAY, align=PP_ALIGN.CENTER, bold=True)
        add_tb(slide, right_l + Inches(0.18), img_top + Inches(2.70),
               right_w - Inches(0.35), Inches(0.40),
               '이 공간에 직접 이미지를 삽입해주세요',
               size=10, color=GRAY, align=PP_ALIGN.CENTER, italic=True)
        print(f'    → 이미지2 없음: Placeholder 박스 생성')

    # ── 이미지 하단 캡션
    cap1_t = img_top + left_h + Inches(0.08)
    cap2_t = img_top + right_h + Inches(0.08)

    add_rect(slide, left_l, cap1_t, left_w, Inches(0.32), GREEN)
    add_tb(slide, left_l + Inches(0.10), cap1_t + Inches(0.05),
           left_w - Inches(0.18), Inches(0.25),
           '① 챗봇 플로우 구성 화면  (Start → 답변 생성 → 추가 질문)', 
           size=9.5, color=WHITE, bold=False)

    add_rect(slide, right_l, cap2_t, right_w, Inches(0.32), GREEN)
    add_tb(slide, right_l + Inches(0.10), cap2_t + Inches(0.05),
           right_w - Inches(0.18), Inches(0.25),
           '② 실제 업무 내용 질문 결과  (5월 7일 업무 내용 요약 예시)',
           size=9.5, color=WHITE, bold=False)

    print('  [Slide 24] 챗봇 예시 슬라이드 재구성 완료')


# ══════════════════════════════════════════════════════════
if __name__ == '__main__':
    print(f'이미지1 존재: {os.path.exists(IMG1)}')
    print(f'이미지2 존재: {os.path.exists(IMG2)}\n')

    rebuild_slide23()
    rebuild_slide24()
    prs.save(DST)
    print(f'\n✓ 저장: {DST}')
