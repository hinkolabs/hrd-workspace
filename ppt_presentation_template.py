"""
ppt_presentation_template.py — 하나증권 발표용 멀티슬라이드 PPT 템플릿

슬라이드 구성 (10장):
  Slide 1 : 표지 (Cover)
  Slide 2 : 목차 (Agenda)
  Slide 3 : 섹션 인트로 Ⅰ (Section divider)
  Slide 4 : 현황/배경 — 2단 텍스트 (Two-column)
  Slide 5 : 핵심 KPI 수치 (KPI stats)
  Slide 6 : 섹션 인트로 Ⅱ (Section divider)
  Slide 7 : 데이터 표 (Table)
  Slide 8 : 비교 분석 (Comparison)
  Slide 9 : 실행 계획 — 프로세스 (Process steps)
  Slide 10: 마무리 (Closing)

사용법:
  python ppt_presentation_template.py
  python ppt_presentation_template.py --title "2026 HRD 전략" --out "발표자료.pptx"

또는 코드에서:
  from ppt_presentation_template import build_presentation
  build_presentation(config, output_path)
"""

import argparse, os
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pptx.oxml.ns import qn

from ppt_theme import (
    new_presentation, blank_slide,
    add_rect, add_textbox, add_multi_para_textbox, add_table,
    SLIDE_W, SLIDE_H,
    PRIMARY, ACCENT, NAVY, DARK, DARK_BG, GRAY, MGRAY, LGRAY, LGRAY2, WHITE,
    FONT_TITLE, FONT_BODY, LOGO_PATH,
)

# ── 공통 상수 ─────────────────────────────────────────────────────────────────────
HEADER_H  = Inches(0.85)
PAD_X     = Inches(0.4)
BODY_Y    = Inches(0.85 + 0.15)      # 헤더 아래 시작점
BODY_H    = SLIDE_H - BODY_Y - Inches(0.3)
SECTION_COLORS = [PRIMARY, NAVY, ACCENT, RGBColor(0x00, 0x7A, 0x77)]
CARD_SHADOW = False   # python-pptx shadow 미지원 — 생략


# ══════════════════════════════════════════════════════════════════════════════════
# 공통 크롬 (모든 슬라이드에 적용)
# ══════════════════════════════════════════════════════════════════════════════════

def _add_header(slide, title: str, slide_num: int | None = None) -> None:
    """공통 헤더 바 (하나그린 + 레드 액센트)"""
    add_rect(slide, 0, 0, SLIDE_W, HEADER_H, fill=PRIMARY)
    add_rect(slide, 0, 0, Inches(0.14), HEADER_H, fill=ACCENT)
    add_rect(slide, 0, HEADER_H - Inches(0.06), SLIDE_W, Inches(0.06), fill=ACCENT)

    add_textbox(slide, Inches(0.22), 0, SLIDE_W - Inches(1.5), HEADER_H,
                title, font_size=Pt(18), bold=True, color=WHITE,
                font_face=FONT_TITLE, align=PP_ALIGN.LEFT)

    if slide_num is not None:
        add_textbox(slide, SLIDE_W - Inches(1.2), 0, Inches(1.1), HEADER_H,
                    str(slide_num), font_size=Pt(11), color=LGRAY2,
                    font_face=FONT_BODY, align=PP_ALIGN.RIGHT)

    # 로고
    if os.path.exists(LOGO_PATH):
        try:
            slide.shapes.add_picture(LOGO_PATH, SLIDE_W - Inches(1.35), Inches(0.1),
                                     height=Inches(0.62))
        except Exception:
            pass


def _add_footer(slide, dept: str = "", date: str = "") -> None:
    """공통 푸터 바"""
    fy = SLIDE_H - Inches(0.42)
    add_rect(slide, 0, fy, SLIDE_W, Inches(0.42), fill=DARK)
    if dept:
        add_textbox(slide, Inches(0.3), fy, Inches(6), Inches(0.42),
                    dept, font_size=Pt(9), color=MGRAY, font_face=FONT_BODY)
    if date:
        add_textbox(slide, SLIDE_W - Inches(3.3), fy, Inches(3.1), Inches(0.42),
                    date, font_size=Pt(9), color=MGRAY, font_face=FONT_BODY,
                    align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════════
# 슬라이드 빌더
# ══════════════════════════════════════════════════════════════════════════════════

def _slide_cover(prs, config: dict) -> None:
    """Slide 1 — 표지"""
    slide = blank_slide(prs)

    # 배경: 하나그린 풀블리드
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=PRIMARY)
    # 하단 어두운 오버레이
    add_rect(slide, 0, SLIDE_H - Inches(2.0), SLIDE_W, Inches(2.0),
             fill=RGBColor(0, 0, 0))  # 투명도 미지원 → 단색

    # 장식용 원
    for x, y, r in [(9.1, -0.35, 1.8), (1.5, 6.0, 2.2)]:
        shape = slide.shapes.add_shape(3,  # ellipse
                                        Inches(x - r), Inches(y - r),
                                        Inches(r * 2), Inches(r * 2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.fill.background()

    # 메인 타이틀
    add_textbox(slide, Inches(0.9), Inches(1.8), Inches(11.5), Inches(2.0),
                config.get("title", "[프레젠테이션 제목]"),
                font_size=Pt(40), bold=True, color=WHITE, font_face=FONT_TITLE,
                align=PP_ALIGN.LEFT)

    # 서브타이틀
    subtitle = config.get("subtitle", "[부제목 또는 보고부서]")
    add_textbox(slide, Inches(0.9), Inches(3.9), Inches(10), Inches(0.8),
                subtitle, font_size=Pt(20), color=LGRAY2, font_face=FONT_BODY,
                align=PP_ALIGN.LEFT)

    # 날짜 + 부서
    info_line = "  |  ".join(filter(None, [config.get("dept", ""), config.get("date", "")]))
    if info_line:
        add_textbox(slide, Inches(0.9), Inches(5.0), Inches(10), Inches(0.5),
                    info_line, font_size=Pt(14), color=MGRAY, font_face=FONT_BODY,
                    align=PP_ALIGN.LEFT)

    # 로고 (하단 우측)
    if os.path.exists(LOGO_PATH):
        try:
            slide.shapes.add_picture(LOGO_PATH, SLIDE_W - Inches(2.5), Inches(6.6),
                                     height=Inches(0.7))
        except Exception:
            pass


def _slide_agenda(prs, config: dict) -> None:
    """Slide 2 — 목차"""
    slide = blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    _add_header(slide, "목  차", 2)
    _add_footer(slide, config.get("dept", ""), config.get("date", ""))

    sections = config.get("sections", [])
    n = len(sections)
    item_h = Inches(0.72)
    total_h = n * item_h + (n - 1) * Inches(0.12)
    start_y = BODY_Y + (BODY_H - total_h) / 2

    for i, sect in enumerate(sections):
        y = start_y + i * (item_h + Inches(0.12))
        accent = SECTION_COLORS[i % len(SECTION_COLORS)]

        # 번호 박스
        num_box = slide.shapes.add_shape(1, PAD_X, y, Inches(0.72), item_h)
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = accent
        num_box.line.fill.background()
        tf = num_box.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = ["Ⅰ", "Ⅱ", "Ⅲ", "Ⅳ", "Ⅴ"][i] if i < 5 else str(i + 1)
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = FONT_TITLE

        # 섹션 제목
        add_textbox(slide, PAD_X + Inches(0.85), y, SLIDE_W - PAD_X - Inches(1.0), item_h,
                    sect.get("title", f"섹션 {i + 1}"),
                    font_size=Pt(16), bold=True, color=DARK, font_face=FONT_TITLE,
                    align=PP_ALIGN.LEFT)

        # 구분선
        if i < n - 1:
            add_rect(slide, PAD_X + Inches(0.85), y + item_h + Inches(0.04),
                     SLIDE_W - PAD_X * 2 - Inches(0.85), Inches(0.03), fill=LGRAY2)


def _slide_section(prs, number: str, title: str, subtitle: str, config: dict,
                   slide_num: int, accent: RGBColor = PRIMARY) -> None:
    """섹션 인트로 슬라이드 (다크 그린 배경)"""
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=DARK_BG)

    # 장식 원
    c_shape = slide.shapes.add_shape(3,
                                      SLIDE_W - Inches(4.5), Inches(-1.0),
                                      Inches(5.0), Inches(5.0))
    c_shape.fill.solid()
    c_shape.fill.fore_color.rgb = WHITE
    c_shape.line.fill.background()

    # 섹션 번호
    add_textbox(slide, PAD_X, Inches(2.2), Inches(3), Inches(1.2),
                number, font_size=Pt(64), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                font_face=FONT_TITLE, align=PP_ALIGN.LEFT)

    # 섹션 제목
    add_textbox(slide, PAD_X, Inches(3.4), Inches(9), Inches(1.0),
                title, font_size=Pt(28), bold=True, color=WHITE,
                font_face=FONT_TITLE, align=PP_ALIGN.LEFT)

    # 구분선
    add_rect(slide, PAD_X, Inches(4.5), Inches(6), Inches(0.05), fill=ACCENT)

    if subtitle:
        add_textbox(slide, PAD_X, Inches(4.65), Inches(9), Inches(0.6),
                    subtitle, font_size=Pt(14), color=LGRAY2, font_face=FONT_BODY)

    # 슬라이드 번호 (우하단)
    add_textbox(slide, SLIDE_W - Inches(1.0), SLIDE_H - Inches(0.5), Inches(0.8), Inches(0.4),
                str(slide_num), font_size=Pt(10), color=MGRAY, font_face=FONT_BODY,
                align=PP_ALIGN.RIGHT)

    # 로고
    if os.path.exists(LOGO_PATH):
        try:
            slide.shapes.add_picture(LOGO_PATH, SLIDE_W - Inches(2.3), SLIDE_H - Inches(1.0),
                                     height=Inches(0.6))
        except Exception:
            pass


def _slide_two_column(prs, config: dict, slide_data: dict, slide_num: int) -> None:
    """2단 비교 슬라이드"""
    slide = blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    _add_header(slide, slide_data.get("title", "[슬라이드 제목]"), slide_num)
    _add_footer(slide, config.get("dept", ""), config.get("date", ""))

    col_w = (SLIDE_W - PAD_X * 2 - Inches(0.12)) / 2
    footer_h = Inches(0.42)
    body_h = SLIDE_H - BODY_Y - footer_h - Inches(0.1)
    HDR_H = Inches(0.38)
    ACCENTS = [PRIMARY, NAVY]

    for ci, col_data in enumerate(["left", "right"]):
        cx = PAD_X + ci * (col_w + Inches(0.12))
        col = slide_data.get(col_data, {})
        accent = ACCENTS[ci]

        # 칼럼 카드
        add_rect(slide, cx, BODY_Y, col_w, body_h, fill=LGRAY)
        add_rect(slide, cx, BODY_Y, col_w, HDR_H, fill=accent)

        add_textbox(slide, cx + Inches(0.12), BODY_Y, col_w - Inches(0.24), HDR_H,
                    col.get("heading", f"{'현황' if ci == 0 else '분석'}"),
                    font_size=Pt(13), bold=True, color=WHITE, font_face=FONT_TITLE)

        bullets = col.get("bullets", ["• 내용을 입력하세요", "• 내용을 입력하세요"])
        text = "\n".join(
            f"• {b}" if not b.startswith("•") else b for b in bullets
        )
        add_textbox(slide, cx + Inches(0.15), BODY_Y + HDR_H + Inches(0.1),
                    col_w - Inches(0.3), body_h - HDR_H - Inches(0.15),
                    text, font_size=Pt(10), color=DARK, font_face=FONT_BODY,
                    align=PP_ALIGN.LEFT, wrap=True)


def _slide_kpi(prs, config: dict, slide_data: dict, slide_num: int) -> None:
    """KPI 수치 슬라이드 (stats 레이아웃)"""
    slide = blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    _add_header(slide, slide_data.get("title", "핵심 수치"), slide_num)
    _add_footer(slide, config.get("dept", ""), config.get("date", ""))

    kpis = slide_data.get("kpis", [])
    n = max(len(kpis), 1)
    footer_h = Inches(0.42)
    body_h = SLIDE_H - BODY_Y - footer_h - Inches(0.1)
    gap = Inches(0.15)
    card_w = (SLIDE_W - PAD_X * 2 - gap * (n - 1)) / n

    for i, kpi in enumerate(kpis):
        cx = PAD_X + i * (card_w + gap)
        is_pos = kpi.get("positive", True)
        accent = PRIMARY if is_pos else ACCENT

        add_rect(slide, cx, BODY_Y, card_w, body_h, fill=LGRAY)
        add_rect(slide, cx, BODY_Y, card_w, Inches(0.06), fill=accent)

        # 라벨
        add_textbox(slide, cx + Inches(0.15), BODY_Y + Inches(0.12),
                    card_w - Inches(0.3), Inches(0.32),
                    kpi.get("label", f"지표 {i+1}"),
                    font_size=Pt(11), bold=True, color=GRAY, font_face=FONT_TITLE)

        # 수치
        add_textbox(slide, cx + Inches(0.15), BODY_Y + body_h * 0.28,
                    card_w - Inches(0.3), body_h * 0.38,
                    kpi.get("value", "—"),
                    font_size=Pt(36), bold=True, color=accent, font_face=FONT_TITLE,
                    align=PP_ALIGN.CENTER)

        # 캡션
        if kpi.get("caption"):
            add_textbox(slide, cx + Inches(0.15), BODY_Y + body_h * 0.72,
                        card_w - Inches(0.3), Inches(0.35),
                        kpi["caption"], font_size=Pt(9.5), color=GRAY,
                        font_face=FONT_BODY, align=PP_ALIGN.CENTER)

        # 변화 배지
        if kpi.get("change"):
            badge_c = RGBColor(0x16, 0xA3, 0x4A) if is_pos else RGBColor(0xDC, 0x26, 0x26)
            add_textbox(slide, cx + Inches(0.15), BODY_Y + body_h - Inches(0.38),
                        card_w - Inches(0.3), Inches(0.32),
                        kpi["change"], font_size=Pt(11), bold=True, color=badge_c,
                        font_face=FONT_BODY, align=PP_ALIGN.CENTER)


def _slide_table(prs, config: dict, slide_data: dict, slide_num: int) -> None:
    """표 슬라이드"""
    slide = blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    _add_header(slide, slide_data.get("title", "[표 슬라이드]"), slide_num)
    _add_footer(slide, config.get("dept", ""), config.get("date", ""))

    headers = slide_data.get("headers", ["항목", "현황", "목표", "비고"])
    rows    = slide_data.get("rows",    [
        ["[항목 1]", "[현황 값]", "[목표 값]", ""],
        ["[항목 2]", "[현황 값]", "[목표 값]", ""],
        ["[항목 3]", "[현황 값]", "[목표 값]", ""],
        ["[항목 4]", "[현황 값]", "[목표 값]", ""],
    ])
    footer_h = Inches(0.42)
    tbl_h = SLIDE_H - BODY_Y - footer_h - Inches(0.2)
    add_table(slide, PAD_X, BODY_Y, SLIDE_W - PAD_X * 2, tbl_h, headers, rows)


def _slide_comparison(prs, config: dict, slide_data: dict, slide_num: int) -> None:
    """비교 분석 슬라이드 (VS 스타일)"""
    slide = blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    _add_header(slide, slide_data.get("title", "[비교 분석]"), slide_num)
    _add_footer(slide, config.get("dept", ""), config.get("date", ""))

    footer_h = Inches(0.42)
    body_h = SLIDE_H - BODY_Y - footer_h - Inches(0.1)
    vs_w = Inches(0.7)
    col_w = (SLIDE_W - PAD_X * 2 - vs_w - Inches(0.24)) / 2

    for ci, (side, accent, bg) in enumerate([
        ("left",  PRIMARY, LGRAY),
        ("right", ACCENT,  RGBColor(0xFD, 0xED, 0xED)),
    ]):
        col = slide_data.get(side, {})
        cx = PAD_X + ci * (col_w + vs_w + Inches(0.12) * (ci))
        if ci == 1:
            cx = PAD_X + col_w + vs_w + Inches(0.12)

        add_rect(slide, cx, BODY_Y, col_w, body_h, fill=bg)
        add_rect(slide, cx, BODY_Y, col_w, Inches(0.4), fill=accent)
        add_textbox(slide, cx + Inches(0.12), BODY_Y, col_w - Inches(0.24), Inches(0.4),
                    col.get("heading", "현재" if ci == 0 else "목표"),
                    font_size=Pt(13), bold=True, color=WHITE, font_face=FONT_TITLE)

        bullets = col.get("bullets", ["• 내용 1", "• 내용 2", "• 내용 3"])
        text = "\n".join(
            f"• {b}" if not b.startswith("•") else b for b in bullets
        )
        add_textbox(slide, cx + Inches(0.15), BODY_Y + Inches(0.48),
                    col_w - Inches(0.3), body_h - Inches(0.58),
                    text, font_size=Pt(10.5), color=DARK, font_face=FONT_BODY, wrap=True)

    # VS 뱃지
    vs_x = PAD_X + col_w + Inches(0.06)
    vs_y = BODY_Y + (body_h - vs_w) / 2
    vs_shape = slide.shapes.add_shape(3, vs_x, vs_y, vs_w, vs_w)
    vs_shape.fill.solid()
    vs_shape.fill.fore_color.rgb = DARK
    vs_shape.line.fill.background()
    tf = vs_shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = slide_data.get("vs_label", "VS")
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT_TITLE


def _slide_process(prs, config: dict, slide_data: dict, slide_num: int) -> None:
    """프로세스 단계 슬라이드"""
    slide = blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    _add_header(slide, slide_data.get("title", "[실행 계획]"), slide_num)
    _add_footer(slide, config.get("dept", ""), config.get("date", ""))

    steps = slide_data.get("steps", [])[:5]
    n = max(len(steps), 1)
    footer_h = Inches(0.42)
    body_h = SLIDE_H - BODY_Y - footer_h - Inches(0.1)
    gap = Inches(0.1)
    step_w = (SLIDE_W - PAD_X * 2 - gap * (n - 1)) / n
    CIRCLE_R = Inches(0.42)

    for i, step in enumerate(steps):
        sx = PAD_X + i * (step_w + gap)
        accent = SECTION_COLORS[i % len(SECTION_COLORS)]

        # 카드 배경
        add_rect(slide, sx, BODY_Y, step_w, body_h, fill=LGRAY)
        add_rect(slide, sx, BODY_Y, step_w, Inches(0.05), fill=accent)

        # 번호 원
        cx = sx + (step_w - CIRCLE_R) / 2
        cy = BODY_Y + Inches(0.18)
        c_shape = slide.shapes.add_shape(3, cx, cy, CIRCLE_R, CIRCLE_R)
        c_shape.fill.solid()
        c_shape.fill.fore_color.rgb = accent
        c_shape.line.fill.background()
        tf = c_shape.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = str(i + 1)
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = FONT_TITLE

        # 단계 제목
        add_textbox(slide, sx + Inches(0.08), BODY_Y + CIRCLE_R + Inches(0.28),
                    step_w - Inches(0.16), Inches(0.45),
                    step.get("label", f"단계 {i+1}"),
                    font_size=Pt(11), bold=True, color=DARK, font_face=FONT_TITLE,
                    align=PP_ALIGN.CENTER)

        # 설명
        if step.get("description"):
            add_textbox(slide, sx + Inches(0.1), BODY_Y + CIRCLE_R + Inches(0.8),
                        step_w - Inches(0.2), body_h - CIRCLE_R - Inches(0.9),
                        step["description"], font_size=Pt(9), color=GRAY,
                        font_face=FONT_BODY, align=PP_ALIGN.CENTER, wrap=True)

        # 단계 간 화살표
        if i < n - 1:
            arr_x = sx + step_w + gap * 0.1
            arr_y = BODY_Y + body_h / 2 - Inches(0.15)
            add_textbox(slide, arr_x, arr_y, gap * 0.8, Inches(0.3),
                        "▶", font_size=Pt(14), color=GRAY, font_face=FONT_BODY,
                        align=PP_ALIGN.CENTER)


def _slide_closing(prs, config: dict, slide_data: dict, slide_num: int) -> None:
    """마무리 슬라이드"""
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=PRIMARY)

    # 장식 원
    for x, y, r in [(10.5, -1.0, 3.5), (-1.5, 6.5, 3.0)]:
        c = slide.shapes.add_shape(3,
                                    Inches(x - r), Inches(y - r),
                                    Inches(r * 2), Inches(r * 2))
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE
        c.line.fill.background()

    add_textbox(slide, Inches(1.0), Inches(2.0), Inches(11.3), Inches(1.5),
                slide_data.get("title", "감사합니다"),
                font_size=Pt(44), bold=True, color=WHITE, font_face=FONT_TITLE,
                align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(3.5), Inches(3.6), Inches(6.3), Inches(0.06), fill=ACCENT)

    if slide_data.get("subtitle"):
        add_textbox(slide, Inches(1.0), Inches(3.8), Inches(11.3), Inches(0.9),
                    slide_data["subtitle"], font_size=Pt(18), color=LGRAY2,
                    font_face=FONT_BODY, align=PP_ALIGN.CENTER)

    if slide_data.get("contact"):
        add_textbox(slide, Inches(1.0), Inches(5.0), Inches(11.3), Inches(0.6),
                    slide_data["contact"], font_size=Pt(13), color=MGRAY,
                    font_face=FONT_BODY, align=PP_ALIGN.CENTER)

    if os.path.exists(LOGO_PATH):
        try:
            slide.shapes.add_picture(LOGO_PATH, SLIDE_W - Inches(2.6), SLIDE_H - Inches(1.1),
                                     height=Inches(0.75))
        except Exception:
            pass

    add_textbox(slide, SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.5), Inches(1.0), Inches(0.4),
                str(slide_num), font_size=Pt(10), color=MGRAY, font_face=FONT_BODY,
                align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════════
# 전체 발표 덱 생성
# ══════════════════════════════════════════════════════════════════════════════════

def build_presentation(config: dict, output_path: str) -> None:
    """
    config 키:
      title    : str  발표 제목 (필수)
      subtitle : str  부제목
      dept     : str  부서/작성자
      date     : str  날짜
      sections : list[dict]  섹션 목록 [{title, subtitle}]
      kpis     : list[dict]  KPI [{label, value, change, positive, caption}]
      table    : dict  {title, headers, rows}
      comparison: dict  {title, left:{heading,bullets}, right:{heading,bullets}, vs_label}
      process  : dict  {title, steps:[{label, description}]}
      closing  : dict  {title, subtitle, contact}
    """
    prs = new_presentation()

    sections = config.get("sections", [
        {"title": "Ⅰ. 현황 및 배경", "subtitle": "현황 분석 및 주요 지표"},
        {"title": "Ⅱ. 실행 방안",    "subtitle": "추진 과제 및 실행 계획"},
    ])

    # Slide 1: 표지
    _slide_cover(prs, config)
    print("  [1/10] 표지 완료")

    # Slide 2: 목차
    _slide_agenda(prs, config)
    print("  [2/10] 목차 완료")

    # Slide 3: 섹션 인트로 Ⅰ
    s1 = sections[0] if len(sections) > 0 else {"title": "Ⅰ. 현황 및 배경", "subtitle": ""}
    _slide_section(prs, "Ⅰ", s1["title"], s1.get("subtitle", ""), config, 3)
    print("  [3/10] 섹션 인트로 Ⅰ 완료")

    # Slide 4: 현황/배경 2단 비교
    two_col = config.get("two_column", {
        "title": "현황 및 배경",
        "left":  {"heading": "현황",   "bullets": ["• 현황 내용 1", "• 현황 내용 2", "• 현황 내용 3"]},
        "right": {"heading": "주요 이슈", "bullets": ["• 이슈 사항 1", "• 이슈 사항 2", "• 이슈 사항 3"]},
    })
    _slide_two_column(prs, config, two_col, 4)
    print("  [4/10] 2단 비교 완료")

    # Slide 5: KPI 수치
    kpi_slide = {
        "title": "핵심 성과 지표 (KPI)",
        "kpis": config.get("kpis", [
            {"label": "KPI 지표 1", "value": "—", "change": "", "positive": True, "caption": "설명 텍스트"},
            {"label": "KPI 지표 2", "value": "—", "change": "", "positive": True, "caption": "설명 텍스트"},
            {"label": "KPI 지표 3", "value": "—", "change": "", "positive": True, "caption": "설명 텍스트"},
            {"label": "KPI 지표 4", "value": "—", "change": "", "positive": True, "caption": "설명 텍스트"},
        ]),
    }
    _slide_kpi(prs, config, kpi_slide, 5)
    print("  [5/10] KPI 완료")

    # Slide 6: 섹션 인트로 Ⅱ
    s2 = sections[1] if len(sections) > 1 else {"title": "Ⅱ. 실행 방안", "subtitle": ""}
    _slide_section(prs, "Ⅱ", s2["title"], s2.get("subtitle", ""), config, 6, accent=NAVY)
    print("  [6/10] 섹션 인트로 Ⅱ 완료")

    # Slide 7: 표
    tbl = config.get("table", {
        "title": "주요 현황 비교표",
        "headers": ["구분",       "현황",       "목표",       "비고"],
        "rows": [
            ["[항목 1]", "[현황 값]", "[목표 값]", ""],
            ["[항목 2]", "[현황 값]", "[목표 값]", ""],
            ["[항목 3]", "[현황 값]", "[목표 값]", ""],
            ["[항목 4]", "[현황 값]", "[목표 값]", ""],
        ],
    })
    _slide_table(prs, config, tbl, 7)
    print("  [7/10] 표 완료")

    # Slide 8: 비교 분석
    comp = config.get("comparison", {
        "title": "현재 vs 목표",
        "left":  {"heading": "현재 상태", "bullets": ["• 현재 내용 1", "• 현재 내용 2", "• 현재 내용 3"]},
        "right": {"heading": "목표 상태", "bullets": ["• 목표 내용 1", "• 목표 내용 2", "• 목표 내용 3"]},
        "vs_label": "VS",
    })
    _slide_comparison(prs, config, comp, 8)
    print("  [8/10] 비교 분석 완료")

    # Slide 9: 프로세스
    proc = config.get("process", {
        "title": "추진 일정 및 실행 계획",
        "steps": [
            {"label": "1단계\n준비",      "description": "사전 준비 및 계획 수립"},
            {"label": "2단계\n실행",      "description": "실행 및 모니터링"},
            {"label": "3단계\n검토",      "description": "중간 점검 및 조정"},
            {"label": "4단계\n완료",      "description": "결과 분석 및 보고"},
        ],
    })
    _slide_process(prs, config, proc, 9)
    print("  [9/10] 프로세스 완료")

    # Slide 10: 마무리
    closing = config.get("closing", {
        "title":    "감사합니다",
        "subtitle": config.get("title", ""),
        "contact":  config.get("dept", ""),
    })
    _slide_closing(prs, config, closing, 10)
    print("  [10/10] 마무리 완료")

    prs.save(output_path)
    print(f"\n✓ 저장 완료: {output_path}")


# ══════════════════════════════════════════════════════════════════════════════════
# 기본 Placeholder 설정
# ══════════════════════════════════════════════════════════════════════════════════

DEFAULT_CONFIG = {
    "title":    "[프레젠테이션 제목을 입력하세요]",
    "subtitle": "[부제목]",
    "dept":     "[부서명]",
    "date":     "",
    "sections": [
        {"title": "Ⅰ. 현황 및 배경", "subtitle": "현황 분석 및 주요 지표"},
        {"title": "Ⅱ. 실행 방안",    "subtitle": "추진 과제 및 실행 계획"},
    ],
}


# ══════════════════════════════════════════════════════════════════════════════════
# CLI 진입점
# ══════════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    from datetime import date as _date

    parser = argparse.ArgumentParser(description="하나증권 발표용 PPT 템플릿 생성기")
    parser.add_argument("--title",  default="", help="발표 제목")
    parser.add_argument("--dept",   default="", help="부서명")
    parser.add_argument("--date",   default="", help="날짜 (예: 2026.05.01)")
    parser.add_argument("--out",    default="", help="출력 파일명 (.pptx)")
    args = parser.parse_args()

    config = dict(DEFAULT_CONFIG)
    if args.title: config["title"] = args.title
    if args.dept:  config["dept"]  = args.dept
    if args.date:  config["date"]  = args.date
    if not config.get("date"):
        config["date"] = _date.today().strftime("%Y.%m.%d")

    safe_title = (args.title or "발표자료").replace(" ", "_").replace("/", "_")
    output = args.out or f"{safe_title}.pptx"

    print(f"생성 중: {output}")
    build_presentation(config, output)
