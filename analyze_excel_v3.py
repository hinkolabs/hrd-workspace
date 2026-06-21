import sys
sys.stdout.reconfigure(encoding='utf-8')
import openpyxl
from pptx import Presentation

# ── 엑셀 파일 분석
XLS = r'C:\Users\jay\Downloads\hana_ai_excel_practice_instructor_v3.xlsx'
wb = openpyxl.load_workbook(XLS, data_only=True)
print('=== 엑셀 시트 목록 ===')
for name in wb.sheetnames:
    print(f'  {name}')

print()
for shname in wb.sheetnames:
    ws = wb[shname]
    print(f'\n=== 시트: {shname} (최대행={ws.max_row}, 최대열={ws.max_column}) ===')
    for row in ws.iter_rows(min_row=1, max_row=min(ws.max_row, 50), values_only=True):
        vals = [str(c) if c is not None else '' for c in row]
        line = ' | '.join(vals).strip(' |')
        if line.replace('|','').strip():
            print(f'  {line}')

# ── PPT 엑셀 실습 슬라이드 현황 (19, 20, 21)
print('\n\n=== PPT 슬라이드 19~21 텍스트 ===')
SRC = r'C:\Users\jay\Downloads\AI교육_완성판_v11.pptx'
prs = Presentation(SRC)
for si in [19, 20, 21]:
    slide = prs.slides[si - 1]
    print(f'\n--- Slide {si} ---')
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                t = p.text.strip()
                if t:
                    print(f'  "{t}"')
