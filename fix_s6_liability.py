"""
Slide 6:
- item 02: '고의가 아니어도 책임' 내용 정제
- item 03: 삼성전자 ChatGPT 기밀유출 사례 (03번 삭제된 자리에 신규 추가)
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SRC = r'C:\Users\jay\Downloads\AI교육_완성판_v10.pptx'
DST = r'C:\Users\jay\Downloads\AI교육_완성판_v11.pptx'
prs = Presentation(SRC)

GREEN      = RGBColor(0x00, 0x91, 0x78)
DARK_GREEN = RGBColor(0x00, 0x4E, 0x42)
LIGHT_GRAY = RGBColor(0xF7, 0xF9, 0xF8)
ORANGE     = RGBColor(0xE8, 0x6B, 0x00)   # 경고/사례 강조색
DARK_TEXT  = RGBColor(0x22, 0x22, 0x22)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
FONT       = '맑은 고딕'


def rgb_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def set_run(shape, text, size=None, bold=None, color=None, italic=None):
    tf = shape.text_frame
    p  = tf.paragraphs[0]
    # 기존 runs 제거 후 새로 작성
    from pptx.oxml.ns import qn
    for r in p.runs:
        r.text = ''
    if p.runs:
        run = p.runs[0]
    else:
        r_el = p._p.makeelement(qn('a:r'))
        p._p.append(r_el)
        run = p.runs[0]
    run.text = text
    run.font.name = FONT
    if size  is not None: run.font.size  = Pt(size)
    if bold  is not None: run.font.bold  = bold
    if color is not None: run.font.color.rgb = color
    if italic is not None: run.font.italic = italic


def add_rect(slide, l, t, w, h, rgb):
    s = slide.shapes.add_shape(1, l, t, w, h)
    rgb_fill(s, rgb)
    return s


def add_tb(slide, l, t, w, h, text, size=12, bold=False,
           color=DARK_TEXT, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = FONT
    return tb


slide  = prs.slides[5]
shapes = list(slide.shapes)

# ── 1. item 02 텍스트 정제 ─────────────────────────────────
# shape[18] = 제목, shape[19] = 본문
set_run(shapes[18],
        'AI에 데이터 입력할 때는 항상 조심',
        size=13, bold=True, color=DARK_GREEN)

set_run(shapes[19],
        '고의로 유출하려 하지 않았어도 책임을 집니다. '
        '개인정보·MNPI·대외비 자료를 외부 AI에 입력하는 순간 해당 데이터는 '
        'AI 서버에 저장·학습될 수 있으며, 과실에 의한 경우에도 '
        '내부통제 위반 및 징계·법적 책임의 대상이 됩니다.',
        size=10, color=DARK_TEXT)

# item 02 body rect 높이 늘리기 (내용이 길어짐)
shapes[17].height = Inches(1.30)   # body bg
shapes[15].height = Inches(1.30)   # badge bg

print('  item 02 텍스트 정제 완료')

# ── 2. item 03 삼성전자 사례 추가 ─────────────────────────
# 삭제된 03 자리: T=3.95 (02 끝 T=2.55+1.30=3.85 + 여백 0.10)
t03  = Inches(3.95)
H03  = Inches(1.60)   # 사례 내용이 길어 높이 여유

# 번호 배지 (주황색 — 경고/사례 강조)
badge_bg = add_rect(slide, Inches(0.47), t03, Inches(0.98), H03, ORANGE)
add_tb(slide, Inches(0.47), t03, Inches(0.98), H03,
       '03', size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 내용 배경
add_rect(slide, Inches(1.50), t03, Inches(8.03), H03, LIGHT_GRAY)

# 제목
add_tb(slide, Inches(1.65), t03 + Inches(0.10), Inches(7.80), Inches(0.32),
       '실제 사례: 삼성전자 ChatGPT 기밀 유출 (2023)',
       size=13, bold=True, color=ORANGE)

# 본문
samsung_body = (
    '삼성전자 반도체 직원들이 업무 편의를 위해 사내 소스코드·장비 오류 내용·회의 녹취를 '
    'ChatGPT에 입력 → 의도치 않은 기밀 3건 유출 발생. '
    '이후 삼성전자는 전사 외부 생성형 AI 사용을 전면 금지 조치. '
    '→  "실수"도 보안 사고이며, 책임은 피할 수 없습니다.'
)
add_tb(slide, Inches(1.65), t03 + Inches(0.46), Inches(7.80), Inches(1.08),
       samsung_body, size=10, color=DARK_TEXT, wrap=True)

print('  item 03 삼성전자 사례 추가 완료')

prs.save(DST)
print(f'\n✓ 저장: {DST}')
