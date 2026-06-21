"""
슬라이드 12~17 재디자인
- 기존 텍스트만 나열된 상태 → 기존 슬라이드 스타일과 통일
"""

import sys, shutil
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SRC = r'C:\Users\jay\Downloads\AI교육_완성판_v2.pptx'
DST = r'C:\Users\jay\Downloads\AI교육_완성판_v2.pptx'

prs = Presentation(SRC)

# ── 공통 색상 ──────────────────────────────────────────
GREEN       = RGBColor(0x00, 0x91, 0x78)
DARK_GREEN  = RGBColor(0x00, 0x4E, 0x42)
LIGHT_GREEN = RGBColor(0xD7, 0xED, 0xE6)
LIGHT_GRAY  = RGBColor(0xF7, 0xF9, 0xF8)
DARK_TEXT   = RGBColor(0x22, 0x22, 0x22)
RED         = RGBColor(0xCC, 0x00, 0x00)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY        = RGBColor(0xCC, 0xCC, 0xCC)
SLIDE_W     = prs.slide_width
SLIDE_H     = prs.slide_height
FONT        = '맑은 고딕'


def rgb_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def add_rect(slide, l, t, w, h, rgb):
    s = slide.shapes.add_shape(1, l, t, w, h)
    rgb_fill(s, rgb)
    return s


def add_tb(slide, l, t, w, h, text, size=12, bold=False,
           color=DARK_TEXT, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    return tb


def add_header(slide, section_label, title, page_num):
    """공통 헤더 (왼쪽 바 + 섹션라벨 + 제목 + 구분선 + 푸터)"""
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    # 왼쪽 녹색 바
    add_rect(slide, 0, 0, Inches(0.14), SLIDE_H, GREEN)
    # 섹션 라벨
    add_tb(slide, Inches(0.47), Inches(0.24), Inches(7.87), Inches(0.28),
           section_label, size=10, color=GREEN)
    # 타이틀
    add_tb(slide, Inches(0.47), Inches(0.48), Inches(9.06), Inches(0.54),
           title, size=24, bold=True, color=DARK_TEXT)
    # 구분선
    add_rect(slide, Inches(0.47), Inches(1.07), Inches(9.06), Inches(0.02), GREEN)
    # 푸터
    add_rect(slide, 0, Inches(7.05), SLIDE_W, Inches(0.45), GREEN)
    add_tb(slide, Inches(1.57), Inches(7.11), Inches(6.30), Inches(0.33),
           'Copyright ⓒ 2025 Hana Financial Group. All rights reserved.',
           size=6, color=WHITE)
    add_tb(slide, Inches(8.98), Inches(7.09), Inches(0.91), Inches(0.37),
           f'| {page_num}', size=8, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)


def remove_shapes(slide, keep_types=None):
    """슬라이드에서 keep_types 제외한 모든 shape 삭제 (기본=전부 삭제)"""
    to_del = []
    for shape in slide.shapes:
        if keep_types and shape.shape_type in keep_types:
            continue
        to_del.append(shape._element)
    for el in to_del:
        el.getparent().remove(el)


def add_numbered_block(slide, num, heading, desc, left, top, w=Inches(8.56), h=Inches(1.10)):
    """번호 강조 + 내용 row (슬라이드 10 스타일)"""
    num_w = Inches(0.80)
    add_rect(slide, left, top, num_w, h, GREEN)
    add_tb(slide, left, top, num_w, h, num,
           size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, left + num_w, top, w - num_w, h, LIGHT_GRAY)
    add_tb(slide, left + num_w + Inches(0.12), top + Inches(0.08),
           w - num_w - Inches(0.20), Inches(0.35),
           heading, size=13, bold=True, color=DARK_GREEN)
    if desc:
        add_tb(slide, left + num_w + Inches(0.12), top + Inches(0.46),
               w - num_w - Inches(0.20), h - Inches(0.52),
               desc, size=10, color=DARK_TEXT, wrap=True)


def add_step_row(slide, step_label, content, left, top, w=Inches(8.56), h=Inches(0.85)):
    """STEP 라벨 + 내용 row"""
    lbl_w = Inches(1.05)
    add_rect(slide, left, top, lbl_w, h, DARK_GREEN)
    add_tb(slide, left, top, lbl_w, h, step_label,
           size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, left + lbl_w, top, w - lbl_w, h, LIGHT_GRAY)
    add_tb(slide, left + lbl_w + Inches(0.12), top + Inches(0.18),
           w - lbl_w - Inches(0.20), h - Inches(0.24),
           content, size=11, color=DARK_TEXT, wrap=True)


# ══════════════════════════════════════════════════════════
#  Slide 12 — DART 실습 섹션 소개 (프롤로그)
# ══════════════════════════════════════════════════════════
def reformat_slide12():
    slide = prs.slides[11]
    remove_shapes(slide)

    add_header(slide, '2교시 · AI 활용 실습',
               'DART 사업보고서 AI 분석 실습', 12)

    # 도전 질문 하이라이트 박스
    add_rect(slide, Inches(0.47), Inches(1.20), Inches(9.06), Inches(0.85), DARK_GREEN)
    add_rect(slide, Inches(0.47), Inches(1.20), Inches(0.06), Inches(0.85), GREEN)
    add_tb(slide, Inches(0.63), Inches(1.35), Inches(8.70), Inches(0.55),
           '💡  SK하이닉스 3개년 사업보고서를 읽고 주가 상승 근거를 찾으려면 몇 시간이 걸릴까요? AI와 함께라면 단 30분!',
           size=13, bold=True, color=WHITE, wrap=True)

    # 실습 구성 3단계 (번호 row)
    items = [
        ('Ⅰ', 'DART 기업 핵심 정보를 정리해보기',
         'DART 사업보고서를 AI에 업로드 → 매출·재무·위험요인 분석 지시 → 원문 검증'),
        ('Ⅱ', 'AI의 결과를 내가 정확하게 이해했는지 확인하기',
         'AI가 생성한 퀴즈로 본인 이해도를 셀프 체크 → 틀린 개념 즉시 보완'),
        ('Ⅲ', 'AI 활용은 업무의 끝이 아닌 시작',
         '리포트·산업 자료를 꾸준히 추가하고, 좋은 질문으로 의미 있는 인사이트 발굴'),
    ]

    row_h  = Inches(1.48)
    gap    = Inches(0.12)
    num_w  = Inches(0.70)
    blk_l  = Inches(0.47)
    blk_w  = Inches(9.06)

    for i, (num, heading, desc) in enumerate(items):
        ty = Inches(2.20) + i * (row_h + gap)
        add_rect(slide, blk_l, ty, num_w, row_h, GREEN)
        add_tb(slide, blk_l, ty, num_w, row_h,
               num, size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, blk_l + num_w, ty, blk_w - num_w, row_h, LIGHT_GRAY)
        add_tb(slide, blk_l + num_w + Inches(0.14), ty + Inches(0.14),
               blk_w - num_w - Inches(0.22), Inches(0.42),
               heading, size=14, bold=True, color=DARK_GREEN)
        add_rect(slide, blk_l + num_w + Inches(0.14), ty + Inches(0.60),
                 blk_w - num_w - Inches(0.28), Inches(0.01), GRAY)
        add_tb(slide, blk_l + num_w + Inches(0.14), ty + Inches(0.68),
               blk_w - num_w - Inches(0.22), Inches(0.72),
               desc, size=11, color=DARK_TEXT, wrap=True)

    print('  [Slide 12] DART 프롤로그 완료')


# ══════════════════════════════════════════════════════════
#  Slide 13 — DART 방문 + AI 활용법
# ══════════════════════════════════════════════════════════
def reformat_slide13():
    slide = prs.slides[12]
    remove_shapes(slide)

    add_header(slide, '2교시 · AI 활용 실습',
               'Ⅰ. DART 기업 핵심 정보를 정리해보기', 13)

    # ── 왼쪽 패널: DART 다운로드 3단계 ──
    panel_w = Inches(4.30)
    lx = Inches(0.47)

    add_rect(slide, lx, Inches(1.20), panel_w, Inches(0.42), GREEN)
    add_tb(slide, lx + Inches(0.10), Inches(1.23), panel_w - Inches(0.20), Inches(0.38),
           'STEP 1  DART 방문 & 다운로드', size=12, bold=True, color=WHITE)

    steps_left = [
        ('①', 'dart.fss.or.kr 접속'),
        ('②', 'SK하이닉스 검색 → 사업보고서 선택'),
        ('③', '2023 · 2024 · 2025년 보고서 각각 다운로드'),
    ]
    for i, (n, t) in enumerate(steps_left):
        ty = Inches(1.72) + i * Inches(0.68)
        add_rect(slide, lx, ty, Inches(0.38), Inches(0.55), DARK_GREEN)
        add_tb(slide, lx, ty, Inches(0.38), Inches(0.55),
               n, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, lx + Inches(0.40), ty, panel_w - Inches(0.40), Inches(0.55), LIGHT_GRAY)
        add_tb(slide, lx + Inches(0.54), ty + Inches(0.10),
               panel_w - Inches(0.60), Inches(0.38),
               t, size=11, color=DARK_TEXT)

    add_rect(slide, lx, Inches(3.80), panel_w, Inches(0.42), GREEN)
    add_tb(slide, lx + Inches(0.10), Inches(3.83), panel_w - Inches(0.20), Inches(0.38),
           'STEP 2  AI에 업로드 & 분석 지시', size=12, bold=True, color=WHITE)

    add_rect(slide, lx, Inches(4.32), panel_w, Inches(2.40), LIGHT_GRAY)
    upload_text = (
        '• AI 채팅창에 PDF 파일 3개 업로드\n'
        '• 아래 항목들을 분석해달라고 지시:\n\n'
        '  1. 주요 사업부문 및 제품\n'
        '  2. 시장점유율\n'
        '  3. 영업현황 (매출·영업이익·영업이익률)\n'
        '  4. 재무현황 (자산·부채·부채비율·현금흐름)\n'
        '  5. 반도체 업황 및 위험요인'
    )
    add_tb(slide, lx + Inches(0.14), Inches(4.42), panel_w - Inches(0.24), Inches(2.20),
           upload_text, size=10.5, color=DARK_TEXT, wrap=True)

    # ── 오른쪽 패널: 사업보고서에서 중요한 정보 ──
    rx = Inches(5.00)
    rw = Inches(4.56)

    add_rect(slide, rx, Inches(1.20), rw, Inches(0.42), DARK_GREEN)
    add_tb(slide, rx + Inches(0.10), Inches(1.23), rw - Inches(0.20), Inches(0.38),
           '사업보고서에서 꼭 확인할 5가지', size=12, bold=True, color=WHITE)

    key_info = [
        ('01', '주요 사업부문 및 제품', '어떤 사업 영역에서 매출을 올리는가'),
        ('02', '시장점유율', '글로벌 반도체 시장에서의 위치'),
        ('03', '영업현황', '매출액·영업이익·영업이익률 추이'),
        ('04', '재무현황', '자산·부채비율·현금흐름의 건전성'),
        ('05', '업황·위험요인', '산업 전망 및 중단기 리스크'),
    ]
    item_h = Inches(1.08)
    for i, (num, ttl, sub) in enumerate(key_info):
        ty = Inches(1.72) + i * item_h
        bg_color = LIGHT_GREEN if i % 2 == 0 else LIGHT_GRAY
        add_rect(slide, rx, ty, rw, item_h - Inches(0.05), bg_color)
        add_rect(slide, rx, ty, Inches(0.42), item_h - Inches(0.05), GREEN)
        add_tb(slide, rx, ty, Inches(0.42), item_h - Inches(0.05),
               num, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_tb(slide, rx + Inches(0.50), ty + Inches(0.08),
               rw - Inches(0.58), Inches(0.35),
               ttl, size=12, bold=True, color=DARK_GREEN)
        add_tb(slide, rx + Inches(0.50), ty + Inches(0.48),
               rw - Inches(0.58), Inches(0.45),
               sub, size=9.5, color=DARK_TEXT)

    print('  [Slide 13] DART 방문 & AI 활용법 완료')


# ══════════════════════════════════════════════════════════
#  Slide 14 — 프롬프트 예시 (AI 분석 지시문)
# ══════════════════════════════════════════════════════════
def reformat_slide14():
    slide = prs.slides[13]
    # 말풍선 자동도형 보존, 일반 텍스트박스만 삭제
    to_del = []
    for shape in slide.shapes:
        if shape.shape_type == 17:  # TEXT_BOX
            to_del.append(shape._element)
    for el in to_del:
        el.getparent().remove(el)

    # 기존 말풍선(AUTO_SHAPE) 위치/크기 조정 → 텍스트 내용 유지
    callout = None
    for shape in slide.shapes:
        if shape.shape_type == 1 and shape.has_text_frame:
            callout = shape
            break

    add_header(slide, '2교시 · AI 활용 실습',
               'Ⅰ. DART 기업 핵심 정보를 정리해보기', 14)

    # ── 왼쪽: 프롬프트 작성 포인트 ──
    lx = Inches(0.47)
    lw = Inches(2.55)

    add_rect(slide, lx, Inches(1.20), lw, Inches(0.42), GREEN)
    add_tb(slide, lx + Inches(0.10), Inches(1.23), lw - Inches(0.20), Inches(0.38),
           '✍️ 프롬프트 작성 포인트', size=11, bold=True, color=WHITE)

    points = [
        ('역할 부여', '증권사 직원 역할 지정'),
        ('데이터 명시', '3개년 사업보고서 기반'),
        ('항목 구체화', '실적·재무·설비·위험 4개 분류'),
        ('수치 원칙', '원문 기준, 억원 단위 통일'),
        ('불확실 표시', '"확인 필요" 로 명시 요청'),
        ('근거 제시', '항목별 출처 문구 요청'),
    ]
    for i, (ttl, sub) in enumerate(points):
        ty = Inches(1.72) + i * Inches(0.85)
        add_rect(slide, lx, ty, lw, Inches(0.78), LIGHT_GRAY)
        add_rect(slide, lx, ty, Inches(0.06), Inches(0.78), GREEN)
        add_tb(slide, lx + Inches(0.14), ty + Inches(0.06),
               lw - Inches(0.20), Inches(0.32),
               ttl, size=11, bold=True, color=DARK_GREEN)
        add_tb(slide, lx + Inches(0.14), ty + Inches(0.40),
               lw - Inches(0.20), Inches(0.32),
               sub, size=9.5, color=DARK_TEXT)

    # ── 말풍선 위치 재설정 (오른쪽에 크게 배치) ──
    if callout:
        callout.left   = Inches(3.20)
        callout.top    = Inches(1.20)
        callout.width  = Inches(6.42)
        callout.height = Inches(5.60)
        # 텍스트 스타일 적용
        try:
            rgb_fill(callout, LIGHT_GRAY)
            callout.line.color.rgb = GREEN
            callout.line.width     = Pt(1.5)
            for para in callout.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size  = Pt(10)
                    run.font.name  = FONT
                    run.font.color.rgb = DARK_TEXT
        except Exception as e:
            print(f'    말풍선 스타일 오류 (무시): {e}')

    # 말풍선 라벨
    add_rect(slide, Inches(3.20), Inches(1.20), Inches(6.42), Inches(0.38), DARK_GREEN)
    add_tb(slide, Inches(3.30), Inches(1.23), Inches(6.20), Inches(0.35),
           '📋 AI 프롬프트 예시 — 복사해서 바로 사용하세요',
           size=11, bold=True, color=WHITE)

    print('  [Slide 14] 프롬프트 예시 완료')


# ══════════════════════════════════════════════════════════
#  Slide 15 — 할루시네이션 방지 원칙 + 검증
# ══════════════════════════════════════════════════════════
def reformat_slide15():
    slide = prs.slides[14]
    # 이미지(13) + 자동도형(1) 보존, 텍스트박스(17)만 삭제
    to_del = []
    for shape in slide.shapes:
        if shape.shape_type == 17:
            to_del.append(shape._element)
    for el in to_del:
        el.getparent().remove(el)

    add_header(slide, '2교시 · AI 활용 실습',
               'Ⅰ. 할루시네이션 방지를 위한 핵심 원칙', 15)

    # 4가지 원칙 카드 (2×2 그리드)
    principles = [
        ('01', '수치는 원문 기준',
         'AI가 임의로 계산하지 않도록, 사업보고서 원문에 있는 수치만 그대로 추출'),
        ('02', '단위 통일',
         '매출·이익·자산 모두 "억원" 단위로 통일 명시 — 혼용 시 비교 불가'),
        ('03', '불확실 내용 표시',
         '확인되지 않는 내용은 "확인 필요"로 명시 요청 — 그럴듯한 오답 차단'),
        ('04', '근거 문구 요청',
         '각 항목마다 "사업보고서 몇 페이지, 어느 문장"인지 출처 함께 제시'),
    ]

    card_w = Inches(4.37)
    card_h = Inches(1.45)
    gap    = Inches(0.25)
    positions = [
        (Inches(0.47), Inches(1.20)),
        (Inches(0.47) + card_w + gap, Inches(1.20)),
        (Inches(0.47), Inches(1.20) + card_h + gap),
        (Inches(0.47) + card_w + gap, Inches(1.20) + card_h + gap),
    ]
    for (num, ttl, desc), (lx, ty) in zip(principles, positions):
        add_rect(slide, lx, ty, card_w, card_h, LIGHT_GRAY)
        add_rect(slide, lx, ty, Inches(0.65), card_h, GREEN)
        add_tb(slide, lx, ty, Inches(0.65), card_h,
               num, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_tb(slide, lx + Inches(0.73), ty + Inches(0.10),
               card_w - Inches(0.80), Inches(0.38),
               ttl, size=13, bold=True, color=DARK_GREEN)
        add_rect(slide, lx + Inches(0.73), ty + Inches(0.52),
                 card_w - Inches(0.80), Inches(0.01), GRAY)
        add_tb(slide, lx + Inches(0.73), ty + Inches(0.60),
               card_w - Inches(0.80), Inches(0.78),
               desc, size=10, color=DARK_TEXT, wrap=True)

    # 검증 안내 (이미지들 위쪽에 배치)
    add_rect(slide, Inches(0.47), Inches(4.30), Inches(9.06), Inches(0.38), DARK_GREEN)
    add_tb(slide, Inches(0.60), Inches(4.33), Inches(8.80), Inches(0.33),
           '✅ 분석 결과가 나왔다면? — DART 원문과 직접 맞춰보세요',
           size=12, bold=True, color=WHITE)

    print('  [Slide 15] 할루시네이션 방지 완료')


# ══════════════════════════════════════════════════════════
#  Slide 17 — AI 활용은 업무의 끝이 아닌 시작
# ══════════════════════════════════════════════════════════
def reformat_slide17():
    slide = prs.slides[16]
    # 이미지(13) 보존, 텍스트박스(17)만 삭제
    to_del = []
    for shape in slide.shapes:
        if shape.shape_type == 17:
            to_del.append(shape._element)
    for el in to_del:
        el.getparent().remove(el)

    add_header(slide, '2교시 · AI 활용 실습',
               'Ⅲ. AI 활용은 업무의 끝이 아닌 시작', 17)

    # 이미지는 왼쪽(L:0.27" ~ 4.26")에 있으므로 오른쪽에 배치
    rx = Inches(4.50)
    rw = Inches(5.08)

    messages = [
        ('📊', 'AI는 속도, 인사이트는 여러분의 몫',
         'AI는 자료 정리 속도를 높여줍니다. 의미 있는 해석과 판단은 항상 여러분이 해야 합니다.'),
        ('📚', '지식 베이스를 꾸준히 쌓으세요',
         '신뢰성 있는 증권사 리포트, 산업 분석 자료를 AI에 추가하면 답변 품질이 높아집니다.'),
        ('⚖️', '경쟁사 자료도 함께 비교하세요',
         '동종 업계 사업보고서를 함께 업로드하면 상대적 강·약점 분석이 가능합니다.'),
        ('🙋', '기초 질문을 두려워 말 것',
         '"이게 무슨 뜻이야?" 같은 기초 질문이 오히려 더 정확한 이해를 만들어줍니다.'),
    ]

    msg_h = Inches(1.35)
    gap   = Inches(0.10)
    for i, (icon, ttl, desc) in enumerate(messages):
        ty = Inches(1.20) + i * (msg_h + gap)
        add_rect(slide, rx, ty, rw, msg_h, LIGHT_GRAY)
        add_rect(slide, rx, ty, rw, Inches(0.04), GREEN)
        add_tb(slide, rx + Inches(0.10), ty + Inches(0.10),
               Inches(0.45), Inches(0.50),
               icon, size=22, color=GREEN)
        add_tb(slide, rx + Inches(0.62), ty + Inches(0.10),
               rw - Inches(0.72), Inches(0.38),
               ttl, size=13, bold=True, color=DARK_GREEN)
        add_rect(slide, rx + Inches(0.62), ty + Inches(0.52),
                 rw - Inches(0.72), Inches(0.01), GRAY)
        add_tb(slide, rx + Inches(0.62), ty + Inches(0.60),
               rw - Inches(0.72), Inches(0.68),
               desc, size=10, color=DARK_TEXT, wrap=True)

    print('  [Slide 17] AI 활용의 가치 완료')


# ══════════════════════════════════════════════════════════
#  Slide 16 — AI 결과 검증 (헤더/푸터 보강)
# ══════════════════════════════════════════════════════════
def reformat_slide16():
    slide = prs.slides[15]
    # 기존 원본 텍스트박스 2개 (title, body) 와 "AI 질문 샘플" 라벨 삭제
    # 새로 추가한 Q1~Q5, 배경박스들은 유지
    to_del = []
    for shape in slide.shapes:
        if shape.shape_type == 17:
            name = shape.name
            # 이전에 추가한 새 텍스트박스는 이름이 'TextBox 14', 'TextBox 16' 등
            # 원본은 'TextBox 3', 'TextBox 9', 'TextBox 11'
            if name in ('TextBox 3', 'TextBox 9', 'TextBox 11'):
                to_del.append(shape._element)
    for el in to_del:
        el.getparent().remove(el)

    # 공통 헤더/푸터 추가
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_rect(slide, 0, 0, Inches(0.14), SLIDE_H, GREEN)
    add_tb(slide, Inches(0.47), Inches(0.24), Inches(7.87), Inches(0.28),
           '2교시 · AI 활용 실습', size=10, color=GREEN)
    add_tb(slide, Inches(0.47), Inches(0.48), Inches(9.06), Inches(0.54),
           'Ⅱ. AI의 결과를 내가 정확하게 이해했는지 확인하기',
           size=20, bold=True, color=DARK_TEXT)
    add_rect(slide, Inches(0.47), Inches(1.07), Inches(9.06), Inches(0.02), GREEN)
    add_rect(slide, 0, Inches(7.05), SLIDE_W, Inches(0.45), GREEN)
    add_tb(slide, Inches(1.57), Inches(7.11), Inches(6.30), Inches(0.33),
           'Copyright ⓒ 2025 Hana Financial Group. All rights reserved.',
           size=6, color=WHITE)
    add_tb(slide, Inches(8.98), Inches(7.09), Inches(0.91), Inches(0.37),
           '| 16', size=8, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    # 왼쪽 지시 패널
    lx, lw = Inches(0.47), Inches(4.00)
    add_rect(slide, lx, Inches(1.20), lw, Inches(0.42), DARK_GREEN)
    add_tb(slide, lx + Inches(0.10), Inches(1.23), lw - Inches(0.20), Inches(0.38),
           '1. AI에게 질문 생성 요청', size=12, bold=True, color=WHITE)
    add_rect(slide, lx, Inches(1.72), lw, Inches(3.05), LIGHT_GRAY)
    prompt = (
        '"내가 너(AI)의 분석 결과를\n'
        '정확하게 이해했는지 검증해줘.\n\n'
        'SK하이닉스 분석 내용을 바탕으로\n'
        '나에게 질문해주고,\n'
        '내 답변도 채점해줘."'
    )
    add_tb(slide, lx + Inches(0.15), Inches(1.82), lw - Inches(0.25), Inches(2.85),
           prompt, size=12, color=DARK_GREEN, wrap=True)

    add_rect(slide, lx, Inches(4.87), lw, Inches(0.38), GREEN)
    add_tb(slide, lx + Inches(0.10), Inches(4.90), lw - Inches(0.20), Inches(0.33),
           '2. 직접 답하고 피드백 받기', size=12, bold=True, color=WHITE)
    add_rect(slide, lx, Inches(5.35), lw, Inches(1.45), LIGHT_GRAY)
    add_tb(slide, lx + Inches(0.15), Inches(5.45), lw - Inches(0.25), Inches(1.25),
           '• AI가 제시한 질문에 직접 답변 입력\n'
           '• AI가 즉시 채점 + 오답 해설 제공\n'
           '• 틀린 개념은 바로 다시 질문해서 보완',
           size=11, color=DARK_TEXT, wrap=True)

    print('  [Slide 16] AI 결과 검증 헤더 보강 완료')


# ══════════════════════════════════════════════════════════
if __name__ == '__main__':
    print('=== 슬라이드 12~17 재디자인 시작 ===\n')
    reformat_slide12()
    reformat_slide13()
    reformat_slide14()
    reformat_slide15()
    reformat_slide16()
    reformat_slide17()

    prs.save(DST)
    print(f'\n✓ 저장 완료: {DST}')
