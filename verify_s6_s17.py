import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches

SRC = r'C:\Users\jay\Downloads\AI교육_완성판_v6.pptx'
prs = Presentation(SRC)

def show(idx):
    slide = prs.slides[idx]
    print(f'\n=== Slide {idx+1} ({len(slide.shapes)} shapes) ===')
    for i, shape in enumerate(slide.shapes):
        L = shape.left/914400
        T = shape.top/914400
        W = shape.width/914400
        H = shape.height/914400
        txt = ''
        if shape.has_text_frame:
            txt = ' | ' + ' / '.join(p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip())[:80]
        print(f'  [{i:2d}] T={T:.2f}" L={L:.2f}" W={W:.2f}" H={H:.2f}" type={shape.shape_type}{txt}')

show(5)
show(16)
