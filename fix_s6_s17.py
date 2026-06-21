"""
Slide 6 : item 02 내용 교체
Slide 17: 사진 뒤로 빠짐 + 레이아웃 겹침 수정
"""
import sys, copy
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

SRC = r'C:\Users\jay\Downloads\AI교육_완성판_v5.pptx'
DST = r'C:\Users\jay\Downloads\AI교육_완성판_v6.pptx'
prs = Presentation(SRC)

GREEN       = RGBColor(0x00, 0x91, 0x78)
DARK_GREEN  = RGBColor(0x00, 0x4E, 0x42)
LIGHT_GREEN = RGBColor(0xD7, 0xED, 0xE6)
LIGHT_GRAY  = RGBColor(0xF7, 0xF9, 0xF8)
DARK_TEXT   = RGBColor(0x22, 0x22, 0x22)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY        = RGBColor(0xCC, 0xCC, 0xCC)
MGRAY       = RGBColor(0x88, 0x88, 0x88)
FONT        = '맑은 고딕'
SLIDE_W     = prs.slide_width
SLIDE_H     = prs.slide_height


def rgb_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def add_rect(slide, l, t, w, h, rgb):
    s = slide.shapes.add_shape(1, l, t, w, h)
    rgb_fill(s, rgb)
    return s


def add_tb(slide, l, t, w, h, text, size=12, bold=False,
           color=DARK_TEXT, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = FONT
    return tb


def remove_shapes_by_idx(slide, indices):
    sp_tree = slide.shapes._spTree
    shapes = list(slide.shapes)
    for i in sorted(indices, reverse=True):
        el = shapes[i]._element
        el.getparent().remove(el)


# ══════════════════════════════════════════════════════════
#  Slide 6 — item 02 교체
#   현재 item 02 shapes: [12] rect badge, [13] tb "02",
#                        [14] rect body, [15] tb title, [16] tb content
#   이것들을 제거하고 새 내용으로 재생성
# ══════════════════════════════════════════════════════════
def fix_slide6():
    slide = prs.slides[5]

    # 기존 item 02 관련 5개 shape 제거 (indices 12~16)
    remove_shapes_by_idx(slide, [12, 13, 14, 15, 16])

    # 새 item 02 추가  (item 01은 T1.30, item 03은 T3.80 → item 02 위치 T2.55 유지)
    t = Inches(2.55)
    H = Inches(1.10)

    # 번호 배지
    badge = add_rect(slide, Inches(0.47), t, Inches(0.98), H, GREEN)
    tb_num = add_tb(slide, Inches(0.47), t, Inches(0.98), H,
                    '02', size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # 내용 배경
    body_bg = add_rect(slide, Inches(1.50), t, Inches(8.03), H, LIGHT_GRAY)

    # 제목
    add_tb(slide, Inches(1.65), t + Inches(0.10), Inches(7.80), Inches(0.32),
           '사내 허용 AI 서비스 안내',
           size=13, bold=True, color=DARK_GREEN)

    # 내용
    content = ('공식 업무용: Alli Works(사내 AI 플랫폼), MS Copilot(도입 예정)  '
               '│  외부 AI(ChatGPT·Gemini 등) 사용 시: 비업무·비공개 계정 이용 권장  '
               '│  AI 생성 콘텐츠는 저작권·사내 가이드라인 준수하여 활용')
    add_tb(slide, Inches(1.65), t + Inches(0.42), Inches(7.80), Inches(0.58),
           content, size=10, color=DARK_TEXT, wrap=True)

    print('  [Slide 6] item 02 교체 완료')


# ══════════════════════════════════════════════════════════
#  Slide 17 — 레이아웃 재구성
#   문제: 사진(shape[0]) 이 L3.55 T2.91 에 있어서 양쪽 패널과 겹침
#   해결: 사진을 좌측 패널 하단에 맞게 재배치 + 좌측 텍스트 높이 조정
# ══════════════════════════════════════════════════════════
def fix_slide17():
    slide = prs.slides[16]

    # ── 기존 shape 인덱스 (원본 v5 기준):
    # [0]  사진  L3.55 T2.91 W3.96 H3.17
    # [26] 좌측 녹색 세로 바  L0.00 T0.00 W0.14 H7.50  ← z-order 문제 원인
    # [33] Banner bg "1. AI에게 질문 생성 요청"  T1.20
    # [34] Banner tb                             T1.23
    # [35] Prompt bg                             T1.72 H3.05
    # [36] Prompt tb                             T1.82
    # [37] Banner bg "2. 직접 답하고..."           T4.87
    # [38] Banner tb                             T4.90
    # [39] Content bg                            T5.35 H1.45
    # [40] Content tb                            T5.45

    shapes = list(slide.shapes)
    sp_tree = slide.shapes._spTree

    # step1: 좌측 텍스트 요소들 위치/크기 조정 ─────────────────────
    # Prompt bg (shape[35]): 높이 H3.05 → H1.85
    shapes[35].height = Inches(1.85)   # T1.72 ~ T3.57

    # Banner2 bg (shape[37]): T4.87 → T3.70
    shapes[37].top = Inches(3.70)
    # Banner2 tb (shape[38]): T4.90 → T3.73  ← 누락 수정
    shapes[38].top = Inches(3.73)

    # Content bg (shape[39]): T5.35 → T4.15, H1.45 → H0.90
    shapes[39].top    = Inches(4.15)
    shapes[39].height = Inches(0.90)
    # Content tb (shape[40]): T5.45 → T4.22, H→H0.80
    shapes[40].top    = Inches(4.22)
    shapes[40].height = Inches(0.80)

    # step2: 사진 재배치 ─────────────────────────────────────
    # 원본 비율 W3.96:H3.17 = 1.249
    # 가용 공간 T5.15 ~ T7.00 = H1.85" → W = 1.85*1.249 = 2.31"
    # 중앙: L = 0.47 + (4.00-2.31)/2 = 1.315"
    s_pic = shapes[0]
    s_pic.left   = Inches(1.00)
    s_pic.top    = Inches(5.12)
    s_pic.width  = Inches(2.31)
    s_pic.height = Inches(1.85)

    # step3: z-order 정리 ─────────────────────────────────────
    # 순서: 녹색 바 → 헤더/푸터 → 좌측 패널 → 사진 → 우측 Q1-Q5
    # 녹색 바(shape[26])를 맨 앞으로 이동
    green_bar = shapes[26]._element
    green_bar.getparent().remove(green_bar)
    sp_tree.insert(2, green_bar)   # spTree 앞쪽에 삽입

    # 사진을 중간으로 (헤더/푸터·좌측 이후, 우측 이전)
    # 현재 shapes 재로드
    shapes2 = list(slide.shapes)
    pic_shape = None
    for s in shapes2:
        if s.shape_type == 13:   # PICTURE
            pic_shape = s
            break
    pic_el = pic_shape._element
    pic_el.getparent().remove(pic_el)
    sp_tree.append(pic_el)

    # 우측 Q1-Q5 패널 (Rectangle12, TextBox14 등) 맨 위로
    # 우측 패널은 L≥4.80 에 위치한 것들 (헤더/푸터 제외)
    def is_right_panel(s):
        return (s.left >= Inches(4.70)
                and s.top >= Inches(1.20)
                and s.top <= Inches(7.00))
    shapes3 = list(slide.shapes)
    right_els = [s._element for s in shapes3 if is_right_panel(s)]
    for el in right_els:
        el.getparent().remove(el)
    for el in right_els:
        sp_tree.append(el)

    print('  [Slide 17] 사진 재배치 + z-order + 레이아웃 정리 완료')


# ══════════════════════════════════════════════════════════
fix_slide6()
fix_slide17()
prs.save(DST)
print(f'\n✓ 저장: {DST}')
