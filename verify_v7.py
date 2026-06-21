import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

SRC = r'C:\Users\jay\Downloads\AI교육_완성판_v7.pptx'
prs = Presentation(SRC)

print('=== Slide 6 item 02 확인 ===')
slide6 = prs.slides[5]
for i, shape in enumerate(slide6.shapes):
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            t = p.text.strip()
            if t and i >= 20:
                print(f'  [{i}] "{t}"')

print('\n=== Slide 18 항목 확인 ===')
slide18 = prs.slides[17]
for i, shape in enumerate(slide18.shapes):
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            t = p.text.strip()
            if t and i >= 16:
                print(f'  [{i}] "{t}"')

print('\n=== 교시 잔존 여부 검색 ===')
for si, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if '교시' in p.text:
                    print(f'  Slide {si+1}: "{p.text.strip()}"')
print('  (없으면 교체 완료)')

print('\n=== 섹션 라벨 확인 (교시 교체 결과) ===')
check_slides = [10, 16, 19, 20, 21, 23, 24]
for si in check_slides:
    slide = prs.slides[si - 1]
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                t = p.text.strip()
                if '이론' in t or '실습' in t or '마무리' in t:
                    print(f'  Slide {si}: "{t}"')
                    break
