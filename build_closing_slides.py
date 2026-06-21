"""
마무리 섹션 3슬라이드 재구성
- Slide 25 교체: 이미지 카드 만들기 실습
- Slide 26 신규: Q & A
- Slide 27 신규: 고생하셨습니다
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SRC = r'C:\Users\jay\Downloads\AI교육_완성판_v4.pptx'
DST = r'C:\Users\jay\Downloads\AI교육_완성판_v5.pptx'
prs = Presentation(SRC)

GREEN       = RGBColor(0x00, 0x91, 0x78)
DARK_GREEN  = RGBColor(0x00, 0x4E, 0x42)
LIGHT_GREEN = RGBColor(0xD7, 0xED, 0xE6)
LIGHT_GRAY  = RGBColor(0xF7, 0xF9, 0xF8)
DARK_TEXT   = RGBColor(0x22, 0x22, 0x22)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY        = RGBColor(0xCC, 0xCC, 0xCC)
MGRAY       = RGBColor(0x88, 0x88, 0x88)
SLIDE_W     = prs.slide_width
SLIDE_H     = prs.slide_height
FONT        = '맑은 고딕'


def rgb_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def add_rect(slide, l, t, w, h, rgb):
    s = slide.shapes.add_shape(1, l, t, w, h)
    rgb_fill(s, rgb)
    return s


def add_tb(slide, l, t, w, h, text, size=12, bold=False,
           color=DARK_TEXT, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = FONT
    return tb


def remove_all(slide):
    for el in [s._element for s in slide.shapes]:
        el.getparent().remove(el)


def add_blank_slide():
    """맨 뒤에 빈 슬라이드 추가 (blank layout)"""
    layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(layout)
    for ph in new_slide.placeholders:
        try:
            ph._element.getparent().remove(ph._element)
        except Exception:
            pass
    return new_slide


# ══════════════════════════════════════════════════════════
#  Slide 25 교체 — 이미지 카드 만들기 실습
# ══════════════════════════════════════════════════════════
def rebuild_slide25():
    slide = prs.slides[24]   # 0-based index 24 = slide 25
    remove_all(slide)

    # 흰 배경 + 공통 헤더
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_rect(slide, 0, 0, Inches(0.14), SLIDE_H, GREEN)
    add_tb(slide, Inches(0.47), Inches(0.24), Inches(7.87), Inches(0.28),
           '마무리 · AI 콘텐츠 만들기', size=10, color=GREEN)
    add_tb(slide, Inches(0.47), Inches(0.48), Inches(9.06), Inches(0.54),
           '오늘 배운 것, 이미지 카드로 소장하기',
           size=24, bold=True, color=DARK_TEXT)
    add_rect(slide, Inches(0.47), Inches(1.07), Inches(9.06), Inches(0.02), GREEN)

    # 상단 설명 배너
    add_rect(slide, Inches(0.47), Inches(1.20), Inches(9.06), Inches(0.46), DARK_GREEN)
    add_rect(slide, Inches(0.47), Inches(1.20), Inches(0.06), Inches(0.46), GREEN)
    add_tb(slide, Inches(0.62), Inches(1.26), Inches(8.80), Inches(0.36),
           '오늘 교육에서 가장 기억하고 싶은 내용 1가지를 골라 나만의 AI 이미지 카드로 만들어 소장해보세요.',
           size=11, color=WHITE, wrap=True)

    # STEP 카드 3개 (가로 배치)
    steps = [
        ('STEP 1', '기억하고 싶은 내용 선택',
         '오늘 배운 것 중 한 가지를 고르세요.\n\n'
         '예)\n'
         '• "AI에게 역할을 먼저 줘라"\n'
         '• "처리상태=정상이어도 오류가 있을 수 있다"\n'
         '• "항상 원문으로 검증하라"'),
        ('STEP 2', 'AI에게 이미지 카드 요청',
         'ChatGPT · Gemini · Canva AI 등을 열고\n아래 프롬프트를 붙여넣어 실행해보세요.\n\n'
         '오른쪽 프롬프트 박스를 그대로 복사해서\n사용하면 됩니다.'),
        ('STEP 3', '저장 & 팀원 공유',
         '결과 이미지를 저장하고\n옆 동료와 어떤 내용을 선택했는지\n가볍게 공유해보세요.\n\n'
         '소요시간: 약 15~20분'),
    ]

    card_w = Inches(2.85)
    gap    = Inches(0.22)
    start_x = Inches(0.47)
    card_h = Inches(3.30)
    card_top = Inches(1.80)

    for i, (step_lbl, heading, desc) in enumerate(steps):
        lx = start_x + i * (card_w + gap)
        add_rect(slide, lx, card_top, card_w, card_h, LIGHT_GRAY)
        add_rect(slide, lx, card_top, card_w, Inches(0.06), GREEN)
        # 스텝 배지
        add_rect(slide, lx + Inches(0.15), card_top + Inches(0.15),
                 Inches(1.00), Inches(0.34), GREEN)
        add_tb(slide, lx + Inches(0.15), card_top + Inches(0.15),
               Inches(1.00), Inches(0.34),
               step_lbl, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_tb(slide, lx + Inches(0.15), card_top + Inches(0.60),
               card_w - Inches(0.28), Inches(0.42),
               heading, size=13, bold=True, color=DARK_GREEN)
        add_rect(slide, lx + Inches(0.15), card_top + Inches(1.06),
                 card_w - Inches(0.28), Inches(0.01), GRAY)
        add_tb(slide, lx + Inches(0.15), card_top + Inches(1.14),
               card_w - Inches(0.28), card_h - Inches(1.22),
               desc, size=10, color=DARK_TEXT, wrap=True)

    # 프롬프트 예시 박스 (하단)
    prompt_top = Inches(5.25)
    add_rect(slide, Inches(0.47), prompt_top, Inches(9.06), Inches(0.36), DARK_GREEN)
    add_tb(slide, Inches(0.60), prompt_top + Inches(0.04),
           Inches(8.80), Inches(0.30),
           '📋  복붙용 AI 프롬프트', size=11, bold=True, color=WHITE)

    prompt_text = (
        '"오늘 AI 교육에서 내가 기억하고 싶은 내용: [여기에 선택한 내용을 입력]. '
        '이것을 감각적인 카드 이미지로 만들어줘. '
        '핵심 문구 1줄 + 간단한 시각 요소 포함. '
        '밝고 긍정적인 톤으로, 금융 직장인이 스마트폰에 저장해두고 싶은 디자인으로."'
    )
    add_rect(slide, Inches(0.47), prompt_top + Inches(0.36),
             Inches(9.06), Inches(0.82), LIGHT_GREEN)
    add_tb(slide, Inches(0.60), prompt_top + Inches(0.42),
           Inches(8.80), Inches(0.72),
           prompt_text, size=10.5, color=DARK_GREEN, italic=True, wrap=True)

    # 푸터
    add_rect(slide, 0, Inches(7.05), SLIDE_W, Inches(0.45), GREEN)
    add_tb(slide, Inches(1.57), Inches(7.11), Inches(6.30), Inches(0.33),
           'Copyright ⓒ 2025 Hana Financial Group. All rights reserved.',
           size=6, color=WHITE)
    add_tb(slide, Inches(8.98), Inches(7.09), Inches(0.91), Inches(0.37),
           '| 25', size=8, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    print('  [Slide 25] 이미지 카드 만들기 실습 완료')


# ══════════════════════════════════════════════════════════
#  Slide 26 신규 — Q & A
# ══════════════════════════════════════════════════════════
def add_slide26():
    slide = add_blank_slide()

    # 짙은 녹색 배경
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_GREEN

    # 장식 원형 (우상단, 좌하단)
    for (cx, cy, r) in [(9.5, -0.8, 3.2), (-1.8, 7.5, 2.8)]:
        c = slide.shapes.add_shape(3,
                                   Inches(cx - r), Inches(cy - r),
                                   Inches(r * 2), Inches(r * 2))
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE
        c.line.fill.background()

    # 중앙 구분선 (위쪽)
    add_rect(slide, Inches(2.00), Inches(2.65), Inches(6.30), Inches(0.06), GREEN)

    # 대형 Q & A 텍스트
    add_tb(slide, Inches(0.80), Inches(2.05), Inches(11.20), Inches(1.60),
           'Q & A',
           size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # 서브 텍스트
    add_tb(slide, Inches(0.80), Inches(3.50), Inches(11.20), Inches(0.70),
           '궁금했던 것, 헷갈렸던 것 모두 질문해주세요',
           size=20, color=LIGHT_GREEN, align=PP_ALIGN.CENTER)

    # 하단 안내
    add_tb(slide, Inches(0.80), Inches(4.50), Inches(11.20), Inches(0.50),
           '오늘 배운 내용 외에도 AI 활용 관련 질문은 언제든지 환영합니다',
           size=13, color=MGRAY, align=PP_ALIGN.CENTER, italic=True)

    # 질문 유도 아이콘 라인
    add_tb(slide, Inches(1.50), Inches(5.40), Inches(10.30), Inches(0.50),
           '💬   어떤 질문도 좋습니다   💬',
           size=14, color=RGBColor(0x00, 0x91, 0x78), align=PP_ALIGN.CENTER)

    # 푸터
    add_tb(slide, Inches(1.57), Inches(7.11), Inches(6.30), Inches(0.33),
           'Copyright ⓒ 2025 Hana Financial Group. All rights reserved.',
           size=6, color=MGRAY)
    add_tb(slide, Inches(8.98), Inches(7.09), Inches(0.91), Inches(0.37),
           '| 26', size=8, bold=True, color=MGRAY, align=PP_ALIGN.RIGHT)

    print('  [Slide 26] Q&A 슬라이드 완료')
    return slide


# ══════════════════════════════════════════════════════════
#  Slide 27 신규 — 고생하셨습니다
# ══════════════════════════════════════════════════════════
def add_slide27():
    slide = add_blank_slide()

    # 하나그린 풀블리드 배경
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = GREEN

    # 장식 원형 (우상단, 좌하단)
    for (cx, cy, r) in [(10.2, -1.0, 3.8), (-1.8, 7.8, 3.2)]:
        c = slide.shapes.add_shape(3,
                                   Inches(cx - r), Inches(cy - r),
                                   Inches(r * 2), Inches(r * 2))
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE
        c.line.fill.background()

    # 작은 상단 라벨
    add_tb(slide, Inches(0.80), Inches(1.30), Inches(11.20), Inches(0.45),
           '하나증권 신입사원 AI 활용 교육',
           size=13, color=LIGHT_GREEN, align=PP_ALIGN.CENTER)

    # 메인 타이틀
    add_tb(slide, Inches(0.80), Inches(2.10), Inches(11.20), Inches(1.60),
           '수고하셨습니다!',
           size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # 구분선
    add_rect(slide, Inches(2.80), Inches(3.85), Inches(7.70), Inches(0.06),
             DARK_GREEN)

    # 서브 메시지
    add_tb(slide, Inches(0.80), Inches(4.05), Inches(11.20), Inches(0.65),
           'AI를 올바르게, 안전하게, 스마트하게 활용하는 하나인이 되어주세요.',
           size=17, color=LIGHT_GREEN, align=PP_ALIGN.CENTER)

    # 핵심 3원칙 요약 (아이콘 + 텍스트)
    principles = [
        ('🛡️', '윤리적으로'),
        ('🔧', '기술적으로'),
        ('✅', '항상 검증'),
    ]
    p_w = Inches(2.60)
    for i, (icon, label) in enumerate(principles):
        lx = Inches(1.70) + i * (p_w + Inches(0.40))
        ty = Inches(5.05)
        add_tb(slide, lx, ty, p_w, Inches(0.55),
               icon, size=28, align=PP_ALIGN.CENTER,
               color=WHITE)
        add_tb(slide, lx, ty + Inches(0.55), p_w, Inches(0.40),
               label, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Copyright
    add_tb(slide, Inches(1.57), Inches(7.11), Inches(7.00), Inches(0.33),
           'Copyright ⓒ 2025 Hana Financial Group. All rights reserved.',
           size=7, color=DARK_GREEN, align=PP_ALIGN.CENTER)

    print('  [Slide 27] 고생하셨습니다 슬라이드 완료')
    return slide


# ══════════════════════════════════════════════════════════
if __name__ == '__main__':
    print(f'현재 슬라이드 수: {len(prs.slides)}\n')
    print('=== 마무리 섹션 3슬라이드 재구성 ===\n')

    rebuild_slide25()
    add_slide26()
    add_slide27()

    prs.save(DST)
    print(f'\n✓ 저장 완료: {DST}')
    print(f'  총 슬라이드 수: {len(prs.slides)}')
