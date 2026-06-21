import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches

SRC = r'C:\Users\jay\Downloads\AI교육_완성판_v6.pptx'
prs = Presentation(SRC)

print('=== Slide 2 (목차) ===')
slide2 = prs.slides[1]
for i, shape in enumerate(slide2.shapes):
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            txt = para.text.strip()
            if txt:
                print(f'  [{i}] "{txt}"')

print('\n=== Slide 6 현재 item 02 내용 ===')
slide6 = prs.slides[5]
for i, shape in enumerate(slide6.shapes):
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            txt = para.text.strip()
            if txt:
                print(f'  [{i}] "{txt}"')

print('\n=== Slide 18 전체 shapes ===')
slide18 = prs.slides[17]
for i, shape in enumerate(slide18.shapes):
    print(f'  [{i}] type={shape.shape_type} L={shape.left/914400:.2f}" T={shape.top/914400:.2f}" W={shape.width/914400:.2f}" H={shape.height/914400:.2f}"')
    if shape.has_text_frame:
        for j, para in enumerate(shape.text_frame.paragraphs):
            txt = para.text.strip()
            if txt:
                print(f'       para[{j}]: "{txt}"')

print('\n=== 전체 슬라이드 "교시" 텍스트 검색 ===')
for si, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                txt = para.text
                if '교시' in txt:
                    print(f'  Slide {si+1}: "{txt.strip()}"')
