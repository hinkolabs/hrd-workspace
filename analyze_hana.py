import sys
import os
sys.stdout.reconfigure(encoding='utf-8')
sys.stderr.reconfigure(encoding='utf-8')
from pptx import Presentation
import docx
import openpyxl

BASE = r'C:\Users\jay\Downloads'

# ============================
# PPT 분석
# ============================
ppt_files = [
    '하나증권_AI교육_Part1_오프닝+1교시.pptx',
    '하나증권_AI교육_Part2_2교시+3교시.pptx',
    '하나증권_AI교육_Part3_4교시+5교시+마무리.pptx',
]

ppt_summary = {}
for fname in ppt_files:
    fpath = os.path.join(BASE, fname)
    prs = Presentation(fpath)
    slides_data = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        slides_data.append(texts)
    ppt_summary[fname] = slides_data
    print(f'\n{"="*60}')
    print(f'파일: {fname}')
    print(f'총 슬라이드 수: {len(slides_data)}')
    for i, slide_texts in enumerate(slides_data):
        if slide_texts:
            joined = ' | '.join(slide_texts[:6])
            print(f'  슬라이드 {i+1:2d}: {joined[:120]}')

# ============================
# Word 강사대본 분석
# ============================
print(f'\n{"="*60}')
print('강사대본 분석')
doc_path = os.path.join(BASE, '하나증권_AI교육_강사대본.docx')
doc = docx.Document(doc_path)
print(f'총 단락 수: {len(doc.paragraphs)}')
for i, para in enumerate(doc.paragraphs[:100]):
    if para.text.strip():
        style = para.style.name if para.style else ''
        print(f'  [{i:3d}] [{style}] {para.text[:120]}')

# ============================
# Excel 분석
# ============================
for excel_fname in [
    'hana_ai_excel_practice_student_clean.xlsx',
    'hana_ai_excel_practice_instructor_clean.xlsx',
]:
    print(f'\n{"="*60}')
    print(f'Excel: {excel_fname}')
    wb = openpyxl.load_workbook(os.path.join(BASE, excel_fname))
    print(f'시트 목록: {wb.sheetnames}')
    for shname in wb.sheetnames:
        ws = wb[shname]
        print(f'\n  [시트: {shname}] 크기: {ws.max_row}행 x {ws.max_column}열')
        for row_idx, row in enumerate(ws.iter_rows(values_only=True)):
            if row_idx >= 20:
                print('  ...(이하 생략)...')
                break
            if any(v is not None for v in row):
                vals = [str(v)[:30] if v is not None else '' for v in row[:10]]
                print(f'    행{row_idx+1:3d}: {" | ".join(vals)}')
