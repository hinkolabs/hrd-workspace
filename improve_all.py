"""
하나증권 AI 교육 자료 통합 개선 스크립트
- 큐시트 DOCX 개선
- PPT Part2 개선 (허용 예시 슬라이드 추가)
- PPT Part3 개선 (5교시 필수/선택 레이블)
- Excel 학생용 파일 개선 (README 내비게이션)
"""
import sys
import os
import copy
sys.stdout.reconfigure(encoding='utf-8')
sys.stderr.reconfigure(encoding='utf-8')

BASE = r'C:\Users\jay\Downloads'
OUT = os.path.join(BASE, '개선판')
os.makedirs(OUT, exist_ok=True)

# ─────────────────────────────────────────────
# 공통 헬퍼 (DOCX)
# ─────────────────────────────────────────────
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Pt, RGBColor


def _make_run_xml(text, bold=False, color_hex=None):
    r = OxmlElement('w:r')
    rPr = OxmlElement('w:rPr')
    if bold:
        b = OxmlElement('w:b')
        rPr.append(b)
    if color_hex:
        clr = OxmlElement('w:color')
        clr.set(qn('w:val'), color_hex)
        rPr.append(clr)
    if bold or color_hex:
        r.insert(0, rPr)
    t = OxmlElement('w:t')
    t.text = text
    t.set('{http://www.w3.org/XML/1998/namespace}space', 'preserve')
    r.append(t)
    return r


def prepend_para_to_cell(cell, text, bold=False, color_hex=None):
    """셀 맨 앞에 새 단락 삽입"""
    tc = cell._tc
    p = OxmlElement('w:p')
    p.append(_make_run_xml(text, bold=bold, color_hex=color_hex))
    first_p = tc.find(qn('w:p'))
    if first_p is not None:
        tc.insert(list(tc).index(first_p), p)
    else:
        tc.append(p)


def append_para_to_cell(cell, text, bold=False, color_hex=None):
    """셀 끝에 새 단락 추가"""
    p = OxmlElement('w:p')
    p.append(_make_run_xml(text, bold=bold, color_hex=color_hex))
    cell._tc.append(p)


def replace_text_in_para(para, old, new):
    """단락 텍스트에서 old → new 교체 (run 분리 대응)"""
    full = ''.join(r.text for r in para.runs)
    if old not in full:
        return False
    new_full = full.replace(old, new)
    if para.runs:
        para.runs[0].text = new_full
        for r in para.runs[1:]:
            r.text = ''
    else:
        para.add_run(new_full)
    return True


def replace_text_in_cell(cell, old, new):
    """셀 내 모든 단락에서 old → new 교체"""
    changed = False
    for para in cell.paragraphs:
        if replace_text_in_para(para, old, new):
            changed = True
    return changed


# ─────────────────────────────────────────────
# 1. 큐시트 DOCX 개선
# ─────────────────────────────────────────────
def improve_cuesheet():
    from docx import Document

    src = os.path.join(BASE, '하나증권_신입연수_AI교육_강사용_진행큐시트.docx')
    dst = os.path.join(OUT, '하나증권_신입연수_AI교육_강사용_진행큐시트_개선판.docx')
    doc = Document(src)

    # ── Fix 1~4: 목차 단락 타이밍 표기 보정 ──
    timing_fixes = {
        '09:00-09:05': '09:00-09:05 (1교시 오프닝 포함)',
        # 1교시 줄에서 겹침 해소
        '3. 1교시 시작: 09:00-09:50': '3. 1교시 시작: 09:00-09:50 (오프닝 5분 포함)',
        '13:00-13:03': '13:00-13:03 (4교시 도입부 포함)',
        '7. 4교시 시작: 13:00-13:50': '7. 4교시 시작: 13:00-13:50 (집중 복귀 3분 포함)',
    }
    for para in doc.paragraphs:
        for old, new in timing_fixes.items():
            if old in para.text:
                replace_text_in_para(para, old, new)

    # ── Fix 5: Table[6] 2교시 퀴즈 운영팁 → 퀴즈 A~E 시나리오 기준표 앞에 추가 ──
    quiz_header_lines = [
        '[ 퀴즈 A~E 판단 기준 — 강사 참고 ]',
        'A. 고객 이름 + 실거래내역 입력  →  ❌ 불가  (실제 고객개인정보)',
        'B. 가상이름 + 가상포트폴리오 시나리오  →  ✅ 가능  (실명·실계좌 없으면 허용, 실제 고객 힌트 주의)',
        'C. 공시된 사업보고서 + "호재야 악재야?" 판단 요청  →  ⚠️ 주의  (공개자료라도 AI 해석 = 미공개 의견 위험)',
        'D. 내일 배포 예정 리서치 초안 붙여넣기  →  ❌ 불가  (발간 전 자료: 가장 흔한 실수)',
        'E. 어제 공시된 분기실적 요약 요청  →  ✅ 가능  (공개자료 기반, 해석은 본인의견임을 명시)',
        '─────────────────────────────',
    ]
    quiz_cell = doc.tables[6].cell(0, 0)
    for line in reversed(quiz_header_lines):
        bold = line.startswith('[')
        prepend_para_to_cell(quiz_cell, line, bold=bold,
                              color_hex='1F3864' if bold else None)

    # ── Fix 6: Table[8] 3교시 운영 기준 → STEP 목표 표현 통일 ──
    step_cell = doc.tables[8].cell(0, 0)
    replace_text_in_cell(
        step_cell,
        '전원 완료 목표는 STEP1 + 샘플 검증이다',
        '최소 목표(전원): STEP1 + 샘플 검증 (= 필수). 선택: STEP2. 심화(빠른 수강생 전용): STEP3'
    )
    # 기존 표현이 약간 다를 수 있으므로 유사 표현도 처리
    replace_text_in_cell(
        step_cell,
        '전원 완료 목표는 STEP1 + 샘플 검증',
        '최소 목표(전원): STEP1 + 샘플 검증 (= 필수). 선택: STEP2. 심화(빠른 수강생): STEP3'
    )

    # ── Fix 7: Table[12] 5교시 진행표 → 형식 선택에 소요시간 추가 ──
    five_table = doc.tables[12]
    format_replacements = [
        ('카드뉴스 5장',     '카드뉴스 3장 (약 20분)'),
        ('카드뉴스',         '카드뉴스 3장 (약 20분)'),  # 기타 표현
        ('30초 영상 대본',   '30초 영상 대본 (약 20분)'),
        ('1페이지 안내문',   '1페이지 안내문 (약 15분)'),
        ('캠페인 문구',      '캠페인 문구 (약 10분)'),
        # 추천 문구도 업데이트
        ('처음 하는 분은 카드뉴스를 추천합니다',
         '처음 하는 분은 캠페인 문구 또는 1페이지 안내문을 추천합니다'),
    ]
    for row in five_table.rows:
        for cell in row.cells:
            if '카드뉴스' in cell.text or '영상 대본' in cell.text:
                for old, new in format_replacements:
                    # 중복 치환 방지 (이미 분/약 이 들어간 경우 skip)
                    if '약' not in old:
                        replace_text_in_cell(cell, old, new)
                    else:
                        replace_text_in_cell(cell, old, new)

    # ── Fix 8: Table[15] 체크리스트 → D-1/D-day 타이밍 그룹 추가 ──
    check_table = doc.tables[15]
    # 각 행의 col[1] 내용을 보고 col[0] 레이블을 세분화
    timing_map = {
        'PPT': 'D-1 전날',
        '강사 파일': 'D-1 전날',
        '학생용 파일': 'D-day 30분 전',
        'ALLI 접속 계정': 'D-day 30분 전 ★',
        'ALLI 접속 실패': 'D-day 30분 전 ★',
    }
    for row in check_table.rows:
        if len(row.cells) < 2:
            continue
        label_cell = row.cells[0]
        content_text = row.cells[1].text
        for keyword, new_label in timing_map.items():
            if keyword in content_text and '강의 전' in label_cell.text:
                replace_text_in_cell(label_cell, '강의 전', new_label)
                break

    doc.save(dst)
    print(f'[완료] 큐시트 개선판: {dst}')


# ─────────────────────────────────────────────
# PPT 공통 헬퍼
# ─────────────────────────────────────────────
def _insert_slide_after(prs, after_idx):
    """after_idx 슬라이드 바로 뒤에 동일 레이아웃으로 빈 슬라이드 삽입"""
    layout = prs.slides[after_idx].slide_layout
    new_slide = prs.slides.add_slide(layout)
    # 끝에 추가된 슬라이드를 after_idx+1 위치로 이동
    xml_slides = prs.slides._sldIdLst
    moved = xml_slides[-1]
    xml_slides.remove(moved)
    xml_slides.insert(after_idx + 1, moved)
    return new_slide


def _add_textbox(slide, left, top, width, height, text,
                 font_size=18, bold=False, color=(0x1F, 0x38, 0x64),
                 bg_color=None, wrap=True):
    from pptx.util import Emu, Pt as PPt
    from pptx.dml.color import RGBColor as PRGB
    from pptx.enum.text import PP_ALIGN

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = PPt(font_size)
    run.font.bold = bold
    run.font.color.rgb = PRGB(*color)
    if bg_color:
        from pptx.oxml.ns import qn
        from lxml import etree
        spPr = txBox._element.spPr
        solidFill = etree.SubElement(spPr, qn('a:solidFill'))
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', '%02X%02X%02X' % bg_color)
    return txBox


# ─────────────────────────────────────────────
# 2. PPT Part2 개선: 허용 예시 슬라이드 삽입
# ─────────────────────────────────────────────
def improve_ppt_part2():
    from pptx import Presentation
    from pptx.util import Inches, Pt as PPt, Emu
    from pptx.dml.color import RGBColor as PRGB

    src = os.path.join(BASE, '하나증권_AI교육_Part2_2교시+3교시.pptx')
    dst = os.path.join(OUT, '하나증권_AI교육_Part2_2교시+3교시_개선판.pptx')
    prs = Presentation(src)

    # 슬라이드 5 (index 4) = 시간 제약 금지 슬라이드
    # 슬라이드 6 (index 5) = AI 입력 3가지 금지 요약
    # ▶ 허용 예시 슬라이드를 index 5 바로 뒤(index 6)에 삽입
    new_slide = _insert_slide_after(prs, 5)

    # 기존 placeholder가 있으면 초기화
    for ph in new_slide.placeholders:
        try:
            ph.text = ''
        except Exception:
            pass

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # ── 섹션 레이블 ──
    _add_textbox(new_slide,
                 left=Inches(0.4), top=Inches(0.25),
                 width=Inches(3), height=Inches(0.45),
                 text='2교시 · 증권의 AI 컴플라이언스',
                 font_size=10, bold=False,
                 color=(0x80, 0x80, 0x80))

    # ── 타이틀 ──
    _add_textbox(new_slide,
                 left=Inches(0.4), top=Inches(0.7),
                 width=Inches(9), height=Inches(0.8),
                 text='반면, 이런 것들은 AI로 해도 됩니다',
                 font_size=28, bold=True,
                 color=(0x1F, 0x38, 0x64))

    # ── 허용 예시 4가지 ──
    allowed = [
        ('✅', '공개된 자료 기반 요약',
         '공시자료 · 뉴스 · IR 자료 등 이미 공개된 정보를 붙여넣고 요약·정리'),
        ('✅', '가상 시나리오 작성',
         '실명·실 계좌번호 없이 "가상 고객 A"로 설정해 투자 시나리오 초안 작성'),
        ('✅', '업무 초안 작성 (검증 전 참고용)',
         '이메일·보고서 초안, 회의록 요약 등 — 반드시 최종 제출 전 검증'),
        ('✅', '반복 작업 구조화',
         '엑셀 수식 설계, 데이터 집계 로직, 정형화된 문서 구조 설계'),
    ]

    box_w = Inches(2.15)
    box_h = Inches(2.0)
    gap = Inches(0.15)
    top_y = Inches(1.7)
    start_x = Inches(0.4)

    for i, (icon, title, desc) in enumerate(allowed):
        x = start_x + i * (box_w + gap)
        # 아이콘
        _add_textbox(new_slide, x, top_y,
                     Inches(0.4), Inches(0.45),
                     icon, font_size=20, bold=False,
                     color=(0x22, 0x8B, 0x22))
        # 제목
        _add_textbox(new_slide, x + Inches(0.38), top_y,
                     Inches(1.7), Inches(0.45),
                     title, font_size=13, bold=True,
                     color=(0x1F, 0x38, 0x64))
        # 설명
        _add_textbox(new_slide, x, top_y + Inches(0.48),
                     box_w, Inches(1.4),
                     desc, font_size=11, bold=False,
                     color=(0x33, 0x33, 0x33))

    # ── 주의 문구 (하단) ──
    _add_textbox(new_slide,
                 left=Inches(0.4), top=Inches(3.9),
                 width=Inches(9), height=Inches(0.55),
                 text='⚠️  단, 모든 출력물은 최종 제출 전 반드시 검증하세요. '
                      'AI 결과물에는 오류·환각이 포함될 수 있습니다.',
                 font_size=12, bold=False,
                 color=(0xC5, 0x50, 0x00))

    # ── 슬라이드 번호 (우하단) ──
    total = len(prs.slides)
    _add_textbox(new_slide,
                 left=slide_w - Inches(0.8), top=slide_h - Inches(0.4),
                 width=Inches(0.6), height=Inches(0.35),
                 text=str(6 + 1),   # 삽입 후 슬라이드 번호
                 font_size=9, color=(0xAA, 0xAA, 0xAA))

    prs.save(dst)
    print(f'[완료] Part2 개선판: {dst}')


# ─────────────────────────────────────────────
# 3. PPT Part3 개선: 5교시 8개 과제 필수/선택 레이블
# ─────────────────────────────────────────────
def improve_ppt_part3():
    from pptx import Presentation
    from pptx.util import Inches, Pt as PPt

    src = os.path.join(BASE, '하나증권_AI교육_Part3_4교시+5교시+마무리.pptx')
    dst = os.path.join(OUT, '하나증권_AI교육_Part3_4교시+5교시+마무리_개선판.pptx')
    prs = Presentation(src)

    # 슬라이드 9 (index 8) = "제안 과제 8가지"
    target_slide = prs.slides[8]

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # ── 필수/선택 범례 박스 (슬라이드 상단) ──
    _add_textbox(target_slide,
                 left=Inches(0.4), top=Inches(0.3),
                 width=Inches(9.2), height=Inches(0.45),
                 text='★ 필수 (1~3번): 오늘 반드시 완성     ○ 선택 (4~8번): 시간이 허락되면 추가 도전',
                 font_size=12, bold=True,
                 color=(0xFF, 0xFF, 0xFF),
                 bg_color=(0x1F, 0x38, 0x64))

    # ── 기존 번호 텍스트박스를 찾아 필수/선택 마커 추가 ──
    required_nums = {'1', '2', '3'}
    optional_nums = {'4', '5', '6', '7', '8'}

    for shape in target_slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()
        if txt in required_nums:
            # 번호 앞에 ★ 삽입
            p = shape.text_frame.paragraphs[0]
            if p.runs:
                p.runs[0].text = '★' + p.runs[0].text
        elif txt in optional_nums:
            p = shape.text_frame.paragraphs[0]
            if p.runs:
                p.runs[0].text = '○' + p.runs[0].text

    prs.save(dst)
    print(f'[완료] Part3 개선판: {dst}')


# ─────────────────────────────────────────────
# 4. Excel 학생용 파일 개선: README 내비게이션 강화
# ─────────────────────────────────────────────
def improve_excel_student():
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    src = os.path.join(BASE, 'hana_ai_excel_practice_student_clean.xlsx')
    dst = os.path.join(OUT, 'hana_ai_excel_practice_student_clean_개선판.xlsx')
    wb = openpyxl.load_workbook(src)

    # ── 00_시작전_README 시트 개선 ──
    ws = wb['00_수강생_README']

    # 1) 헤더 위에 "★ 지금 바로 시작하는 법" 안내 행 삽입
    ws.insert_rows(1)
    ws.insert_rows(1)

    guide_cell = ws.cell(row=1, column=1)
    guide_cell.value = '★ 지금 바로 여기서 시작하세요  →  먼저 [거래원장_RAW] 시트를 열어 데이터를 확인하세요'
    guide_cell.font = Font(name='맑은 고딕', size=13, bold=True, color='FFFFFF')
    guide_cell.fill = PatternFill(fill_type='solid', fgColor='1F3864')
    guide_cell.alignment = Alignment(horizontal='left', vertical='center', wrap_text=True)
    ws.row_dimensions[1].height = 30
    # 가로 병합
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=5)

    goal_cell = ws.cell(row=2, column=1)
    goal_cell.value = '✅ 오늘의 목표: STEP1까지 완성하면 충분합니다. STEP2·STEP3은 선택 과제입니다.'
    goal_cell.font = Font(name='맑은 고딕', size=11, bold=False, color='1F3864')
    goal_cell.fill = PatternFill(fill_type='solid', fgColor='D9E1F2')
    goal_cell.alignment = Alignment(horizontal='left', vertical='center', wrap_text=True)
    ws.row_dimensions[2].height = 22
    ws.merge_cells(start_row=2, start_column=1, end_row=2, end_column=5)

    # 2) STEP별 사용 시트 열(D열)에 시간 안내 추가
    step_time_map = {
        'STEP 1': '(8~16분)',
        'STEP 2': '(23~33분)',
        'STEP 3': '(33~43분)',
    }
    for row in ws.iter_rows():
        for cell in row:
            if cell.value and isinstance(cell.value, str):
                for step, timing in step_time_map.items():
                    if step in cell.value and timing not in cell.value:
                        cell.value = cell.value.replace(step, f'{step} {timing}')

    # 3) 시트 탭 색상으로 시작 시트 강조 (거래명세_RAW)
    if '거래원장_RAW' in wb.sheetnames:
        wb['거래원장_RAW'].sheet_properties.tabColor = '1F3864'
    if '결과작성_STEP1' in wb.sheetnames:
        wb['결과작성_STEP1'].sheet_properties.tabColor = '70AD47'
    if '결과작성_STEP2' in wb.sheetnames:
        wb['결과작성_STEP2'].sheet_properties.tabColor = 'ED7D31'
    if '결과작성_STEP3' in wb.sheetnames:
        wb['결과작성_STEP3'].sheet_properties.tabColor = 'FF0000'

    # ── 01_실습흐름 시트: 오늘의 목표 기준 강조 ──
    schedule_name = None
    for name in wb.sheetnames:
        if '실습흐름' in name or '01_' in name:
            schedule_name = name
            break
    if schedule_name:
        ws2 = wb[schedule_name]
        # 최상단에 목표 안내 행 삽입
        ws2.insert_rows(1)
        c = ws2.cell(row=1, column=1)
        c.value = '최소 목표(전원): STEP1 + 샘플 검증  |  선택: STEP2  |  심화(빠른 수강생): STEP3'
        c.font = Font(name='맑은 고딕', size=11, bold=True, color='FFFFFF')
        c.fill = PatternFill(fill_type='solid', fgColor='375623')
        c.alignment = Alignment(horizontal='left', vertical='center')
        ws2.row_dimensions[1].height = 22
        ws2.merge_cells(start_row=1, start_column=1, end_row=1, end_column=5)

    wb.save(dst)
    print(f'[완료] Excel 학생용 개선판: {dst}')


# ─────────────────────────────────────────────
# 실행
# ─────────────────────────────────────────────
if __name__ == '__main__':
    print('=== 하나증권 AI 교육 자료 개선 시작 ===\n')
    try:
        improve_cuesheet()
    except Exception as e:
        print(f'[오류] 큐시트: {e}')
        import traceback; traceback.print_exc()

    try:
        improve_ppt_part2()
    except Exception as e:
        print(f'[오류] PPT Part2: {e}')
        import traceback; traceback.print_exc()

    try:
        improve_ppt_part3()
    except Exception as e:
        print(f'[오류] PPT Part3: {e}')
        import traceback; traceback.print_exc()

    try:
        improve_excel_student()
    except Exception as e:
        print(f'[오류] Excel: {e}')
        import traceback; traceback.print_exc()

    print('\n=== 완료 ===')
    print(f'출력 폴더: {OUT}')
