import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

SRC = r'C:\Users\jay\Downloads\AI교육_완성판_v6.pptx'
prs = Presentation(SRC)

target_slides = [10, 16, 19, 20, 21, 23, 24]  # 1-based
for si in target_slides:
    slide = prs.slides[si - 1]
    print(f'\n=== Slide {si} ===')
    for i, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            for j, para in enumerate(shape.text_frame.paragraphs):
                txt = para.text.strip()
                if '교시' in txt or (si == 10 and txt):
                    print(f'  shape[{i}] para[{j}]: "{txt}"')
