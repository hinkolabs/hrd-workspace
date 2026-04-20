"""
Create executive-style PPT applying design from '22년 공통시스템 운영안' template
to the content of '경영지원그룹장 CEO업무보고' presentation.
"""
import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Cm
import copy
from lxml import etree

# ─── Color Palette (from design template) ───────────────────────────────────
RED     = RGBColor(0xEF, 0x35, 0x35)   # primary accent
DARK    = RGBColor(0x26, 0x26, 0x26)   # main text
GRAY    = RGBColor(0x65, 0x61, 0x67)   # secondary text
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xF5, 0xF5, 0xF5)  # light gray bg
DGRAY   = RGBColor(0x44, 0x44, 0x44)  # dark gray
MGRAY   = RGBColor(0x9E, 0x9E, 0x9E)  # medium gray
BLUE    = RGBColor(0x1E, 0x4D, 0x9B)  # secondary accent (dark navy blue)
LBLUE   = RGBColor(0xE8, 0xEE, 0xF8)  # light blue bg
LTRED   = RGBColor(0xFD, 0xED, 0xED)  # light red bg

# ─── Slide Dimensions ────────────────────────────────────────────────────────
W = Inches(13.33)   # 10" → 13.33" widescreen
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK_LAYOUT = prs.slide_layouts[6]  # blank layout


# ═══════════════════════════════════════════════════════════════════════════════
# Helper Functions
# ═══════════════════════════════════════════════════════════════════════════════

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=Pt(11),
                bold=False, color=DARK, align=PP_ALIGN.LEFT,
                wrap=True, italic=False, line_spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def set_para(para, text, font_size=Pt(11), bold=False, color=DARK,
             align=PP_ALIGN.LEFT, italic=False, space_before=None, space_after=None):
    para.clear()
    para.alignment = align
    if space_before:
        para.space_before = space_before
    if space_after:
        para.space_after = space_after
    run = para.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_multi_para_textbox(slide, left, top, width, height, paras, wrap=True):
    """paras: list of dicts with keys: text, font_size, bold, color, align, italic, space_before, space_after"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    for i, p_info in enumerate(paras):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        set_para(
            para,
            p_info.get('text', ''),
            font_size=p_info.get('font_size', Pt(11)),
            bold=p_info.get('bold', False),
            color=p_info.get('color', DARK),
            align=p_info.get('align', PP_ALIGN.LEFT),
            italic=p_info.get('italic', False),
            space_before=p_info.get('space_before'),
            space_after=p_info.get('space_after'),
        )
    return txBox


def add_slide_header(slide, title_text, slide_num=None, subtitle=None):
    """Add top header bar with red accent line"""
    # White header background
    add_rect(slide, 0, 0, W, Inches(0.9), fill_color=WHITE)
    # Red left accent bar
    add_rect(slide, 0, 0, Inches(0.18), Inches(0.9), fill_color=RED)
    # Title text
    add_textbox(slide, Inches(0.3), Inches(0.12), Inches(11), Inches(0.65),
                title_text, font_size=Pt(18), bold=True, color=DARK)
    # Red separator line
    line_shape = slide.shapes.add_shape(1, 0, Inches(0.9), W, Pt(2))
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = RED
    line_shape.line.fill.background()
    # Slide number (top right)
    if slide_num:
        add_textbox(slide, Inches(12.3), Inches(0.15), Inches(0.8), Inches(0.5),
                    str(slide_num), font_size=Pt(11), bold=False, color=MGRAY,
                    align=PP_ALIGN.RIGHT)
    # Bottom footer bar
    add_rect(slide, 0, Inches(7.1), W, Inches(0.4), fill_color=DARK)
    add_textbox(slide, Inches(0.3), Inches(7.12), Inches(6), Inches(0.3),
                '경영지원그룹  정기환 전무', font_size=Pt(9), bold=False, color=MGRAY)
    add_textbox(slide, Inches(10), Inches(7.12), Inches(3), Inches(0.3),
                '2026.04.20', font_size=Pt(9), bold=False, color=MGRAY, align=PP_ALIGN.RIGHT)


def add_section_tag(slide, left, top, text, bg=RED, fg=WHITE, width=Inches(2.5), height=Inches(0.35)):
    shape = add_rect(slide, left, top, width, height, fill_color=bg)
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = fg
    return shape


def add_bullet_box(slide, left, top, width, height, title, bullets,
                   title_bg=RED, title_fg=WHITE, body_bg=LGRAY,
                   title_height=Inches(0.42)):
    """Box with colored title bar and bullet list body"""
    # Title bar
    title_rect = add_rect(slide, left, top, width, title_height, fill_color=title_bg)
    tf = title_rect.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = title_fg
    # Body
    body_height = height - title_height
    body_rect = add_rect(slide, left, top + title_height, width, body_height, fill_color=body_bg)
    body_rect.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    body_rect.line.width = Pt(0.5)
    # Text in body
    padding = Inches(0.12)
    para_infos = []
    for b in bullets:
        if isinstance(b, dict):
            para_infos.append(b)
        else:
            para_infos.append({'text': b, 'font_size': Pt(10.5), 'color': DARK})
    add_multi_para_textbox(slide, left + padding, top + title_height + padding,
                           width - padding*2, body_height - padding*2, para_infos)
    return body_rect


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — 표지 (Cover)
# ═══════════════════════════════════════════════════════════════════════════════
def create_slide1():
    slide = prs.slides.add_slide(BLANK_LAYOUT)

    # Full background
    add_rect(slide, 0, 0, W, H, fill_color=WHITE)

    # Left dark sidebar
    sidebar_w = Inches(4.2)
    add_rect(slide, 0, 0, sidebar_w, H, fill_color=DARK)

    # Red accent strip on left sidebar
    add_rect(slide, sidebar_w - Inches(0.08), 0, Inches(0.08), H, fill_color=RED)

    # Title area on left sidebar
    add_textbox(slide, Inches(0.35), Inches(1.2), Inches(3.6), Inches(0.5),
                '그룹 CEO 업무보고', font_size=Pt(13), bold=False, color=MGRAY)

    add_textbox(slide, Inches(0.35), Inches(1.7), Inches(3.6), Inches(1.2),
                '관계사 협업\n활성화 방안', font_size=Pt(28), bold=True, color=WHITE)

    # Red underline
    add_rect(slide, Inches(0.35), Inches(3.0), Inches(2.5), Pt(3), fill_color=RED)

    # Subtitle info
    add_textbox(slide, Inches(0.35), Inches(3.2), Inches(3.6), Inches(0.4),
                '시너지 + 기업문화 + 인사의 삼위일체', font_size=Pt(11), bold=False, color=LGRAY)

    # Presenter
    add_textbox(slide, Inches(0.35), Inches(5.8), Inches(3.6), Inches(0.35),
                '경영지원그룹장', font_size=Pt(10), bold=False, color=MGRAY)
    add_textbox(slide, Inches(0.35), Inches(6.15), Inches(3.6), Inches(0.45),
                '정기환 전무', font_size=Pt(16), bold=True, color=WHITE)
    add_textbox(slide, Inches(0.35), Inches(6.6), Inches(3.6), Inches(0.35),
                '2026. 04. 20', font_size=Pt(10), bold=False, color=MGRAY)

    # Right side — content overview
    right_x = sidebar_w + Inches(0.5)
    right_w = W - sidebar_w - Inches(0.6)

    # WHY? Block
    add_textbox(slide, right_x, Inches(0.5), right_w, Inches(0.45),
                '관계사 협업?  WHY?', font_size=Pt(15), bold=True, color=RED)
    add_rect(slide, right_x, Inches(0.95), right_w, Pt(1.5), fill_color=RGBColor(0xDD, 0xDD, 0xDD))

    why_left_w = (right_w - Inches(0.3)) / 2
    # WHY Left box
    left_box = add_rect(slide, right_x, Inches(1.05), why_left_w, Inches(2.0), fill_color=LBLUE)
    left_box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    left_box.line.width = Pt(0.5)
    add_multi_para_textbox(slide, right_x + Inches(0.12), Inches(1.1), why_left_w - Inches(0.2), Inches(1.85), [
        {'text': '매년 사업계획 회의장표에 빠지지 않았던 아젠다', 'font_size': Pt(10), 'color': DGRAY},
        {'text': '', 'font_size': Pt(6)},
        {'text': '→ 그룹 콜라보 활성화!', 'font_size': Pt(11.5), 'bold': True, 'color': BLUE},
        {'text': '→ 관계사 협업 확대!', 'font_size': Pt(11.5), 'bold': True, 'color': BLUE},
    ])

    # WHY Right box
    right_box_x = right_x + why_left_w + Inches(0.3)
    right_box = add_rect(slide, right_box_x, Inches(1.05), why_left_w, Inches(2.0), fill_color=LTRED)
    right_box.line.color.rgb = RGBColor(0xEE, 0xBB, 0xBB)
    right_box.line.width = Pt(0.5)
    add_multi_para_textbox(slide, right_box_x + Inches(0.12), Inches(1.1), why_left_w - Inches(0.2), Inches(1.85), [
        {'text': '그러나, 현실은?', 'font_size': Pt(10), 'color': DGRAY, 'bold': True},
        {'text': '', 'font_size': Pt(6)},
        {'text': '✗  관계사 상품 이해·관심 無!', 'font_size': Pt(10.5), 'color': RGBColor(0xC0,0x30,0x30)},
        {'text': '✗  목표부여! 실적체크!', 'font_size': Pt(10.5), 'color': RGBColor(0xC0,0x30,0x30)},
        {'text': '✗  내가 왜? 카드와 보험을?', 'font_size': Pt(10.5), 'color': RGBColor(0xC0,0x30,0x30)},
    ])

    # Competition data
    comp_y = Inches(3.15)
    comp_box = add_rect(slide, right_x, comp_y, right_w, Inches(0.75), fill_color=DARK)
    add_multi_para_textbox(slide, right_x + Inches(0.15), comp_y + Inches(0.08), right_w - Inches(0.3), Inches(0.6), [
        {'text': '경쟁사 비은행 실적 기여도', 'font_size': Pt(9.5), 'color': MGRAY},
        {'text': '38%    29%    12%', 'font_size': Pt(14), 'bold': True, 'color': WHITE, 'align': PP_ALIGN.CENTER},
    ])

    # HOW section
    add_textbox(slide, right_x, Inches(4.05), right_w, Inches(0.45),
                '관계사 협업?  HOW?  —  삼위일체 전략', font_size=Pt(15), bold=True, color=RED)
    add_rect(slide, right_x, Inches(4.5), right_w, Pt(1.5), fill_color=RGBColor(0xDD, 0xDD, 0xDD))

    box_w = (right_w - Inches(0.4)) / 3
    how_y = Inches(4.6)
    how_h = Inches(2.3)
    how_items = [
        ('시너지', ['콜라보북 제작·배포', '그룹별 협업 마케팅', '실적 주간 점검 (CEO 직접 관리)'], BLUE),
        ('기업문화', ['칭찬 웨이브', '현장 격려 이벤트', '본업 캠페인'], RED),
        ('인사', ['협업 우수직원 인사 가점', '우수 연수 기회 우선배정', '분기별 표창'], RGBColor(0x2E,0x7D,0x32)),
    ]
    for k, (label, items, col) in enumerate(how_items):
        bx = right_x + k * (box_w + Inches(0.2))
        # Title
        title_rect = add_rect(slide, bx, how_y, box_w, Inches(0.38), fill_color=col)
        tf = title_rect.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = WHITE
        # Body
        body_rect = add_rect(slide, bx, how_y + Inches(0.38), box_w, how_h - Inches(0.38), fill_color=LGRAY)
        body_rect.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        body_rect.line.width = Pt(0.5)
        bullet_paras = []
        for item in items:
            bullet_paras.append({'text': f'• {item}', 'font_size': Pt(10), 'color': DGRAY})
        add_multi_para_textbox(slide, bx + Inches(0.1), how_y + Inches(0.45), box_w - Inches(0.15), how_h - Inches(0.5), bullet_paras)

    # Page number
    add_textbox(slide, W - Inches(0.7), Inches(7.1), Inches(0.5), Inches(0.35),
                '1', font_size=Pt(9), color=MGRAY, align=PP_ALIGN.CENTER)

    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — HOW? 기업문화
# ═══════════════════════════════════════════════════════════════════════════════
def create_slide2():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_rect(slide, 0, 0, W, H, fill_color=WHITE)
    add_slide_header(slide, '관계사 협업  HOW?  ① 기업문화', slide_num=2)

    content_top = Inches(1.05)
    content_h   = Inches(5.85)
    pad = Inches(0.25)
    col_w = (W - pad*3) / 2  # two-column wait full width - now use 3 columns

    # Timeline table (top full-width)
    tl_h = Inches(1.3)
    tl_y = content_top
    tl_titles = ['2023 ~ 2024', '2025', '2026']
    tl_bodies = [
        '소통과 참여, 변화와 혁신의 시작\nOne Team, One Spirit',
        "하나금융의 DNA로\n'다시, 하나답게' 캠페인",
        '역동적인 하나증권\nCheer Up! Level Up!\n협업문화 장착'
    ]
    tl_colors = [DGRAY, BLUE, RED]
    col3_w = (W - pad*4) / 3
    for k in range(3):
        bx = pad + k * (col3_w + pad)
        # Active indicator
        if k == 2:
            active_rect = add_rect(slide, bx, tl_y, col3_w, tl_h, fill_color=LTRED)
            active_rect.line.color.rgb = RED
            active_rect.line.width = Pt(1.5)
        else:
            body_rect = add_rect(slide, bx, tl_y, col3_w, tl_h, fill_color=LGRAY)
            body_rect.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            body_rect.line.width = Pt(0.5)
        # Year label
        yr_rect = add_rect(slide, bx, tl_y, col3_w, Inches(0.36), fill_color=tl_colors[k])
        tf = yr_rect.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = tl_titles[k]
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = WHITE
        # Body text
        body_c = RED if k == 2 else DGRAY
        add_textbox(slide, bx + Inches(0.1), tl_y + Inches(0.42), col3_w - Inches(0.15), tl_h - Inches(0.5),
                    tl_bodies[k], font_size=Pt(10.5), bold=(k==2), color=body_c, align=PP_ALIGN.CENTER)

    # Three strategy boxes below
    box_y = content_top + tl_h + Inches(0.2)
    box_h = Inches(4.2)
    strategies = [
        {
            'title': '① 칭찬 웨이브',
            'color': RED,
            'points': [
                '매주 협업 우수사례 공유',
                '상위 우수 직원에 대한 칭찬 포인트 지급',
                '  (월·분기·기단위 포상)',
                '사내 복지포인트 전환',
                '휴양소 지원',
                'CES 등의 교육기회 부여',
            ]
        },
        {
            'title': '② 현장 격려 이벤트',
            'color': BLUE,
            'points': [
                '증권-은행 복합점포 등 거점점포 중심의',
                '  격려 이벤트 추진',
                '커피/간식 차량이 현장에 직접 방문하여 응원',
            ]
        },
        {
            'title': '③ 관계사협업은 본업 캠페인',
            'color': RGBColor(0x2E,0x7D,0x32),
            'points': [
                '아이디어 도출 필요',
                '회사내 POP, 플랙카드, 별프로 활용한 캠페인',
                '매월 첫번째 수요일은 콜라보 DAY',
            ]
        },
    ]
    for k, s in enumerate(strategies):
        bx = pad + k * (col3_w + pad)
        add_bullet_box(slide, bx, box_y, col3_w, box_h,
                       title=s['title'],
                       bullets=[{'text': f'• {p}', 'font_size': Pt(10.5), 'color': DARK} for p in s['points']],
                       title_bg=s['color'], body_bg=LGRAY)

    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — HOW? 인사 + 시너지
# ═══════════════════════════════════════════════════════════════════════════════
def create_slide3():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_rect(slide, 0, 0, W, H, fill_color=WHITE)
    add_slide_header(slide, '관계사 협업  HOW?  ② 인사  /  ③ 시너지', slide_num=3)

    pad = Inches(0.25)
    content_top = Inches(1.05)
    col_w = (W - pad*3) / 2
    box_h = Inches(5.8)

    # ── 인사 (left) ──────────────────────────────────────────────────────────
    left_x = pad
    # section header
    sec = add_rect(slide, left_x, content_top, col_w, Inches(0.42), fill_color=BLUE)
    tf = sec.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = '■ 인사'; run.font.size = Pt(13); run.font.bold = True; run.font.color.rgb = WHITE

    body_rect = add_rect(slide, left_x, content_top + Inches(0.42), col_w, box_h - Inches(0.42), fill_color=LGRAY)
    body_rect.line.color.rgb = RGBColor(0xCC,0xCC,0xCC); body_rect.line.width = Pt(0.5)

    add_multi_para_textbox(slide, left_x + Inches(0.18), content_top + Inches(0.55),
                           col_w - Inches(0.3), box_h - Inches(0.65), [
        {'text': '1. 협업에 대한 信賞 (신상필벌)', 'font_size': Pt(12), 'bold': True, 'color': BLUE, 'space_after': Pt(4)},
        {'text': '• 우수직원에 대한 인사 가점', 'font_size': Pt(11), 'color': DARK},
        {'text': '  → 분기별 표창을 통하여 CEO표창에 준하는 수준으로 부여', 'font_size': Pt(10.5), 'color': DGRAY},
        {'text': '', 'font_size': Pt(5)},
        {'text': '• 우수 외부연수 기회 우선배정', 'font_size': Pt(11), 'color': DARK},
        {'text': '  → 배정인원의 30% 우선 배정', 'font_size': Pt(10.5), 'color': DGRAY},
    ])

    # ── 시너지 (right) ───────────────────────────────────────────────────────
    right_x = pad*2 + col_w
    sec2 = add_rect(slide, right_x, content_top, col_w, Inches(0.42), fill_color=RED)
    tf2 = sec2.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run(); run2.text = '■ 시너지'; run2.font.size = Pt(13); run2.font.bold = True; run2.font.color.rgb = WHITE

    body2 = add_rect(slide, right_x, content_top + Inches(0.42), col_w, box_h - Inches(0.42), fill_color=LGRAY)
    body2.line.color.rgb = RGBColor(0xCC,0xCC,0xCC); body2.line.width = Pt(0.5)

    add_multi_para_textbox(slide, right_x + Inches(0.18), content_top + Inches(0.55),
                           col_w - Inches(0.3), box_h - Inches(0.65), [
        {'text': '1. 콜라보북 제작 및 전직원 배포', 'font_size': Pt(12), 'bold': True, 'color': RED, 'space_after': Pt(3)},
        {'text': '  "하나로 연결된 모두의 금융" 실천', 'font_size': Pt(11), 'bold': True, 'color': DARK},
        {'text': '• 관계사 상품정보 및 상품 리플렛이 담긴 콜라보 포켓북 제작·배포', 'font_size': Pt(10.5), 'color': DGRAY},
        {'text': '• 손님을 만날 때 자동적으로 관계사 상품을 제안하는 문화 만들기', 'font_size': Pt(10.5), 'color': DGRAY},
        {'text': '', 'font_size': Pt(5)},
        {'text': '2. 그룹별 협업 마케팅 활성화 (구조화)', 'font_size': Pt(12), 'bold': True, 'color': RED, 'space_after': Pt(3)},
        {'text': '• WM그룹 : The센터필드W 등 증권 패밀리오피스 자산가 전용 하나카드 제작', 'font_size': Pt(10.5), 'color': DGRAY},
        {'text': '  → 자산가의 증권 로열티 함양', 'font_size': Pt(10), 'color': MGRAY},
        {'text': '• IB그룹 : 인프라, 부동산, 물류창고 공사 등 Deal 발주단계에서 관계사 상품 패키지 involve', 'font_size': Pt(10.5), 'color': DGRAY},
        {'text': '• S&T그룹 : 당사 법인금상손님  하나저축은행 금리상품 판매 확대', 'font_size': Pt(10.5), 'color': DGRAY},
        {'text': '  (금리형상품 600억 매각)', 'font_size': Pt(10), 'color': MGRAY},
        {'text': '', 'font_size': Pt(5)},
        {'text': '3. 관계사 협업 실적 점검 (주간단위)', 'font_size': Pt(12), 'bold': True, 'color': RED, 'space_after': Pt(3)},
        {'text': '  → CEO께서 직접 관리', 'font_size': Pt(11), 'bold': True, 'color': DARK},
    ])

    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — 협업 우수사례
# ═══════════════════════════════════════════════════════════════════════════════
def create_slide4():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_rect(slide, 0, 0, W, H, fill_color=WHITE)
    add_slide_header(slide, '협업 우수사례', slide_num=4)

    pad = Inches(0.25)
    content_top = Inches(1.1)
    col_w = (W - pad*4) / 3
    cases = [
        {
            'group': 'WM그룹',
            'color': BLUE,
            'summary': '영업1부 해외사업을 정리하고\n국내로 복귀하는 손님에게 하나카드 제안',
            'results': [
                '하나Club1카드 개설',
                '하나은행 예금 150억 유치',
            ]
        },
        {
            'group': 'IB그룹',
            'color': RED,
            'summary': '㈜디엠티 브릿지론 인수금융 딜\n발주단계에 관계사 상품 involve',
            'results': [
                '에스크로예금 하나은행 130억',
                '인수금융 이자적립계좌 하나은행 33억',
                '주거래은행 변경 (신한 → 하나)',
                '하나법인카드 유치 (12좌)',
            ]
        },
        {
            'group': 'S&T그룹',
            'color': RGBColor(0x2E,0x7D,0x32),
            'summary': '시너지팀 협업 사례 수집 중\n(김주형 실장)',
            'results': [
                '추가 사례 발굴 진행 중',
            ]
        },
    ]

    for k, c in enumerate(cases):
        bx = pad + k * (col_w + pad)
        box_h = Inches(5.7)
        # Group header
        hdr = add_rect(slide, bx, content_top, col_w, Inches(0.5), fill_color=c['color'])
        tf = hdr.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = c['group']; run.font.size = Pt(14); run.font.bold = True; run.font.color.rgb = WHITE

        # Summary box
        sum_h = Inches(1.5)
        sum_box = add_rect(slide, bx, content_top + Inches(0.5), col_w, sum_h,
                           fill_color=LBLUE if c['color']==BLUE else (LTRED if c['color']==RED else RGBColor(0xE8,0xF5,0xE9)))
        sum_box.line.color.rgb = RGBColor(0xCC,0xCC,0xCC); sum_box.line.width = Pt(0.5)
        add_textbox(slide, bx + Inches(0.12), content_top + Inches(0.58),
                    col_w - Inches(0.22), sum_h - Inches(0.15),
                    c['summary'], font_size=Pt(11), color=DGRAY, wrap=True)

        # Results
        res_y = content_top + Inches(0.5) + sum_h + Inches(0.1)
        res_h = box_h - Inches(0.5) - sum_h - Inches(0.15)
        res_box = add_rect(slide, bx, res_y, col_w, res_h, fill_color=LGRAY)
        res_box.line.color.rgb = RGBColor(0xCC,0xCC,0xCC); res_box.line.width = Pt(0.5)

        # "성과" label
        add_textbox(slide, bx + Inches(0.12), res_y + Inches(0.1), col_w - Inches(0.22), Inches(0.3),
                    '주요 성과', font_size=Pt(10), bold=True, color=c['color'])
        add_rect(slide, bx + Inches(0.12), res_y + Inches(0.38), col_w - Inches(0.28), Pt(1), fill_color=RGBColor(0xCC,0xCC,0xCC))

        bullet_paras = []
        for r in c['results']:
            bullet_paras.append({'text': f'✓  {r}', 'font_size': Pt(10.5), 'bold': True, 'color': DARK})
        add_multi_para_textbox(slide, bx + Inches(0.12), res_y + Inches(0.45),
                               col_w - Inches(0.22), res_h - Inches(0.55), bullet_paras)

    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — 결론 / 실천 선언
# ═══════════════════════════════════════════════════════════════════════════════
def create_slide5():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_rect(slide, 0, 0, W, H, fill_color=WHITE)
    add_slide_header(slide, '실천 선언 — 관계사 협업!  HOW!', slide_num=5)

    pad = Inches(0.25)
    content_top = Inches(1.1)

    # Three circles for 삼위일체
    circle_data = [
        ('시너지기획', BLUE, Inches(1.5)),
        ('기업문화', RED, Inches(5.0)),
        ('인사', RGBColor(0x2E,0x7D,0x32), Inches(8.5)),
    ]
    c_size = Inches(1.6)
    c_top  = content_top + Inches(0.2)
    for label, col, cx in circle_data:
        ellipse = slide.shapes.add_shape(9, cx - c_size/2, c_top, c_size, c_size)  # oval=9
        ellipse.fill.solid(); ellipse.fill.fore_color.rgb = col
        ellipse.line.fill.background()
        tf = ellipse.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = label
        run.font.size = Pt(13); run.font.bold = True; run.font.color.rgb = WHITE

    # Plus signs between circles
    for cx in [Inches(3.35), Inches(6.85)]:
        add_textbox(slide, cx - Inches(0.3), c_top + Inches(0.5), Inches(0.6), Inches(0.6),
                    '+', font_size=Pt(24), bold=True, color=MGRAY, align=PP_ALIGN.CENTER)

    # = 삼위일체
    add_textbox(slide, Inches(9.4), c_top + Inches(0.5), Inches(0.3), Inches(0.6),
                '=', font_size=Pt(24), bold=True, color=MGRAY, align=PP_ALIGN.CENTER)
    three_rect = add_rect(slide, Inches(9.8), c_top, Inches(2.8), c_size, fill_color=DARK)
    tf3 = three_rect.text_frame; p3 = tf3.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run(); run3.text = '삼위일체\n협업문화 장착'; run3.font.size = Pt(14); run3.font.bold = True; run3.font.color.rgb = WHITE

    # Mid divider
    div_y = content_top + c_size + Inches(0.35)
    add_rect(slide, pad, div_y, W - pad*2, Pt(2), fill_color=RED)

    # Main declaration text
    decl_y = div_y + Inches(0.15)
    decl_box = add_rect(slide, pad, decl_y, W - pad*2, Inches(1.0), fill_color=DARK)
    add_multi_para_textbox(slide, pad + Inches(0.3), decl_y + Inches(0.12), W - pad*2 - Inches(0.4), Inches(0.8), [
        {'text': '■ 관계사 협업!  HOW!   —  하나증권이 솔선수범하여 그룹 전체를 움직이겠습니다', 'font_size': Pt(14), 'bold': True, 'color': WHITE},
    ])

    # HOW details
    how_y = decl_y + Inches(1.1)
    col_w = (W - pad*3) / 2

    # Left block
    left_b = add_rect(slide, pad, how_y, col_w, Inches(2.8), fill_color=LGRAY)
    left_b.line.color.rgb = RGBColor(0xDD,0xDD,0xDD); left_b.line.width = Pt(0.5)
    add_multi_para_textbox(slide, pad + Inches(0.18), how_y + Inches(0.15), col_w - Inches(0.3), Inches(2.55), [
        {'text': '■ 손님에게! 관계사 협업에! 1년만 미쳐라!', 'font_size': Pt(12), 'bold': True, 'color': RED, 'space_after': Pt(6)},
        {'text': '   Why?  하나증권의 생존을 위해서!', 'font_size': Pt(11.5), 'bold': True, 'color': DARK, 'space_after': Pt(8)},
        {'text': '→ 하나증권이 솔선수범하여 그룹 전체를 움직이게 할 것입니다!', 'font_size': Pt(11), 'color': DGRAY},
        {'text': '→ 증권을 위해서!  손님을 위해서!  관계사 협업은 선택이 아니라', 'font_size': Pt(11), 'color': DGRAY},
    ])

    # Right block — 핵심 문구
    right_x = pad*2 + col_w
    right_b = add_rect(slide, right_x, how_y, col_w, Inches(2.8), fill_color=DARK)
    add_multi_para_textbox(slide, right_x + Inches(0.2), how_y + Inches(0.25), col_w - Inches(0.3), Inches(2.3), [
        {'text': '"하! 卽! 生 !"    "하! 卽! 死 !"', 'font_size': Pt(18), 'bold': True, 'color': RED, 'align': PP_ALIGN.CENTER},
        {'text': '', 'font_size': Pt(6)},
        {'text': '하면 살고, 하지 않으면 죽는다', 'font_size': Pt(13), 'bold': True, 'color': WHITE, 'align': PP_ALIGN.CENTER},
        {'text': '', 'font_size': Pt(8)},
        {'text': '손님을 만날 때, 은행은 어느 곳을 거래하는지?\n카드는 어떤 카드를 사용하는지?\n궁금해 하며 관계사 상품을 제안하는 문화를\n반드시 만들겠습니다!', 'font_size': Pt(10.5), 'color': LGRAY, 'align': PP_ALIGN.CENTER},
    ])

    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# Build & Save
# ═══════════════════════════════════════════════════════════════════════════════
print('Building slides...')
create_slide1()
print('  Slide 1 done')
create_slide2()
print('  Slide 2 done')
create_slide3()
print('  Slide 3 done')
create_slide4()
print('  Slide 4 done')
create_slide5()
print('  Slide 5 done')

output_path = r'C:\Users\jay\Downloads\경영지원그룹장_CEO업무보고_임원보고용.pptx'
prs.save(output_path)
print(f'\n✓ Saved to: {output_path}')
