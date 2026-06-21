import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

SRC = r'C:\Users\jay\Downloads\AI교육_완성판_v13.pptx'
prs = Presentation(SRC)
print(f'총 슬라이드 수: {len(prs.slides)}\n')

for si, slide in enumerate(prs.slides):
    print(f'\n=== Slide {si+1} ===')
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    print(f'  {t}')
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            print(f'  [기존 노트] {notes[:80]}')
